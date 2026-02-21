"""Tests for the PPTX renderer module - DSL to PPTX conversion."""

import io
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

from backend.dsl.schema import (
    BoundingBox,
    Canvas,
    Effects,
    SolidFill,
    SlideMetadata,
    Shadow,
    Shape,
    SlideScene,
    Stroke,
    TextContent,
    TextRun,
    ThemeColors,
    Transform,
)
from backend.renderer import PPTXWriter, ShapeRenderer, StyleRenderer, TextRenderer
from backend.constraints import ConstraintEngine, ArchetypeRules


class TestShapeRenderer:
    """Tests for ShapeRenderer."""

    def test_render_rectangle(self):
        """Test rendering a basic rectangle."""
        shape = Shape(
            id="test_rect",
            type="autoShape",
            bbox=BoundingBox(x=914400, y=914400, width=2743200, height=914400),
            transform=Transform(),
            fill=SolidFill(color="#0D9488"),
            auto_shape_type="rect",
        )

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        theme = ThemeColors()

        renderer = ShapeRenderer()
        pptx_shape = renderer.render(slide, shape, theme)

        # ShapeRenderer.render returns None (adds shape to slide)
        # Check that a shape was added to the slide
        assert len(slide.shapes) >= 1

    def test_render_roundrect(self):
        """Test rendering a rounded rectangle."""
        shape = Shape(
            id="test_roundrect",
            type="autoShape",
            bbox=BoundingBox(x=1000000, y=1000000, width=2000000, height=800000),
            transform=Transform(),
            fill=SolidFill(color="#14B8A6"),
            auto_shape_type="roundRect",
        )

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = ThemeColors()

        renderer = ShapeRenderer()
        renderer.render(slide, shape, theme)

        assert len(slide.shapes) >= 1

    def test_render_chevron(self):
        """Test rendering a chevron shape."""
        shape = Shape(
            id="test_chevron",
            type="autoShape",
            bbox=BoundingBox(x=500000, y=500000, width=1500000, height=600000),
            transform=Transform(),
            fill=SolidFill(color="#2DD4BF"),
            auto_shape_type="chevron",
        )

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = ThemeColors()

        renderer = ShapeRenderer()
        renderer.render(slide, shape, theme)

        assert len(slide.shapes) >= 1

    def test_render_ellipse(self):
        """Test rendering an ellipse."""
        shape = Shape(
            id="test_ellipse",
            type="autoShape",
            bbox=BoundingBox(x=2000000, y=2000000, width=1000000, height=1000000),
            transform=Transform(),
            fill=SolidFill(color="#5EEAD4"),
            auto_shape_type="ellipse",
        )

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = ThemeColors()

        renderer = ShapeRenderer()
        renderer.render(slide, shape, theme)

        assert len(slide.shapes) >= 1

    def test_render_with_text(self):
        """Test rendering a shape with text."""
        shape = Shape(
            id="test_text_shape",
            type="autoShape",
            bbox=BoundingBox(x=914400, y=914400, width=2743200, height=914400),
            transform=Transform(),
            fill=SolidFill(color="#0D9488"),
            auto_shape_type="rect",
            text=TextContent(
                runs=[TextRun(text="Hello World", font_size=1800)],
                alignment="center",
            ),
        )

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = ThemeColors()

        renderer = ShapeRenderer()
        renderer.render(slide, shape, theme)

        # Check text was rendered
        pptx_shape = slide.shapes[0]
        assert pptx_shape.has_text_frame
        assert "Hello World" in pptx_shape.text


class TestStyleRenderer:
    """Tests for StyleRenderer."""

    def test_apply_solid_fill(self):
        """Test applying solid fill."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pptx_shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(1000000), Emu(500000))

        fill = SolidFill(color="#0D9488", alpha=1.0)
        theme = ThemeColors()

        renderer = StyleRenderer()
        renderer.apply_fill(pptx_shape, fill, theme)

        # Verify fill was applied
        assert pptx_shape.fill.type is not None

    def test_apply_stroke(self):
        """Test applying stroke/border."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pptx_shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(1000000), Emu(500000))

        stroke = Stroke(color="#000000", width=25400, alpha=1.0)
        theme = ThemeColors()

        renderer = StyleRenderer()
        renderer.apply_stroke(pptx_shape, stroke, theme)

        # Verify stroke was applied
        assert pptx_shape.line.width is not None

    def test_apply_shadow_effect(self):
        """Test applying shadow effect."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pptx_shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(1000000), Emu(500000))

        shadow = Shadow(
            enabled=True,
            blur_radius=50800,
            distance=38100,
            direction=45.0,
            color="#000000",
            alpha=0.5,
        )
        effects = Effects(shadow=shadow)

        renderer = StyleRenderer()
        renderer.apply_effects(pptx_shape, effects)

        # Shadow is applied via XML, so just verify no errors


class TestTextRenderer:
    """Tests for TextRenderer."""

    def test_render_simple_text(self):
        """Test rendering simple text."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pptx_shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(2000000), Emu(500000))

        text_content = TextContent(
            runs=[TextRun(text="Test Text", font_size=1400)],
            alignment="left",
        )
        theme = ThemeColors()

        renderer = TextRenderer()
        renderer.render(pptx_shape, text_content, theme)

        assert pptx_shape.text == "Test Text"

    def test_render_multiple_runs(self):
        """Test rendering multiple text runs."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pptx_shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(2000000), Emu(1000000))

        text_content = TextContent(
            runs=[
                TextRun(text="First "),
                TextRun(text="Second"),
            ],
            alignment="left",
        )
        theme = ThemeColors()

        renderer = TextRenderer()
        renderer.render(pptx_shape, text_content, theme)

        assert "First" in pptx_shape.text
        assert "Second" in pptx_shape.text

    def test_render_formatted_text(self):
        """Test rendering bold and italic text."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pptx_shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(2000000), Emu(500000))

        text_content = TextContent(
            runs=[
                TextRun(text="Bold ", bold=True),
                TextRun(text="Italic ", italic=True),
                TextRun(text="Normal"),
            ],
            alignment="center",
        )
        theme = ThemeColors()

        renderer = TextRenderer()
        renderer.render(pptx_shape, text_content, theme)

        assert "Bold" in pptx_shape.text
        assert "Italic" in pptx_shape.text
        assert "Normal" in pptx_shape.text


class TestPPTXWriter:
    """Tests for PPTXWriter - full scene to PPTX conversion."""

    def test_write_single_slide(self, sample_slide_scene):
        """Test writing a single slide scene."""
        scene = SlideScene(**sample_slide_scene)
        writer = PPTXWriter()

        output = io.BytesIO()
        writer.write_single(scene, output)

        # Verify PPTX was created
        output.seek(0)
        prs = Presentation(output)
        assert len(prs.slides) == 1

    def test_write_multiple_slides(self):
        """Test writing multiple slides."""
        canvas = Canvas()
        theme = ThemeColors()

        scenes = []
        for i in range(3):
            shapes = [
                Shape(
                    id=f"shape_{i}",
                    type="autoShape",
                    bbox=BoundingBox(x=914400, y=914400 + i * 1000000, width=2000000, height=500000),
                    transform=Transform(),
                    fill=SolidFill(color="#0D9488"),
                    auto_shape_type="rect",
                )
            ]
            scenes.append(
                SlideScene(
                    canvas=canvas,
                    shapes=shapes,
                    theme=theme,
                    metadata=SlideMetadata(slide_number=i + 1),
                )
            )

        writer = PPTXWriter()
        output = io.BytesIO()
        writer.write(scenes, output)

        output.seek(0)
        prs = Presentation(output)
        assert len(prs.slides) == 3

    def test_write_to_file(self, sample_slide_scene):
        """Test writing to a file path."""
        scene = SlideScene(**sample_slide_scene)
        writer = PPTXWriter()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            temp_path = Path(f.name)

        try:
            writer.write_single(scene, str(temp_path))
            assert temp_path.exists()
            assert temp_path.stat().st_size > 0

            # Verify it's a valid PPTX
            prs = Presentation(str(temp_path))
            assert len(prs.slides) == 1
        finally:
            temp_path.unlink()

    def test_write_bytes(self, sample_slide_scene):
        """Test returning bytes instead of writing to file."""
        scene = SlideScene(**sample_slide_scene)
        writer = PPTXWriter()

        result = writer.write_single(scene, None)

        assert isinstance(result, bytes)
        assert len(result) > 0

        # Verify it's a valid PPTX
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 1


class TestConstraintEngine:
    """Tests for the constraint engine."""

    def test_validate_scene(self, sample_slide_scene):
        """Test validating a scene."""
        scene = SlideScene(**sample_slide_scene)
        engine = ConstraintEngine()

        result = engine.validate(scene)

        assert result.is_valid is True
        assert result.score > 0

    def test_detect_out_of_bounds(self):
        """Test detecting out-of-bounds shapes."""
        shape = Shape(
            id="out_of_bounds",
            type="autoShape",
            bbox=BoundingBox(x=-100000, y=500000, width=1000000, height=500000),
            transform=Transform(),
            fill=SolidFill(color="#FF0000"),
            auto_shape_type="rect",
        )

        scene = SlideScene(
            canvas=Canvas(),
            shapes=[shape],
            theme=ThemeColors(),
            metadata=SlideMetadata(),
        )

        engine = ConstraintEngine()
        result = engine.validate(scene)

        assert result.is_valid is False
        assert len([v for v in result.violations if v.rule == "bounds"]) > 0

    def test_fix_out_of_bounds(self):
        """Test fixing out-of-bounds shapes."""
        shape = Shape(
            id="out_of_bounds",
            type="autoShape",
            bbox=BoundingBox(x=-100000, y=500000, width=1000000, height=500000),
            transform=Transform(),
            fill=SolidFill(color="#FF0000"),
            auto_shape_type="rect",
        )

        scene = SlideScene(
            canvas=Canvas(),
            shapes=[shape],
            theme=ThemeColors(),
            metadata=SlideMetadata(),
        )

        engine = ConstraintEngine()
        fixed_scene = engine.fix(scene)

        # Verify shape is now in bounds
        fixed_shape = fixed_scene.shapes[0]
        assert fixed_shape.bbox.x >= 0

    def test_detect_overlaps(self):
        """Test detecting overlapping shapes."""
        shapes = [
            Shape(
                id="shape1",
                type="autoShape",
                bbox=BoundingBox(x=1000000, y=1000000, width=1000000, height=500000),
                transform=Transform(),
                fill=SolidFill(color="#FF0000"),
                auto_shape_type="rect",
            ),
            Shape(
                id="shape2",
                type="autoShape",
                bbox=BoundingBox(x=1500000, y=1200000, width=1000000, height=500000),
                transform=Transform(),
                fill=SolidFill(color="#00FF00"),
                auto_shape_type="rect",
            ),
        ]

        scene = SlideScene(
            canvas=Canvas(),
            shapes=shapes,
            theme=ThemeColors(),
            metadata=SlideMetadata(),
        )

        engine = ConstraintEngine()
        result = engine.validate(scene)

        overlap_violations = [v for v in result.violations if v.rule == "overlap"]
        assert len(overlap_violations) > 0


class TestArchetypeRules:
    """Tests for archetype-specific layout rules."""

    def test_funnel_rules(self):
        """Test funnel archetype rules."""
        rules = ArchetypeRules.get_rules("funnel")
        assert len(rules) > 0
        rule_names = [r.name for r in rules]
        assert "center_horizontal" in rule_names
        assert "decreasing_width" in rule_names

    def test_pyramid_rules(self):
        """Test pyramid archetype rules."""
        rules = ArchetypeRules.get_rules("pyramid")
        assert len(rules) > 0
        rule_names = [r.name for r in rules]
        assert "center_horizontal" in rule_names
        assert "increasing_width" in rule_names

    def test_timeline_rules(self):
        """Test timeline archetype rules."""
        rules = ArchetypeRules.get_rules("timeline")
        assert len(rules) > 0
        rule_names = [r.name for r in rules]
        assert "distribute_horizontal" in rule_names

    def test_hub_spoke_rules(self):
        """Test hub and spoke archetype rules."""
        rules = ArchetypeRules.get_rules("hub_spoke")
        assert len(rules) > 0
        rule_names = [r.name for r in rules]
        assert "center_hub" in rule_names
        assert "radial_spokes" in rule_names

    def test_cycle_rules(self):
        """Test cycle archetype rules."""
        rules = ArchetypeRules.get_rules("cycle")
        assert len(rules) > 0
        rule_names = [r.name for r in rules]
        assert "circular_layout" in rule_names

    def test_apply_funnel_rules(self):
        """Test applying funnel rules to shapes."""
        shapes = [
            Shape(
                id=f"stage_{i}",
                type="autoShape",
                bbox=BoundingBox(x=1000000, y=1000000 + i * 800000, width=3000000, height=600000),
                transform=Transform(),
                fill=SolidFill(color="#0D9488"),
                auto_shape_type="trapezoid",
            )
            for i in range(4)
        ]

        scene = SlideScene(
            canvas=Canvas(),
            shapes=shapes,
            theme=ThemeColors(),
            metadata=SlideMetadata(archetype="funnel"),
        )

        fixed_scene = ArchetypeRules.apply_rules(scene)

        # Verify widths decrease from top to bottom
        widths = [s.bbox.width for s in sorted(fixed_scene.shapes, key=lambda s: s.bbox.y)]
        assert all(widths[i] >= widths[i + 1] for i in range(len(widths) - 1))

    def test_apply_pyramid_rules(self):
        """Test applying pyramid rules to shapes."""
        shapes = [
            Shape(
                id=f"level_{i}",
                type="autoShape",
                bbox=BoundingBox(x=2000000, y=500000 + i * 600000, width=2000000, height=500000),
                transform=Transform(),
                fill=SolidFill(color="#0D9488"),
                auto_shape_type="trapezoid",
            )
            for i in range(4)
        ]

        scene = SlideScene(
            canvas=Canvas(),
            shapes=shapes,
            theme=ThemeColors(),
            metadata=SlideMetadata(archetype="pyramid"),
        )

        fixed_scene = ArchetypeRules.apply_rules(scene)

        # Verify widths increase from top to bottom
        widths = [s.bbox.width for s in sorted(fixed_scene.shapes, key=lambda s: s.bbox.y)]
        assert all(widths[i] <= widths[i + 1] for i in range(len(widths) - 1))


class TestRoundTrip:
    """Round-trip tests: DSL → PPTX → DSL should produce similar results."""

    def test_basic_roundtrip(self, sample_slide_scene):
        """Test basic round-trip conversion."""
        # Create a scene
        original_scene = SlideScene(**sample_slide_scene)

        # Write to PPTX
        writer = PPTXWriter()
        pptx_bytes = writer.write_single(original_scene, None)

        # Verify PPTX is valid
        assert pptx_bytes is not None
        assert len(pptx_bytes) > 0

        # Parse PPTX back (using python-pptx directly for now)
        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) == 1

        # Verify shapes exist
        slide = prs.slides[0]
        assert len(slide.shapes) == len(original_scene.shapes)

    def test_multiple_shapes_roundtrip(self):
        """Test round-trip with multiple shapes."""
        shapes = [
            Shape(
                id=f"shape_{i}",
                type="autoShape",
                bbox=BoundingBox(
                    x=1000000 + i * 2500000,
                    y=2000000,
                    width=2000000,
                    height=1000000,
                ),
                transform=Transform(),
                fill=SolidFill(color=f"#{i+1:02x}{i+2:02x}{i+3:02x}"),
                auto_shape_type="rect",
            )
            for i in range(3)
        ]

        scene = SlideScene(
            canvas=Canvas(),
            shapes=shapes,
            theme=ThemeColors(),
            metadata=SlideMetadata(slide_number=1),
        )

        # Write and read back
        writer = PPTXWriter()
        pptx_bytes = writer.write_single(scene, None)

        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides[0].shapes) == 3

    def test_text_roundtrip(self):
        """Test round-trip with text content."""
        shape = Shape(
            id="text_shape",
            type="autoShape",
            bbox=BoundingBox(x=1000000, y=1000000, width=4000000, height=1500000),
            transform=Transform(),
            fill=SolidFill(color="#0D9488"),
            auto_shape_type="rect",
            text=TextContent(
                runs=[
                    TextRun(text="Title Text", font_size=2400, bold=True),
                    TextRun(text=" Body text goes here", font_size=1400),
                ],
                alignment="center",
            ),
        )

        scene = SlideScene(
            canvas=Canvas(),
            shapes=[shape],
            theme=ThemeColors(),
            metadata=SlideMetadata(),
        )

        writer = PPTXWriter()
        pptx_bytes = writer.write_single(scene, None)

        prs = Presentation(io.BytesIO(pptx_bytes))
        slide_shape = prs.slides[0].shapes[0]

        assert slide_shape.has_text_frame
        assert "Title Text" in slide_shape.text
        assert "Body text" in slide_shape.text


class TestIntegration:
    """Integration tests for the full rendering pipeline."""

    def test_funnel_diagram_generation(self):
        """Test generating a complete funnel diagram."""
        stages = ["Awareness", "Interest", "Decision", "Action"]
        shapes = []

        for i, stage in enumerate(stages):
            shapes.append(
                Shape(
                    id=f"funnel_{i}",
                    type="autoShape",
                    bbox=BoundingBox(
                        x=2000000,
                        y=1000000 + i * 800000,
                        width=6000000 - i * 1000000,
                        height=600000,
                    ),
                    transform=Transform(),
                    fill=SolidFill(color="#0D9488"),
                    auto_shape_type="trapezoid",
                    text=TextContent(
                        runs=[TextRun(text=stage, font_size=1800, bold=True)],
                        alignment="center",
                    ),
                )
            )

        scene = SlideScene(
            canvas=Canvas(),
            shapes=shapes,
            theme=ThemeColors(),
            metadata=SlideMetadata(archetype="funnel", tags=["marketing", "sales"]),
        )

        # Apply archetype rules
        scene = ArchetypeRules.apply_rules(scene)

        # Validate
        engine = ConstraintEngine()
        result = engine.validate(scene)
        assert result.score >= 50  # Reasonable quality

        # Render
        writer = PPTXWriter()
        pptx_bytes = writer.write_single(scene, None)

        # Verify output
        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) == 1
        assert len(prs.slides[0].shapes) == 4

    def test_timeline_diagram_generation(self):
        """Test generating a complete timeline diagram."""
        events = ["Q1 2024", "Q2 2024", "Q3 2024", "Q4 2024"]
        shapes = []

        for i, event in enumerate(events):
            shapes.append(
                Shape(
                    id=f"timeline_{i}",
                    type="autoShape",
                    bbox=BoundingBox(
                        x=1000000 + i * 2800000,
                        y=3000000,
                        width=800000,
                        height=800000,
                    ),
                    transform=Transform(),
                    fill=SolidFill(color="#14B8A6"),
                    auto_shape_type="ellipse",
                    text=TextContent(
                        runs=[TextRun(text=event, font_size=1200)],
                        alignment="center",
                    ),
                )
            )

        scene = SlideScene(
            canvas=Canvas(),
            shapes=shapes,
            theme=ThemeColors(),
            metadata=SlideMetadata(archetype="timeline"),
        )

        # Apply rules and validate
        scene = ArchetypeRules.apply_rules(scene)
        engine = ConstraintEngine()
        result = engine.validate(scene)

        # Render
        writer = PPTXWriter()
        pptx_bytes = writer.write_single(scene, None)

        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides[0].shapes) == 4
