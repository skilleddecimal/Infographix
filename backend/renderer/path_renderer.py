"""Render freeform path shapes to PowerPoint.

Handles Bezier curves, arcs, and complex path commands by building
custom geometry XML for freeform shapes.
"""

from typing import Any

from lxml import etree
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu

from backend.dsl.schema import PathCommand, PathCommandType, Shape


# XML namespaces
NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


class PathRenderer:
    """Renders freeform path shapes with Bezier curve support."""

    def render_path(
        self,
        slide: Slide,
        shape: Shape,
    ) -> Any | None:
        """Render a freeform path shape.

        Creates a custom geometry shape with proper Bezier curve support
        by building the XML directly.

        Args:
            slide: The PowerPoint slide.
            shape: The DSL shape with path commands.

        Returns:
            The created PowerPoint shape or None if failed.
        """
        if not shape.path:
            return None

        try:
            # Use python-pptx's freeform builder for simple paths
            if self._is_simple_path(shape.path):
                return self._render_simple_path(slide, shape)
            else:
                return self._render_complex_path(slide, shape)
        except Exception:
            return None

    def _is_simple_path(self, path: list[PathCommand]) -> bool:
        """Check if path only contains simple commands.

        Simple paths only have moveTo, lineTo, and close - no curves.

        Args:
            path: List of path commands.

        Returns:
            True if path is simple (no Bezier curves).
        """
        for cmd in path:
            if cmd.type in (
                PathCommandType.CURVE_TO,
                PathCommandType.QUAD_TO,
                PathCommandType.ARC_TO,
            ):
                return False
        return True

    def _render_simple_path(self, slide: Slide, shape: Shape) -> Any | None:
        """Render a simple path using python-pptx's freeform builder.

        Args:
            slide: The PowerPoint slide.
            shape: The DSL shape.

        Returns:
            The created shape or None.
        """
        if not shape.path:
            return None

        # Find starting point
        start_x = shape.bbox.x
        start_y = shape.bbox.y

        for cmd in shape.path:
            if cmd.type == PathCommandType.MOVE_TO and cmd.x is not None:
                start_x = shape.bbox.x + cmd.x
                start_y = shape.bbox.y + (cmd.y or 0)
                break

        builder = slide.shapes.build_freeform(Emu(start_x), Emu(start_y))

        # Track current position for relative moves
        curr_x = start_x
        curr_y = start_y

        for cmd in shape.path:
            if cmd.type == PathCommandType.MOVE_TO:
                # Move is handled by starting position or line segment
                if cmd.x is not None:
                    curr_x = shape.bbox.x + cmd.x
                    curr_y = shape.bbox.y + (cmd.y or 0)

            elif cmd.type == PathCommandType.LINE_TO:
                if cmd.x is not None and cmd.y is not None:
                    abs_x = shape.bbox.x + cmd.x
                    abs_y = shape.bbox.y + cmd.y
                    builder.add_line_segments(
                        [(Emu(abs_x), Emu(abs_y))],
                        close=False,
                    )
                    curr_x = abs_x
                    curr_y = abs_y

            elif cmd.type == PathCommandType.CLOSE:
                pass  # Close is implicit

        return builder.convert_to_shape(Emu(start_x), Emu(start_y))

    def _render_complex_path(self, slide: Slide, shape: Shape) -> Any | None:
        """Render a complex path with Bezier curves using XML.

        Creates a shape with custom geometry XML to support cubic and
        quadratic Bezier curves.

        Args:
            slide: The PowerPoint slide.
            shape: The DSL shape.

        Returns:
            The created shape or None.
        """
        if not shape.path:
            return None

        # First create a placeholder rectangle shape
        from pptx.enum.shapes import MSO_SHAPE

        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(shape.bbox.x),
            Emu(shape.bbox.y),
            Emu(shape.bbox.width),
            Emu(shape.bbox.height),
        )

        # Now replace its geometry with custom geometry
        sp = placeholder._element
        spPr = sp.find(qn("p:spPr"))

        if spPr is None:
            return placeholder

        # Remove the preset geometry
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            spPr.remove(prstGeom)

        # Create custom geometry
        custGeom = etree.SubElement(spPr, qn("a:custGeom"))

        # Add adjustment values (empty for now)
        etree.SubElement(custGeom, qn("a:avLst"))

        # Add guide definitions (empty)
        etree.SubElement(custGeom, qn("a:gdLst"))

        # Add connection sites (empty)
        etree.SubElement(custGeom, qn("a:cxnLst"))

        # Add shape rectangle
        rect = etree.SubElement(custGeom, qn("a:rect"))
        rect.set("l", "0")
        rect.set("t", "0")
        rect.set("r", "r")
        rect.set("b", "b")

        # Add path list
        pathLst = etree.SubElement(custGeom, qn("a:pathLst"))

        # Create path with shape dimensions
        path_elem = etree.SubElement(pathLst, qn("a:path"))
        path_elem.set("w", str(shape.bbox.width))
        path_elem.set("h", str(shape.bbox.height))

        # Build path commands
        self._build_path_commands(path_elem, shape.path)

        return placeholder

    def _build_path_commands(
        self,
        path_elem: Any,
        commands: list[PathCommand],
    ) -> None:
        """Build XML path command elements.

        Args:
            path_elem: The <a:path> element to populate.
            commands: List of PathCommand objects.
        """
        for cmd in commands:
            if cmd.type == PathCommandType.MOVE_TO:
                self._add_move_to(path_elem, cmd)
            elif cmd.type == PathCommandType.LINE_TO:
                self._add_line_to(path_elem, cmd)
            elif cmd.type == PathCommandType.CURVE_TO:
                self._add_cubic_bezier(path_elem, cmd)
            elif cmd.type == PathCommandType.QUAD_TO:
                self._add_quad_bezier(path_elem, cmd)
            elif cmd.type == PathCommandType.ARC_TO:
                self._add_arc_to(path_elem, cmd)
            elif cmd.type == PathCommandType.CLOSE:
                etree.SubElement(path_elem, qn("a:close"))

    def _add_move_to(self, path_elem: Any, cmd: PathCommand) -> None:
        """Add moveTo command to path.

        Args:
            path_elem: The path element.
            cmd: The PathCommand.
        """
        if cmd.x is None or cmd.y is None:
            return

        moveTo = etree.SubElement(path_elem, qn("a:moveTo"))
        pt = etree.SubElement(moveTo, qn("a:pt"))
        pt.set("x", str(cmd.x))
        pt.set("y", str(cmd.y))

    def _add_line_to(self, path_elem: Any, cmd: PathCommand) -> None:
        """Add lineTo command to path.

        Args:
            path_elem: The path element.
            cmd: The PathCommand.
        """
        if cmd.x is None or cmd.y is None:
            return

        lnTo = etree.SubElement(path_elem, qn("a:lnTo"))
        pt = etree.SubElement(lnTo, qn("a:pt"))
        pt.set("x", str(cmd.x))
        pt.set("y", str(cmd.y))

    def _add_cubic_bezier(self, path_elem: Any, cmd: PathCommand) -> None:
        """Add cubic Bezier curve command to path.

        Args:
            path_elem: The path element.
            cmd: The PathCommand with control points.
        """
        if cmd.x is None or cmd.y is None:
            return
        if cmd.x1 is None or cmd.y1 is None:
            return
        if cmd.x2 is None or cmd.y2 is None:
            return

        cubicBezTo = etree.SubElement(path_elem, qn("a:cubicBezTo"))

        # Control point 1
        pt1 = etree.SubElement(cubicBezTo, qn("a:pt"))
        pt1.set("x", str(cmd.x1))
        pt1.set("y", str(cmd.y1))

        # Control point 2
        pt2 = etree.SubElement(cubicBezTo, qn("a:pt"))
        pt2.set("x", str(cmd.x2))
        pt2.set("y", str(cmd.y2))

        # End point
        pt3 = etree.SubElement(cubicBezTo, qn("a:pt"))
        pt3.set("x", str(cmd.x))
        pt3.set("y", str(cmd.y))

    def _add_quad_bezier(self, path_elem: Any, cmd: PathCommand) -> None:
        """Add quadratic Bezier curve command to path.

        Args:
            path_elem: The path element.
            cmd: The PathCommand with control point.
        """
        if cmd.x is None or cmd.y is None:
            return
        if cmd.x1 is None or cmd.y1 is None:
            return

        quadBezTo = etree.SubElement(path_elem, qn("a:quadBezTo"))

        # Control point
        pt1 = etree.SubElement(quadBezTo, qn("a:pt"))
        pt1.set("x", str(cmd.x1))
        pt1.set("y", str(cmd.y1))

        # End point
        pt2 = etree.SubElement(quadBezTo, qn("a:pt"))
        pt2.set("x", str(cmd.x))
        pt2.set("y", str(cmd.y))

    def _add_arc_to(self, path_elem: Any, cmd: PathCommand) -> None:
        """Add arc command to path.

        Args:
            path_elem: The path element.
            cmd: The PathCommand with arc parameters.
        """
        arcTo = etree.SubElement(path_elem, qn("a:arcTo"))

        # Set arc parameters
        if cmd.width_radius is not None:
            arcTo.set("wR", str(cmd.width_radius))
        if cmd.height_radius is not None:
            arcTo.set("hR", str(cmd.height_radius))
        if cmd.start_angle is not None:
            # Convert degrees to 60000ths of a degree
            arcTo.set("stAng", str(int(cmd.start_angle * 60000)))
        if cmd.swing_angle is not None:
            arcTo.set("swAng", str(int(cmd.swing_angle * 60000)))


def apply_transform_to_shape(pptx_shape: Any, shape: Shape) -> None:
    """Apply transform properties (rotation, flip) to a shape.

    Args:
        pptx_shape: The python-pptx shape object.
        shape: The DSL shape with transform info.
    """
    # Apply rotation
    if shape.transform.rotation != 0:
        pptx_shape.rotation = shape.transform.rotation

    # Apply flips via XML (not directly supported by python-pptx)
    if shape.transform.flip_h or shape.transform.flip_v:
        try:
            sp = pptx_shape._element
            spPr = sp.find(qn("p:spPr"))
            if spPr is None:
                return

            xfrm = spPr.find(qn("a:xfrm"))
            if xfrm is None:
                # Create xfrm if it doesn't exist
                xfrm = etree.SubElement(spPr, qn("a:xfrm"))

            if shape.transform.flip_h:
                xfrm.set("flipH", "1")
            if shape.transform.flip_v:
                xfrm.set("flipV", "1")

        except Exception:
            pass
