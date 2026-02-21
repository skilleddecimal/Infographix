"""High-level PPTX generation from DSL scene graphs."""

from io import BytesIO
from pathlib import Path
from typing import BinaryIO, Union

from pptx import Presentation
from pptx.util import Emu

from backend.dsl.schema import SlideScene
from backend.renderer.shape_renderer import ShapeRenderer
from backend.renderer.style_renderer import StyleRenderer


class PPTXWriter:
    """Generates PPTX files from DSL scene graphs."""

    def __init__(self) -> None:
        """Initialize the PPTX writer."""
        self.shape_renderer = ShapeRenderer()
        self.style_renderer = StyleRenderer()

    def write(
        self,
        scenes: list[SlideScene],
        output: Union[str, Path, BinaryIO, None] = None,
    ) -> bytes | None:
        """Write scene graphs to a PPTX file.

        Args:
            scenes: List of SlideScene objects to render.
            output: Output path, file object, or None to return bytes.

        Returns:
            PPTX bytes if output is None, otherwise None.
        """
        if not scenes:
            raise ValueError("At least one scene is required")

        # Use first scene's canvas for presentation dimensions
        first_canvas = scenes[0].canvas
        prs = Presentation()
        prs.slide_width = Emu(first_canvas.width)
        prs.slide_height = Emu(first_canvas.height)

        # Render each scene as a slide
        for scene in scenes:
            self._render_slide(prs, scene)

        # Save or return bytes
        if output is None:
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer.read()
        elif isinstance(output, (str, Path)):
            prs.save(str(output))
            return None
        else:
            prs.save(output)
            return None

    def write_single(
        self,
        scene: SlideScene,
        output: Union[str, Path, BinaryIO, None] = None,
    ) -> bytes | None:
        """Write a single scene to a PPTX file.

        Args:
            scene: SlideScene to render.
            output: Output path, file object, or None to return bytes.

        Returns:
            PPTX bytes if output is None, otherwise None.
        """
        return self.write([scene], output)

    def _render_slide(self, prs: Presentation, scene: SlideScene) -> None:
        """Render a single scene as a slide.

        Args:
            prs: The Presentation object.
            scene: The SlideScene to render.
        """
        # Use blank layout
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Apply background
        self.style_renderer.apply_background(slide, scene.canvas.background)

        # Sort shapes by z_index
        sorted_shapes = sorted(scene.shapes, key=lambda s: s.z_index)

        # Render each shape
        for shape in sorted_shapes:
            self.shape_renderer.render(slide, shape, scene.theme)

    def create_presentation(
        self,
        width: int = 12192000,
        height: int = 6858000,
    ) -> Presentation:
        """Create a new presentation with specified dimensions.

        Args:
            width: Slide width in EMUs (default 16:9).
            height: Slide height in EMUs (default 16:9).

        Returns:
            New Presentation object.
        """
        prs = Presentation()
        prs.slide_width = Emu(width)
        prs.slide_height = Emu(height)
        return prs
