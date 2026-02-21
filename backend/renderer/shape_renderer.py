"""Render shapes from DSL to PowerPoint."""

from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pptx.util import Emu

from backend.dsl.schema import Shape, ShapeType, ThemeColors
from backend.renderer.style_renderer import StyleRenderer
from backend.renderer.text_renderer import TextRenderer


# Map DSL auto shape types to MSO_SHAPE
AUTO_SHAPE_MAP: dict[str, MSO_SHAPE] = {
    # Basic shapes
    "rect": MSO_SHAPE.RECTANGLE,
    "rectangle": MSO_SHAPE.RECTANGLE,
    "roundrect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "roundedrectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "oval": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "isosceles_triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
    "pentagon": MSO_SHAPE.PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "heptagon": MSO_SHAPE.HEPTAGON,
    "octagon": MSO_SHAPE.OCTAGON,
    "decagon": MSO_SHAPE.DECAGON,
    "dodecagon": MSO_SHAPE.DODECAGON,
    # Arrows
    "arrow": MSO_SHAPE.RIGHT_ARROW,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "left_arrow": MSO_SHAPE.LEFT_ARROW,
    "up_arrow": MSO_SHAPE.UP_ARROW,
    "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "chevron": MSO_SHAPE.CHEVRON,
    "notched_right_arrow": MSO_SHAPE.NOTCHED_RIGHT_ARROW,
    "pentagon_arrow": MSO_SHAPE.CHEVRON,
    "block_arc": MSO_SHAPE.BLOCK_ARC,
    # Flowchart
    "flowchart_process": MSO_SHAPE.FLOWCHART_PROCESS,
    "flowchart_decision": MSO_SHAPE.FLOWCHART_DECISION,
    "flowchart_terminator": MSO_SHAPE.FLOWCHART_TERMINATOR,
    "flowchart_data": MSO_SHAPE.FLOWCHART_DATA,
    # Callouts
    "callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "rectangular_callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "rounded_rectangular_callout": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
    "oval_callout": MSO_SHAPE.OVAL_CALLOUT,
    "cloud_callout": MSO_SHAPE.CLOUD_CALLOUT,
    # Stars and banners
    "star4": MSO_SHAPE.STAR_4_POINT,
    "star5": MSO_SHAPE.STAR_5_POINT,
    "star6": MSO_SHAPE.STAR_6_POINT,
    "star_4_point": MSO_SHAPE.STAR_4_POINT,
    "star_5_point": MSO_SHAPE.STAR_5_POINT,
    "star_6_point": MSO_SHAPE.STAR_6_POINT,
    "ribbon": MSO_SHAPE.UP_RIBBON,
    "ribbon_up": MSO_SHAPE.UP_RIBBON,
    "ribbon_down": MSO_SHAPE.DOWN_RIBBON,
    "curved_down_ribbon": MSO_SHAPE.CURVED_DOWN_RIBBON,
    "curved_up_ribbon": MSO_SHAPE.CURVED_UP_RIBBON,
    "wave": MSO_SHAPE.WAVE,
    # Other
    "heart": MSO_SHAPE.HEART,
    "lightning_bolt": MSO_SHAPE.LIGHTNING_BOLT,
    "sun": MSO_SHAPE.SUN,
    "moon": MSO_SHAPE.MOON,
    "cloud": MSO_SHAPE.CLOUD,
    "arc": MSO_SHAPE.ARC,
    "donut": MSO_SHAPE.DONUT,
    "no_symbol": MSO_SHAPE.NO_SYMBOL,
    "cross": MSO_SHAPE.CROSS,
    "cube": MSO_SHAPE.CUBE,
    "can": MSO_SHAPE.CAN,
    "pie": MSO_SHAPE.PIE,
    "chord": MSO_SHAPE.CHORD,
    "frame": MSO_SHAPE.FRAME,
    "bevel": MSO_SHAPE.BEVEL,
    "folded_corner": MSO_SHAPE.FOLDED_CORNER,
    "smiley_face": MSO_SHAPE.SMILEY_FACE,
    "action_button_home": MSO_SHAPE.ACTION_BUTTON_HOME,
}


class ShapeRenderer:
    """Renders DSL shapes to PowerPoint shapes."""

    def __init__(self) -> None:
        """Initialize the shape renderer."""
        self.style_renderer = StyleRenderer()
        self.text_renderer = TextRenderer()

    def render(self, slide: Slide, shape: Shape, theme: ThemeColors) -> None:
        """Render a shape to a slide.

        Args:
            slide: The PowerPoint slide.
            shape: The DSL shape to render.
            theme: Theme colors for resolving color references.
        """
        if shape.type == ShapeType.AUTO_SHAPE:
            self._render_auto_shape(slide, shape, theme)
        elif shape.type == ShapeType.TEXT:
            self._render_text_box(slide, shape, theme)
        elif shape.type == ShapeType.IMAGE:
            self._render_image(slide, shape)
        elif shape.type == ShapeType.GROUP:
            self._render_group(slide, shape, theme)
        elif shape.type == ShapeType.FREEFORM:
            self._render_freeform(slide, shape, theme)
        elif shape.type == ShapeType.CONNECTOR:
            self._render_connector(slide, shape, theme)

    def _render_auto_shape(
        self,
        slide: Slide,
        shape: Shape,
        theme: ThemeColors,
    ) -> None:
        """Render an auto shape.

        Args:
            slide: The PowerPoint slide.
            shape: The DSL shape.
            theme: Theme colors.
        """
        # Determine MSO shape type
        shape_type_name = (shape.auto_shape_type or "rectangle").lower()
        mso_shape = AUTO_SHAPE_MAP.get(shape_type_name, MSO_SHAPE.RECTANGLE)

        # Create shape
        pptx_shape = slide.shapes.add_shape(
            mso_shape,
            Emu(shape.bbox.x),
            Emu(shape.bbox.y),
            Emu(shape.bbox.width),
            Emu(shape.bbox.height),
        )

        # Apply rotation
        if shape.transform.rotation != 0:
            pptx_shape.rotation = shape.transform.rotation

        # Apply fill
        self.style_renderer.apply_fill(pptx_shape, shape.fill, theme)

        # Apply stroke
        if shape.stroke:
            self.style_renderer.apply_stroke(pptx_shape, shape.stroke, theme)
        else:
            # No stroke
            pptx_shape.line.fill.background()

        # Apply effects
        self.style_renderer.apply_effects(pptx_shape, shape.effects)

        # Apply text if present
        if shape.text and shape.text.runs:
            self.text_renderer.render(pptx_shape, shape.text, theme)

    def _render_text_box(
        self,
        slide: Slide,
        shape: Shape,
        theme: ThemeColors,
    ) -> None:
        """Render a text box.

        Args:
            slide: The PowerPoint slide.
            shape: The DSL shape.
            theme: Theme colors.
        """
        # Create text box
        text_box = slide.shapes.add_textbox(
            Emu(shape.bbox.x),
            Emu(shape.bbox.y),
            Emu(shape.bbox.width),
            Emu(shape.bbox.height),
        )

        # Apply rotation
        if shape.transform.rotation != 0:
            text_box.rotation = shape.transform.rotation

        # Render text content
        if shape.text and shape.text.runs:
            self.text_renderer.render(text_box, shape.text, theme)

    def _render_image(self, slide: Slide, shape: Shape) -> None:
        """Render an image.

        Args:
            slide: The PowerPoint slide.
            shape: The DSL shape.
        """
        if not shape.image_path:
            return

        try:
            slide.shapes.add_picture(
                shape.image_path,
                Emu(shape.bbox.x),
                Emu(shape.bbox.y),
                Emu(shape.bbox.width),
                Emu(shape.bbox.height),
            )
        except Exception:
            # If image can't be loaded, create a placeholder rectangle
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(shape.bbox.x),
                Emu(shape.bbox.y),
                Emu(shape.bbox.width),
                Emu(shape.bbox.height),
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = self.style_renderer._parse_color("#CCCCCC")

    def _render_group(
        self,
        slide: Slide,
        shape: Shape,
        theme: ThemeColors,
    ) -> None:
        """Render a group and its children.

        Args:
            slide: The PowerPoint slide.
            shape: The DSL group shape.
            theme: Theme colors.
        """
        # PowerPoint groups are complex; for now, render children directly
        if shape.children:
            for child in shape.children:
                self.render(slide, child, theme)

    def _render_freeform(
        self,
        slide: Slide,
        shape: Shape,
        theme: ThemeColors,
    ) -> None:
        """Render a freeform path shape.

        Args:
            slide: The PowerPoint slide.
            shape: The DSL freeform shape.
            theme: Theme colors.
        """
        if not shape.path:
            return

        # Build freeform from path commands
        builder = slide.shapes.build_freeform(
            Emu(shape.bbox.x),
            Emu(shape.bbox.y),
        )

        for cmd in shape.path:
            if cmd.type.value == "moveTo" and cmd.x is not None and cmd.y is not None:
                # Move is implicit in build_freeform start
                pass
            elif cmd.type.value == "lineTo" and cmd.x is not None and cmd.y is not None:
                builder.add_line_segments(
                    [(Emu(cmd.x), Emu(cmd.y))],
                    close=False,
                )
            elif cmd.type.value == "close":
                pass  # Close handled at end

        try:
            freeform = builder.convert_to_shape(
                Emu(shape.bbox.x),
                Emu(shape.bbox.y),
            )

            # Apply fill
            self.style_renderer.apply_fill(freeform, shape.fill, theme)

            # Apply stroke
            if shape.stroke:
                self.style_renderer.apply_stroke(freeform, shape.stroke, theme)

        except Exception:
            # Fallback to rectangle if freeform fails
            self._render_auto_shape(slide, shape, theme)

    def _render_connector(
        self,
        slide: Slide,
        shape: Shape,
        theme: ThemeColors,
    ) -> None:
        """Render a connector line.

        Args:
            slide: The PowerPoint slide.
            shape: The DSL connector shape.
            theme: Theme colors.
        """
        # Create a simple line connector
        connector = slide.shapes.add_connector(
            1,  # Straight connector
            Emu(shape.bbox.x),
            Emu(shape.bbox.y),
            Emu(shape.bbox.x + shape.bbox.width),
            Emu(shape.bbox.y + shape.bbox.height),
        )

        # Apply stroke
        if shape.stroke:
            self.style_renderer.apply_stroke(connector, shape.stroke, theme)
