"""Render text content to PowerPoint shapes."""

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from backend.dsl.schema import TextContent, TextRun, ThemeColors


# Map DSL alignment to PowerPoint
ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

VERTICAL_ALIGN_MAP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


class TextRenderer:
    """Renders text content to PowerPoint shapes."""

    def render(
        self,
        pptx_shape: Any,
        text_content: TextContent,
        theme: ThemeColors,
    ) -> None:
        """Render text content to a shape.

        Args:
            pptx_shape: The python-pptx shape object.
            text_content: The DSL text content.
            theme: Theme colors for resolving references.
        """
        if not text_content.runs:
            return

        # Get text frame
        if not hasattr(pptx_shape, "text_frame"):
            return

        text_frame = pptx_shape.text_frame

        # Set text frame properties
        text_frame.word_wrap = text_content.word_wrap

        # Set margins
        text_frame.margin_left = Emu(text_content.margin_left)
        text_frame.margin_right = Emu(text_content.margin_right)
        text_frame.margin_top = Emu(text_content.margin_top)
        text_frame.margin_bottom = Emu(text_content.margin_bottom)

        # Set vertical alignment
        vertical_align = VERTICAL_ALIGN_MAP.get(
            text_content.vertical_alignment,
            MSO_ANCHOR.MIDDLE,
        )
        text_frame.anchor = vertical_align

        # Set auto-fit
        if text_content.auto_fit == "shrink":
            text_frame.auto_size = True
        elif text_content.auto_fit == "shape":
            # Resize shape to fit text - not directly supported
            pass

        # Clear existing paragraphs
        for paragraph in text_frame.paragraphs:
            paragraph.clear()

        # Get first paragraph
        if text_frame.paragraphs:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()

        # Set horizontal alignment
        paragraph.alignment = ALIGN_MAP.get(text_content.alignment, PP_ALIGN.LEFT)

        # Add text runs
        first_run = True
        for text_run in text_content.runs:
            if first_run:
                # Use existing run in first paragraph
                if paragraph.runs:
                    run = paragraph.runs[0]
                else:
                    run = paragraph.add_run()
                first_run = False
            else:
                run = paragraph.add_run()

            self._apply_run_formatting(run, text_run, theme)

    def _apply_run_formatting(
        self,
        run: Any,
        text_run: TextRun,
        theme: ThemeColors,
    ) -> None:
        """Apply formatting to a text run.

        Args:
            run: The python-pptx Run object.
            text_run: The DSL text run.
            theme: Theme colors.
        """
        # Set text
        run.text = text_run.text

        # Get font
        font = run.font

        # Set font family
        font.name = text_run.font_family

        # Set font size (convert from hundredths of point to Pt)
        font.size = Pt(text_run.font_size / 100)

        # Set bold
        font.bold = text_run.bold

        # Set italic
        font.italic = text_run.italic

        # Set underline
        font.underline = text_run.underline

        # Set color
        color = self._resolve_color(text_run.color, theme)
        font.color.rgb = color

    def _resolve_color(self, color: str, theme: ThemeColors) -> RGBColor:
        """Resolve a color string to RGBColor.

        Args:
            color: Color string (hex or theme reference).
            theme: Theme colors for resolving references.

        Returns:
            RGBColor object.
        """
        color_lower = color.lower()

        # Check for theme references
        if color_lower.startswith("accent"):
            try:
                accent_num = int(color_lower.replace("accent", "").replace("_", ""))
                theme_colors = {
                    1: theme.accent1,
                    2: theme.accent2,
                    3: theme.accent3,
                    4: theme.accent4,
                    5: theme.accent5,
                    6: theme.accent6,
                }
                color = theme_colors.get(accent_num, theme.accent1)
            except ValueError:
                pass
        elif color_lower in ("dark1", "dk1"):
            color = theme.dark1
        elif color_lower in ("light1", "lt1"):
            color = theme.light1

        return self._parse_color(color)

    def _parse_color(self, color: str) -> RGBColor:
        """Parse a hex color string to RGBColor.

        Args:
            color: Hex color string (e.g., '#000000').

        Returns:
            RGBColor object.
        """
        color = color.lstrip("#")
        if len(color) == 6:
            r = int(color[0:2], 16)
            g = int(color[2:4], 16)
            b = int(color[4:6], 16)
            return RGBColor(r, g, b)
        return RGBColor(0, 0, 0)

    def create_text_content(
        self,
        text: str,
        font_family: str = "Calibri",
        font_size: int = 1800,
        bold: bool = False,
        italic: bool = False,
        color: str = "#000000",
        alignment: str = "center",
    ) -> TextContent:
        """Helper to create simple text content.

        Args:
            text: The text string.
            font_family: Font family name.
            font_size: Font size in hundredths of points.
            bold: Whether text is bold.
            italic: Whether text is italic.
            color: Text color.
            alignment: Horizontal alignment.

        Returns:
            TextContent object.
        """
        return TextContent(
            runs=[
                TextRun(
                    text=text,
                    font_family=font_family,
                    font_size=font_size,
                    bold=bold,
                    italic=italic,
                    color=color,
                )
            ],
            alignment=alignment,  # type: ignore
        )
