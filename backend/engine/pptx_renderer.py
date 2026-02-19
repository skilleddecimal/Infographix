"""
pptx_renderer.py — PowerPoint file generation from PositionedLayout.

This renderer consumes PositionedLayout and produces .pptx files.
It NEVER computes positions — that's the layout engine's job.
Uses learned styles from training when available (priority over hardcoded).
It only converts inches to EMU and creates the actual shapes.

CRITICAL RULES:
1. NEVER set shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
2. ALWAYS use pre-computed font sizes from PositionedText
3. ALWAYS disable word wrap for single-line text
4. Use explicit paragraph formatting (no defaults)

VISUAL EFFECTS SUPPORT:
- Shadows (outer, inner, perspective)
- Gradients (linear, radial)
- Multiple shape types (rectangle, rounded rectangle, trapezoid, chevron, etc.)
- 3D effects (bevel, depth)
"""

from pathlib import Path
from typing import Optional, BinaryIO, Union, Dict, Any
from io import BytesIO
from dataclasses import dataclass
from enum import Enum

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

from .positioned import (
    PositionedLayout,
    PositionedElement,
    PositionedConnector,
    PositionedText,
    ElementType,
    ConnectorStyle,
    TextAlignment,
    MultiSlidePresentation,
    RoutingStyle,
    AnchorPosition,
    IconPosition,
)
from .units import (
    inches_to_emu,
    pt_to_emu,
    SLIDE_WIDTH_EMU,
    SLIDE_HEIGHT_EMU,
)

# Learned styles from template analysis
try:
    from .learned_styles import (
        generate_creative_pyramid_style,
        LearnedShadow,
        SHADOW_PYRAMIDA,
    )
    LEARNED_STYLES_AVAILABLE = True
except ImportError:
    LEARNED_STYLES_AVAILABLE = False

# Comprehensive shape learning system (v2)
try:
    from .comprehensive_shape_learning import (
        ComprehensiveLearnings,
        LearnedArchetype,
        LearnedVisualEffects as ComprehensiveVisualEffects,
    )
    from .shape_generator import ShapeGenerator as ComprehensiveShapeGenerator
    from .composition_engine import CompositionEngine, LearnedShapeRenderer
    COMPREHENSIVE_LEARNINGS_AVAILABLE = True
except ImportError:
    COMPREHENSIVE_LEARNINGS_AVAILABLE = False

# Custom shape rendering system
try:
    from .shape_learning import (
        ShapeLibrary,
        ShapeGenerator,
        get_shape_library,
    )
    from .custom_shape_renderer import (
        CustomShapeRenderer,
        render_pyramid_to_slide,
        add_custom_shape_to_slide,
    )
    CUSTOM_SHAPES_AVAILABLE = True
except ImportError:
    CUSTOM_SHAPES_AVAILABLE = False

# Icon library for SVG icon rendering
try:
    from .icon_library import IconLibrary, IconDefinition
    ICON_LIBRARY_AVAILABLE = True
    _icon_library_instance = None

    def get_icon_library() -> IconLibrary:
        """Get singleton icon library instance."""
        global _icon_library_instance
        if _icon_library_instance is None:
            _icon_library_instance = IconLibrary()
        return _icon_library_instance
except ImportError:
    ICON_LIBRARY_AVAILABLE = False

    def get_icon_library():
        return None

# Visual effects configurations
try:
    from .shape_effects import (
        GlowConfig,
        ReflectionConfig,
        SoftEdgeConfig,
        VisualEffects as ShapeVisualEffects,
    )
    VISUAL_EFFECTS_AVAILABLE = True
except ImportError:
    VISUAL_EFFECTS_AVAILABLE = False


# =============================================================================
# VISUAL EFFECT TYPES
# =============================================================================

class ShapeType(Enum):
    """Available shape types for elements."""
    RECTANGLE = "rectangle"
    ROUNDED_RECTANGLE = "rounded_rectangle"
    TRAPEZOID = "trapezoid"
    CHEVRON = "chevron"
    HEXAGON = "hexagon"
    PARALLELOGRAM = "parallelogram"
    OVAL = "oval"
    DIAMOND = "diamond"
    PENTAGON = "pentagon"
    OCTAGON = "octagon"
    ARROW = "arrow"
    CALLOUT = "callout"


# Map ShapeType to MSO_SHAPE constants
SHAPE_TYPE_MAP = {
    ShapeType.RECTANGLE: MSO_SHAPE.RECTANGLE,
    ShapeType.ROUNDED_RECTANGLE: MSO_SHAPE.ROUNDED_RECTANGLE,
    ShapeType.TRAPEZOID: MSO_SHAPE.TRAPEZOID,
    ShapeType.CHEVRON: MSO_SHAPE.CHEVRON,
    ShapeType.HEXAGON: MSO_SHAPE.HEXAGON,
    ShapeType.PARALLELOGRAM: MSO_SHAPE.PARALLELOGRAM,
    ShapeType.OVAL: MSO_SHAPE.OVAL,
    ShapeType.DIAMOND: MSO_SHAPE.DIAMOND,
    ShapeType.PENTAGON: MSO_SHAPE.PENTAGON,
    ShapeType.OCTAGON: MSO_SHAPE.OCTAGON,
    ShapeType.ARROW: MSO_SHAPE.RIGHT_ARROW,
    ShapeType.CALLOUT: MSO_SHAPE.RECTANGULAR_CALLOUT,
}


@dataclass
class ShadowEffect:
    """Shadow effect configuration."""
    enabled: bool = True
    blur_radius_pt: float = 4.0  # Blur radius in points
    distance_pt: float = 3.0     # Distance from shape
    direction_deg: float = 45.0  # Angle in degrees (0=right, 90=down)
    color: str = "#000000"       # Shadow color
    opacity: float = 0.4         # 0-1, transparency
    shadow_type: str = "outer"   # outer, inner, perspective


@dataclass
class GradientEffect:
    """Gradient fill effect configuration."""
    enabled: bool = True
    gradient_type: str = "linear"  # linear, radial, rectangular, path
    angle_deg: float = 270.0       # For linear gradients (270=top to bottom)
    stops: list = None             # List of (position, color) tuples

    def __post_init__(self):
        if self.stops is None:
            # Default: darker at bottom for 3D effect
            self.stops = [
                (0.0, None),    # Will be filled with base color lightened
                (1.0, None),    # Will be filled with base color darkened
            ]


@dataclass
class VisualStyle:
    """Complete visual style for an element."""
    shape_type: ShapeType = ShapeType.ROUNDED_RECTANGLE
    shadow: Optional[ShadowEffect] = None
    gradient: Optional[GradientEffect] = None
    bevel_enabled: bool = False
    bevel_type: str = "circle"    # circle, relaxedInset, softRound, etc.
    bevel_width_pt: float = 3.0
    bevel_height_pt: float = 3.0
    reflection_enabled: bool = False
    glow_enabled: bool = False
    glow_color: str = "#FFFFFF"
    glow_radius_pt: float = 4.0


# =============================================================================
# COLOR HELPERS
# =============================================================================

def hex_to_rgb_color(hex_color: str) -> RGBColor:
    """Convert hex color string to pptx RGBColor."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return RGBColor(0x33, 0x33, 0x33)  # Default dark gray

    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        return RGBColor(0x33, 0x33, 0x33)


def hex_to_rgb_tuple(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple (0-255)."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return (51, 51, 51)
    try:
        return (
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    except ValueError:
        return (51, 51, 51)


def rgb_to_hex(r: int, g: int, b: int) -> str:
    """Convert RGB values to hex string."""
    return f"#{r:02X}{g:02X}{b:02X}"


def lighten_color(hex_color: str, factor: float = 0.2) -> str:
    """Lighten a color by a factor (0-1)."""
    r, g, b = hex_to_rgb_tuple(hex_color)
    r = min(255, int(r + (255 - r) * factor))
    g = min(255, int(g + (255 - g) * factor))
    b = min(255, int(b + (255 - b) * factor))
    return rgb_to_hex(r, g, b)


def darken_color(hex_color: str, factor: float = 0.2) -> str:
    """Darken a color by a factor (0-1)."""
    r, g, b = hex_to_rgb_tuple(hex_color)
    r = max(0, int(r * (1 - factor)))
    g = max(0, int(g * (1 - factor)))
    b = max(0, int(b * (1 - factor)))
    return rgb_to_hex(r, g, b)


# =============================================================================
# PRESET VISUAL STYLES
# =============================================================================

# Professional shadow preset
SHADOW_SUBTLE = ShadowEffect(
    enabled=True,
    blur_radius_pt=6.0,
    distance_pt=3.0,
    direction_deg=135.0,
    color="#000000",
    opacity=0.25,
    shadow_type="outer"
)

SHADOW_MEDIUM = ShadowEffect(
    enabled=True,
    blur_radius_pt=8.0,
    distance_pt=4.0,
    direction_deg=135.0,
    color="#000000",
    opacity=0.35,
    shadow_type="outer"
)

SHADOW_STRONG = ShadowEffect(
    enabled=True,
    blur_radius_pt=12.0,
    distance_pt=6.0,
    direction_deg=135.0,
    color="#000000",
    opacity=0.45,
    shadow_type="outer"
)

# Professional shadow preset (learned from template analysis)
# Key insight: professional templates use SMALLER blur (8-15pt not 40pt)
# and more vertical direction (67-90 degrees) for natural lighting
SHADOW_PYRAMIDA = ShadowEffect(
    enabled=True,
    blur_radius_pt=15.0,    # Professional: 8-15pt for sharp edges
    distance_pt=5.0,        # Subtle offset
    direction_deg=90.0,     # Straight down (natural top lighting)
    color="#000000",
    opacity=0.45,           # 45% opacity for visible but subtle shadow
    shadow_type="outer"
)

# Gradient presets
# Professional gradient with HIGHLIGHT effect (key learning from templates)
# The 15% stop position creates a "shine" effect at the top that makes
# shapes look 3D and polished - this is what separates professional from amateur
GRADIENT_SUBTLE_3D = GradientEffect(
    enabled=True,
    gradient_type="linear",
    angle_deg=270.0,  # Top to bottom
    stops=[
        (0.0, None),    # Lighter at top (auto-computed)
        (0.15, None),   # KEY: 15% stop creates highlight/shine effect
        (1.0, None)     # Darker at bottom (auto-computed)
    ]
)

GRADIENT_GLASS = GradientEffect(
    enabled=True,
    gradient_type="linear",
    angle_deg=270.0,
    stops=[
        (0.0, None),     # Very light at top
        (0.10, None),    # Quick transition for glass effect
        (0.5, None),     # Base color in middle
        (1.0, None)      # Darker at bottom
    ]
)

# =============================================================================
# STYLE FROM LEARNINGS
# =============================================================================

def load_style_from_learnings(archetype: str) -> Optional["VisualStyle"]:
    """
    Load visual style from learnings for a specific archetype.

    Checks in order:
    1. comprehensive_learnings.json (v2 system with full geometry)
    2. shape_learnings.json (v1 system with basic effects)

    Returns None if no learnings are available (triggers basic mode).
    """
    import json
    import os

    # Try comprehensive learnings first (v2)
    comprehensive_path = os.path.join(
        os.path.dirname(__file__), "..", "data", "comprehensive_learnings.json"
    )

    if os.path.exists(comprehensive_path) and COMPREHENSIVE_LEARNINGS_AVAILABLE:
        try:
            learnings = ComprehensiveLearnings.load(comprehensive_path)
            if archetype in learnings.archetypes:
                arch = learnings.archetypes[archetype]
                return _convert_comprehensive_to_visual_style(arch)
        except Exception:
            pass  # Fall through to v1 learnings

    # Fall back to v1 shape_learnings.json
    learnings_path = os.path.join(
        os.path.dirname(__file__), "..", "data", "shape_learnings.json"
    )

    if not os.path.exists(learnings_path):
        return None

    try:
        with open(learnings_path, "r") as f:
            data = json.load(f)
    except (json.JSONDecodeError, IOError):
        return None

    # Check if learnings exist
    shape_learnings = data.get("shape_learnings", {})
    if not shape_learnings:
        return None  # No learnings = basic mode

    # Find the category for this archetype
    category = None
    learning = None
    for cat_name, cat_data in shape_learnings.items():
        if cat_name.startswith("_"):
            continue
        archetype_types = cat_data.get("archetype_types", [])
        if archetype in archetype_types:
            category = cat_name
            learning = cat_data
            break

    if learning is None:
        return None

    # Build VisualStyle from learning
    shadow = None
    if learning.get("shadow_blur_pt", 0) > 0:
        shadow = ShadowEffect(
            enabled=True,
            blur_radius_pt=learning.get("shadow_blur_pt", 12.0),
            distance_pt=learning.get("shadow_distance_pt", 5.0),
            direction_deg=learning.get("shadow_direction_deg", 90.0),
            color="#000000",
            opacity=learning.get("shadow_opacity", 0.4),
            shadow_type="outer"
        )

    gradient = None
    if learning.get("gradient_type"):
        gradient = GradientEffect(
            enabled=True,
            gradient_type=learning.get("gradient_type", "linear"),
            angle_deg=learning.get("gradient_angle", 270.0),
            stops=[
                {"position": 0.0, "color_offset": 0.30},
                {"position": learning.get("gradient_highlight_position", 0.15), "color_offset": 0.0},
                {"position": 1.0, "color_offset": -0.30},
            ]
        )

    bevel_enabled = learning.get("bevel_width_inches", 0) > 0
    bevel_type = learning.get("bevel_type", "softRound")
    bevel_width_pt = learning.get("bevel_width_inches", 0.028) * 72  # Convert inches to points

    # Determine shape type
    shape_type = ShapeType.ROUNDED_RECTANGLE
    if category == "tapered":
        shape_type = ShapeType.TRAPEZOID
    elif learning.get("corner_radius_ratio", 0) < 0.03:
        shape_type = ShapeType.RECTANGLE

    return VisualStyle(
        shape_type=shape_type,
        shadow=shadow,
        gradient=gradient,
        bevel_enabled=bevel_enabled,
        bevel_type=bevel_type,
        bevel_width_pt=bevel_width_pt,
        bevel_height_pt=bevel_width_pt
    )


def _convert_comprehensive_to_visual_style(arch: "LearnedArchetype") -> Optional["VisualStyle"]:
    """Convert a LearnedArchetype from comprehensive system to VisualStyle."""
    if not arch.effects:
        return None

    effects = arch.effects

    # Convert shadow
    shadow = None
    if effects.shadow and effects.shadow.enabled:
        shadow = ShadowEffect(
            enabled=True,
            blur_radius_pt=effects.shadow.blur_radius_pt,
            distance_pt=effects.shadow.distance_pt,
            direction_deg=effects.shadow.direction_deg,
            color=effects.shadow.color,
            opacity=effects.shadow.opacity,
            shadow_type="outer" if effects.shadow.shadow_type.value == "outer" else "inner"
        )

    # Convert gradient
    gradient = None
    if effects.gradient and effects.gradient.enabled:
        # Convert gradient stops
        stops = []
        for stop in effects.gradient.stops:
            pos = stop.position
            color = None  # Will be computed from base
            stops.append((pos, color))

        gradient = GradientEffect(
            enabled=True,
            gradient_type=effects.gradient.gradient_type.value,
            angle_deg=effects.gradient.angle_deg,
            stops=stops if stops else None
        )

    # Convert bevel
    bevel_enabled = effects.bevel and effects.bevel.enabled
    bevel_type = effects.bevel.bevel_type if bevel_enabled else "softRound"
    bevel_width_pt = effects.bevel.width_pt if bevel_enabled else 0.0

    # Determine shape type based on archetype
    shape_type = ShapeType.ROUNDED_RECTANGLE
    if arch.archetype_name in ["funnel", "pyramid"]:
        shape_type = ShapeType.TRAPEZOID
    elif arch.archetype_name in ["chevron", "process_flow"]:
        shape_type = ShapeType.CHEVRON

    return VisualStyle(
        shape_type=shape_type,
        shadow=shadow,
        gradient=gradient,
        bevel_enabled=bevel_enabled,
        bevel_type=bevel_type,
        bevel_width_pt=bevel_width_pt,
        bevel_height_pt=bevel_width_pt
    )


def has_shape_learnings() -> bool:
    """Check if shape learnings exist (for determining basic vs trained mode)."""
    import json
    import os

    learnings_path = os.path.join(
        os.path.dirname(__file__), "..", "data", "shape_learnings.json"
    )

    if not os.path.exists(learnings_path):
        return False

    try:
        with open(learnings_path, "r") as f:
            data = json.load(f)
        shape_learnings = data.get("shape_learnings", {})
        # Exclude metadata keys starting with "_"
        real_learnings = [k for k in shape_learnings.keys() if not k.startswith("_")]
        return len(real_learnings) > 0
    except (json.JSONDecodeError, IOError):
        return False


# =============================================================================
# COMBINED VISUAL STYLE PRESETS
# =============================================================================

# BASIC style - NO effects (used when no learnings are available)
# This shows the "before training" state
STYLE_BASIC = VisualStyle(
    shape_type=ShapeType.ROUNDED_RECTANGLE,
    shadow=None,
    gradient=None,
    bevel_enabled=False
)

# Combined visual style presets
STYLE_FLAT = VisualStyle(
    shape_type=ShapeType.ROUNDED_RECTANGLE,
    shadow=None,
    gradient=None,
    bevel_enabled=False
)

STYLE_SUBTLE_3D = VisualStyle(
    shape_type=ShapeType.ROUNDED_RECTANGLE,
    shadow=SHADOW_SUBTLE,
    gradient=GRADIENT_SUBTLE_3D,
    bevel_enabled=False
)

STYLE_PROFESSIONAL = VisualStyle(
    shape_type=ShapeType.ROUNDED_RECTANGLE,
    shadow=SHADOW_MEDIUM,
    gradient=GRADIENT_SUBTLE_3D,
    bevel_enabled=True,
    bevel_type="softRound",
    bevel_width_pt=2.0,
    bevel_height_pt=2.0
)

STYLE_EXECUTIVE = VisualStyle(
    shape_type=ShapeType.RECTANGLE,
    shadow=SHADOW_STRONG,
    gradient=GRADIENT_GLASS,
    bevel_enabled=True,
    bevel_type="relaxedInset",
    bevel_width_pt=3.0,
    bevel_height_pt=3.0
)

# Pyramid-specific styles (learned from Pyramida template)
# Template uses: soft shadows (40pt blur), no strokes, no 3D bevels
STYLE_PYRAMID_LEVEL = VisualStyle(
    shape_type=ShapeType.TRAPEZOID,
    shadow=SHADOW_PYRAMIDA,       # Soft diffuse shadow from template
    gradient=GRADIENT_SUBTLE_3D,  # Subtle depth gradient
    bevel_enabled=False,          # Template doesn't use 3D bevels
    bevel_type=None,
    bevel_width_pt=0,
    bevel_height_pt=0
)

# Funnel-specific styles (PROFESSIONAL - learned from template analysis)
# Key learnings from professional templates:
# - Shadow blur: 8-15pt (not 40pt) for sharp, professional look
# - Shadow direction: 90 degrees (straight down) simulates natural top light
# - Bevel: softRound at 0.028" width gives subtle 3D embossed look
# - Gradient stops: 15% position for highlight effect (not 50%)

SHADOW_FUNNEL = ShadowEffect(
    enabled=True,
    blur_radius_pt=12.0,      # Professional: 8-15pt (NOT 40pt!)
    distance_pt=5.0,          # Subtle offset
    direction_deg=90.0,       # Straight down (natural light from above)
    color="#000000",
    opacity=0.40,             # 40% opacity for visible but subtle shadow
    shadow_type="outer"
)

# Professional gradient with highlight effect (key learning from templates)
# The 15% stop position creates a "shine" effect at the top
GRADIENT_PROFESSIONAL = GradientEffect(
    gradient_type="linear",
    angle_deg=270.0,
    stops=[
        {"position": 0.0, "color_offset": 0.30},   # Lighter at top
        {"position": 0.15, "color_offset": 0.0},   # KEY: 15% stop creates highlight
        {"position": 1.0, "color_offset": -0.30},  # Darker at bottom
    ]
)

STYLE_FUNNEL_STAGE = VisualStyle(
    shape_type=ShapeType.ROUNDED_RECTANGLE,  # Use rounded rectangles, NOT custom shapes
    shadow=SHADOW_FUNNEL,                     # Professional shadow (12pt blur, 90 deg)
    gradient=GRADIENT_SUBTLE_3D,              # Top-to-bottom gradient for 3D depth
    bevel_enabled=True,                       # Subtle 3D bevel effect
    bevel_type="softRound",                   # Soft rounded bevel (from templates)
    bevel_width_pt=2.5,                       # ~0.028" width (learned from templates)
    bevel_height_pt=2.5
)


def get_alignment(text_align: TextAlignment) -> PP_ALIGN:
    """Convert TextAlignment to pptx PP_ALIGN."""
    mapping = {
        TextAlignment.LEFT: PP_ALIGN.LEFT,
        TextAlignment.CENTER: PP_ALIGN.CENTER,
        TextAlignment.RIGHT: PP_ALIGN.RIGHT,
    }
    return mapping.get(text_align, PP_ALIGN.CENTER)


# =============================================================================
# XML GENERATION FOR VISUAL EFFECTS
# =============================================================================

# XML namespace for DrawingML
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
A_NSMAP = {'a': A_NS}


def _create_shadow_xml(shadow: ShadowEffect) -> etree.Element:
    """
    Create DrawingML XML for shadow effect.

    Returns an <a:effectLst> element with the shadow definition.
    """
    # Convert values to EMU (1 point = 12700 EMU)
    blur_emu = int(shadow.blur_radius_pt * 12700)
    dist_emu = int(shadow.distance_pt * 12700)

    # Direction in PowerPoint: 0 = right, 5400000 = down (60000 units = 1 degree)
    direction = int(shadow.direction_deg * 60000)

    # Opacity: 0 = transparent, 100000 = opaque
    alpha = int(shadow.opacity * 100000)

    # Color without #
    color_hex = shadow.color.lstrip('#')

    if shadow.shadow_type == "outer":
        effect_xml = f'''
        <a:effectLst xmlns:a="{A_NS}">
            <a:outerShdw blurRad="{blur_emu}" dist="{dist_emu}" dir="{direction}" algn="ctr" rotWithShape="0">
                <a:srgbClr val="{color_hex}">
                    <a:alpha val="{alpha}"/>
                </a:srgbClr>
            </a:outerShdw>
        </a:effectLst>
        '''
    elif shadow.shadow_type == "inner":
        effect_xml = f'''
        <a:effectLst xmlns:a="{A_NS}">
            <a:innerShdw blurRad="{blur_emu}" dist="{dist_emu}" dir="{direction}">
                <a:srgbClr val="{color_hex}">
                    <a:alpha val="{alpha}"/>
                </a:srgbClr>
            </a:innerShdw>
        </a:effectLst>
        '''
    elif shadow.shadow_type == "perspective":
        # Perspective shadow with preset
        effect_xml = f'''
        <a:effectLst xmlns:a="{A_NS}">
            <a:outerShdw blurRad="{blur_emu}" dist="{dist_emu}" dir="{direction}"
                         sx="100000" sy="100000" kx="0" ky="0" algn="b" rotWithShape="0">
                <a:srgbClr val="{color_hex}">
                    <a:alpha val="{alpha}"/>
                </a:srgbClr>
            </a:outerShdw>
        </a:effectLst>
        '''
    else:
        effect_xml = f'<a:effectLst xmlns:a="{A_NS}"/>'

    return etree.fromstring(effect_xml)


def _create_gradient_fill_xml(gradient: GradientEffect, base_color: str) -> etree.Element:
    """
    Create DrawingML XML for gradient fill.

    Returns an <a:gradFill> element.
    """
    # Angle in DrawingML: 0 = horizontal (left to right), 16200000 = 270 degrees
    # angle_val = angle_deg * 60000
    angle_val = int(gradient.angle_deg * 60000)

    # Build gradient stops
    stops_xml = ""
    for i, (pos, color) in enumerate(gradient.stops):
        # Position as percentage * 1000 (0-100000)
        pos_val = int(pos * 100000)

        # If color is None, auto-compute from base color
        if color is None:
            if len(gradient.stops) == 2:
                # Two stops: lighter at top, base at bottom
                if i == 0:
                    color = lighten_color(base_color, 0.25)
                else:
                    color = darken_color(base_color, 0.15)
            elif len(gradient.stops) == 3:
                # Three stops: glass effect
                if i == 0:
                    color = lighten_color(base_color, 0.35)
                elif i == 1:
                    color = base_color
                else:
                    color = darken_color(base_color, 0.20)
            else:
                color = base_color

        color_hex = color.lstrip('#')
        stops_xml += f'''
            <a:gs pos="{pos_val}">
                <a:srgbClr val="{color_hex}"/>
            </a:gs>
        '''

    if gradient.gradient_type == "linear":
        fill_xml = f'''
        <a:gradFill xmlns:a="{A_NS}" rotWithShape="1">
            <a:gsLst>
                {stops_xml}
            </a:gsLst>
            <a:lin ang="{angle_val}" scaled="1"/>
        </a:gradFill>
        '''
    elif gradient.gradient_type == "radial":
        fill_xml = f'''
        <a:gradFill xmlns:a="{A_NS}" rotWithShape="1">
            <a:gsLst>
                {stops_xml}
            </a:gsLst>
            <a:path path="circle">
                <a:fillToRect l="50000" t="50000" r="50000" b="50000"/>
            </a:path>
        </a:gradFill>
        '''
    else:
        # Default to linear
        fill_xml = f'''
        <a:gradFill xmlns:a="{A_NS}" rotWithShape="1">
            <a:gsLst>
                {stops_xml}
            </a:gsLst>
            <a:lin ang="{angle_val}" scaled="1"/>
        </a:gradFill>
        '''

    return etree.fromstring(fill_xml)


def _create_bevel_xml(bevel_type: str, width_pt: float, height_pt: float) -> etree.Element:
    """
    Create DrawingML XML for 3D bevel effect.

    Returns an <a:sp3d> element.
    """
    # Convert points to EMU
    width_emu = int(width_pt * 12700)
    height_emu = int(height_pt * 12700)

    sp3d_xml = f'''
    <a:sp3d xmlns:a="{A_NS}">
        <a:bevelT w="{width_emu}" h="{height_emu}" prst="{bevel_type}"/>
    </a:sp3d>
    '''

    return etree.fromstring(sp3d_xml)


def _create_glow_xml(glow_config) -> etree.Element:
    """
    Create DrawingML XML for outer glow effect.

    The glow effect creates a soft colored glow around the shape edge.

    Returns an <a:glow> element to be added to <a:effectLst>.
    """
    # Convert radius to EMU (1 point = 12700 EMU)
    radius_emu = int(glow_config.radius_pt * 12700)

    # Alpha: 0 = transparent, 100000 = opaque
    alpha = int((1.0 - glow_config.transparency) * 100000)

    # Color without #
    color_hex = glow_config.color.lstrip('#')

    glow_xml = f'''
    <a:glow xmlns:a="{A_NS}" rad="{radius_emu}">
        <a:srgbClr val="{color_hex}">
            <a:alpha val="{alpha}"/>
        </a:srgbClr>
    </a:glow>
    '''

    return etree.fromstring(glow_xml)


def _create_reflection_xml(reflection_config) -> etree.Element:
    """
    Create DrawingML XML for reflection effect.

    The reflection effect creates a mirror image below the shape.

    Returns an <a:reflection> element to be added to <a:effectLst>.
    """
    # Convert blur radius to EMU
    blur_emu = int(reflection_config.blur_radius_pt * 12700)

    # Start and end alpha (opacity * 100000)
    start_alpha = int(reflection_config.start_opacity * 100000)
    end_alpha = int(reflection_config.end_opacity * 100000)

    # Distance in EMU
    distance_emu = int(reflection_config.distance_pt * 12700)

    # Direction: 5400000 = 90 degrees in DrawingML (60000 per degree)
    direction = int(reflection_config.direction_degrees * 60000)

    # Scale Y as percentage (100000 = 100%)
    scale_y = int(reflection_config.scale_y * 100000)

    reflection_xml = f'''
    <a:reflection xmlns:a="{A_NS}" blurRad="{blur_emu}" stA="{start_alpha}" endA="{end_alpha}"
                  dist="{distance_emu}" dir="{direction}" sy="{scale_y}"
                  algn="bl" rotWithShape="0"/>
    '''

    return etree.fromstring(reflection_xml)


def _create_soft_edge_xml(soft_edge_config) -> etree.Element:
    """
    Create DrawingML XML for soft edge effect.

    The soft edge effect blurs the edges of the shape.

    Returns an <a:softEdge> element to be added to <a:effectLst>.
    """
    # Convert radius to EMU
    radius_emu = int(soft_edge_config.radius_pt * 12700)

    soft_edge_xml = f'''
    <a:softEdge xmlns:a="{A_NS}" rad="{radius_emu}"/>
    '''

    return etree.fromstring(soft_edge_xml)


def _create_icon_custgeom_xml(svg_path: str, view_box: str, width_emu: int, height_emu: int) -> str:
    """
    Convert SVG path data to DrawingML custGeom XML for icon rendering.

    Args:
        svg_path: SVG path 'd' attribute data
        view_box: SVG viewBox (e.g., "0 0 24 24")
        width_emu: Target width in EMU
        height_emu: Target height in EMU

    Returns:
        DrawingML custGeom XML string
    """
    # Parse viewBox
    vb_parts = view_box.split()
    if len(vb_parts) == 4:
        vb_x, vb_y, vb_w, vb_h = map(float, vb_parts)
    else:
        vb_x, vb_y, vb_w, vb_h = 0, 0, 24, 24

    # Scale factors
    scale_x = width_emu / vb_w
    scale_y = height_emu / vb_h

    # Convert SVG path to DrawingML path elements
    path_elements = _svg_path_to_drawingml(svg_path, scale_x, scale_y, vb_x, vb_y)

    custgeom_xml = f'''
    <a:custGeom xmlns:a="{A_NS}">
        <a:avLst/>
        <a:gdLst/>
        <a:ahLst/>
        <a:cxnLst/>
        <a:rect l="0" t="0" r="{width_emu}" b="{height_emu}"/>
        <a:pathLst>
            <a:path w="{width_emu}" h="{height_emu}">
                {path_elements}
            </a:path>
        </a:pathLst>
    </a:custGeom>
    '''

    return custgeom_xml


def _svg_path_to_drawingml(svg_path: str, scale_x: float, scale_y: float,
                            offset_x: float = 0, offset_y: float = 0) -> str:
    """
    Convert SVG path data to DrawingML path elements.

    Supports: M (moveTo), L (lineTo), H (horizontal), V (vertical),
              C (cubicBezTo), Q (quadBezTo), A (arcTo), Z (close)

    Args:
        svg_path: SVG path 'd' attribute
        scale_x: X scale factor
        scale_y: Y scale factor
        offset_x: X offset from viewBox
        offset_y: Y offset from viewBox

    Returns:
        DrawingML path elements as XML string
    """
    import re

    elements = []
    current_x = 0.0
    current_y = 0.0

    # Split path into commands with their parameters
    # This regex captures command letter followed by coordinates
    path_re = re.compile(r'([MmLlHhVvCcSsQqTtAaZz])\s*([^MmLlHhVvCcSsQqTtAaZz]*)')

    for match in path_re.finditer(svg_path):
        cmd = match.group(1)
        params_str = match.group(2).strip()

        # Parse numbers from params
        if params_str:
            nums = re.findall(r'-?[\d.]+(?:e[+-]?\d+)?', params_str)
            params = [float(n) for n in nums]
        else:
            params = []

        # Handle each command
        is_relative = cmd.islower()
        cmd_upper = cmd.upper()

        if cmd_upper == 'M':
            # MoveTo
            if len(params) >= 2:
                x, y = params[0], params[1]
                if is_relative:
                    x += current_x
                    y += current_y
                current_x, current_y = x, y
                scaled_x = int((x - offset_x) * scale_x)
                scaled_y = int((y - offset_y) * scale_y)
                elements.append(f'<a:moveTo><a:pt x="{scaled_x}" y="{scaled_y}"/></a:moveTo>')

                # Additional coordinates are lineTo
                for i in range(2, len(params), 2):
                    if i + 1 < len(params):
                        x, y = params[i], params[i + 1]
                        if is_relative:
                            x += current_x
                            y += current_y
                        current_x, current_y = x, y
                        scaled_x = int((x - offset_x) * scale_x)
                        scaled_y = int((y - offset_y) * scale_y)
                        elements.append(f'<a:lnTo><a:pt x="{scaled_x}" y="{scaled_y}"/></a:lnTo>')

        elif cmd_upper == 'L':
            # LineTo
            for i in range(0, len(params), 2):
                if i + 1 < len(params):
                    x, y = params[i], params[i + 1]
                    if is_relative:
                        x += current_x
                        y += current_y
                    current_x, current_y = x, y
                    scaled_x = int((x - offset_x) * scale_x)
                    scaled_y = int((y - offset_y) * scale_y)
                    elements.append(f'<a:lnTo><a:pt x="{scaled_x}" y="{scaled_y}"/></a:lnTo>')

        elif cmd_upper == 'H':
            # Horizontal LineTo
            for x in params:
                if is_relative:
                    x += current_x
                current_x = x
                scaled_x = int((x - offset_x) * scale_x)
                scaled_y = int((current_y - offset_y) * scale_y)
                elements.append(f'<a:lnTo><a:pt x="{scaled_x}" y="{scaled_y}"/></a:lnTo>')

        elif cmd_upper == 'V':
            # Vertical LineTo
            for y in params:
                if is_relative:
                    y += current_y
                current_y = y
                scaled_x = int((current_x - offset_x) * scale_x)
                scaled_y = int((y - offset_y) * scale_y)
                elements.append(f'<a:lnTo><a:pt x="{scaled_x}" y="{scaled_y}"/></a:lnTo>')

        elif cmd_upper == 'C':
            # Cubic Bezier
            for i in range(0, len(params), 6):
                if i + 5 < len(params):
                    x1, y1, x2, y2, x, y = params[i:i + 6]
                    if is_relative:
                        x1 += current_x
                        y1 += current_y
                        x2 += current_x
                        y2 += current_y
                        x += current_x
                        y += current_y
                    current_x, current_y = x, y

                    sx1 = int((x1 - offset_x) * scale_x)
                    sy1 = int((y1 - offset_y) * scale_y)
                    sx2 = int((x2 - offset_x) * scale_x)
                    sy2 = int((y2 - offset_y) * scale_y)
                    sx = int((x - offset_x) * scale_x)
                    sy = int((y - offset_y) * scale_y)

                    elements.append(
                        f'<a:cubicBezTo>'
                        f'<a:pt x="{sx1}" y="{sy1}"/>'
                        f'<a:pt x="{sx2}" y="{sy2}"/>'
                        f'<a:pt x="{sx}" y="{sy}"/>'
                        f'</a:cubicBezTo>'
                    )

        elif cmd_upper == 'Q':
            # Quadratic Bezier
            for i in range(0, len(params), 4):
                if i + 3 < len(params):
                    x1, y1, x, y = params[i:i + 4]
                    if is_relative:
                        x1 += current_x
                        y1 += current_y
                        x += current_x
                        y += current_y
                    current_x, current_y = x, y

                    sx1 = int((x1 - offset_x) * scale_x)
                    sy1 = int((y1 - offset_y) * scale_y)
                    sx = int((x - offset_x) * scale_x)
                    sy = int((y - offset_y) * scale_y)

                    elements.append(
                        f'<a:quadBezTo>'
                        f'<a:pt x="{sx1}" y="{sy1}"/>'
                        f'<a:pt x="{sx}" y="{sy}"/>'
                        f'</a:quadBezTo>'
                    )

        elif cmd_upper == 'A':
            # Arc - convert to line approximation (simplification)
            for i in range(0, len(params), 7):
                if i + 6 < len(params):
                    x, y = params[i + 5], params[i + 6]
                    if is_relative:
                        x += current_x
                        y += current_y
                    current_x, current_y = x, y
                    scaled_x = int((x - offset_x) * scale_x)
                    scaled_y = int((y - offset_y) * scale_y)
                    elements.append(f'<a:lnTo><a:pt x="{scaled_x}" y="{scaled_y}"/></a:lnTo>')

        elif cmd_upper == 'Z':
            # Close path
            elements.append('<a:close/>')

    return '\n                '.join(elements)


def apply_visual_effects_to_shape(shape, visual_style: VisualStyle, base_color: str) -> None:
    """
    Apply visual effects (shadow, gradient, bevel) to a PowerPoint shape.

    This manipulates the underlying XML of the shape to add effects
    that python-pptx doesn't directly support.
    """
    try:
        spPr = shape._element.spPr

        # Apply gradient fill
        if visual_style.gradient and visual_style.gradient.enabled:
            # Remove existing solidFill if present
            existing_fill = spPr.find(qn('a:solidFill'))
            if existing_fill is not None:
                spPr.remove(existing_fill)

            # Add gradient fill
            grad_fill = _create_gradient_fill_xml(visual_style.gradient, base_color)

            # Insert after prstGeom/custGeom (OOXML requires fill after geometry)
            # Order: xfrm -> prstGeom/custGeom -> fill -> ln -> effectLst -> sp3d
            prstGeom = spPr.find(qn('a:prstGeom'))
            custGeom = spPr.find(qn('a:custGeom'))
            geom = prstGeom if prstGeom is not None else custGeom

            if geom is not None:
                geom_index = list(spPr).index(geom)
                spPr.insert(geom_index + 1, grad_fill)
            else:
                # Fallback: insert after xfrm
                xfrm = spPr.find(qn('a:xfrm'))
                if xfrm is not None:
                    xfrm_index = list(spPr).index(xfrm)
                    spPr.insert(xfrm_index + 1, grad_fill)
                else:
                    spPr.insert(0, grad_fill)

        # Apply shadow effect
        if visual_style.shadow and visual_style.shadow.enabled:
            # Remove existing effectLst
            existing_effects = spPr.find(qn('a:effectLst'))
            if existing_effects is not None:
                spPr.remove(existing_effects)

            # Add shadow
            effect_lst = _create_shadow_xml(visual_style.shadow)
            spPr.append(effect_lst)

        # Apply 3D bevel
        if visual_style.bevel_enabled:
            # Remove existing sp3d
            existing_3d = spPr.find(qn('a:sp3d'))
            if existing_3d is not None:
                spPr.remove(existing_3d)

            # Add bevel
            sp3d = _create_bevel_xml(
                visual_style.bevel_type,
                visual_style.bevel_width_pt,
                visual_style.bevel_height_pt
            )
            spPr.append(sp3d)

    except Exception as e:
        # Silently fail - effects are optional enhancements
        pass


def apply_extended_effects_to_shape(
    shape,
    glow: Optional[Any] = None,
    reflection: Optional[Any] = None,
    soft_edge: Optional[Any] = None
) -> None:
    """
    Apply extended visual effects (glow, reflection, soft edge) to a PowerPoint shape.

    These effects are added to the effectLst XML element.

    Args:
        shape: PowerPoint shape to apply effects to
        glow: GlowConfig instance or None
        reflection: ReflectionConfig instance or None
        soft_edge: SoftEdgeConfig instance or None
    """
    if not VISUAL_EFFECTS_AVAILABLE:
        return

    if glow is None and reflection is None and soft_edge is None:
        return

    try:
        spPr = shape._element.spPr

        # Get or create effectLst
        effect_lst = spPr.find(qn('a:effectLst'))
        if effect_lst is None:
            effect_lst = etree.Element(qn('a:effectLst'))
            spPr.append(effect_lst)

        # Apply glow effect
        if glow is not None and glow.enabled:
            # Remove existing glow
            existing_glow = effect_lst.find(qn('a:glow'))
            if existing_glow is not None:
                effect_lst.remove(existing_glow)

            glow_elem = _create_glow_xml(glow)
            effect_lst.append(glow_elem)

        # Apply soft edge effect
        if soft_edge is not None and soft_edge.enabled:
            # Remove existing softEdge
            existing_soft = effect_lst.find(qn('a:softEdge'))
            if existing_soft is not None:
                effect_lst.remove(existing_soft)

            soft_edge_elem = _create_soft_edge_xml(soft_edge)
            effect_lst.append(soft_edge_elem)

        # Apply reflection effect
        if reflection is not None and reflection.enabled:
            # Remove existing reflection
            existing_reflection = effect_lst.find(qn('a:reflection'))
            if existing_reflection is not None:
                effect_lst.remove(existing_reflection)

            reflection_elem = _create_reflection_xml(reflection)
            effect_lst.append(reflection_elem)

    except Exception as e:
        # Silently fail - effects are optional enhancements
        pass


def apply_effects_from_visual_effects(shape, effects: Any, base_color: str = "#4472C4") -> None:
    """
    Apply effects from a VisualEffects (shape_effects.py) instance.

    This is a convenience function that handles both basic effects (shadow, gradient, bevel)
    and extended effects (glow, reflection, soft_edge).

    Args:
        shape: PowerPoint shape to apply effects to
        effects: VisualEffects instance from shape_effects.py
        base_color: Base color for gradient computation
    """
    if not VISUAL_EFFECTS_AVAILABLE:
        return

    if effects is None:
        return

    try:
        spPr = shape._element.spPr

        # Apply gradient fill
        if effects.gradient is not None:
            # Create gradient stops from effects.gradient
            gradient_effect = GradientEffect(
                enabled=True,
                gradient_type=effects.gradient.type.value if hasattr(effects.gradient.type, 'value') else "linear",
                angle_deg=effects.gradient.angle,
                stops=[(s.position, s.color) for s in effects.gradient.stops] if effects.gradient.stops else None
            )

            # Remove existing solidFill
            existing_fill = spPr.find(qn('a:solidFill'))
            if existing_fill is not None:
                spPr.remove(existing_fill)

            grad_fill = _create_gradient_fill_xml(gradient_effect, base_color)

            # Insert after geometry
            prstGeom = spPr.find(qn('a:prstGeom'))
            custGeom = spPr.find(qn('a:custGeom'))
            geom = prstGeom if prstGeom is not None else custGeom

            if geom is not None:
                geom_index = list(spPr).index(geom)
                spPr.insert(geom_index + 1, grad_fill)
            else:
                xfrm = spPr.find(qn('a:xfrm'))
                if xfrm is not None:
                    xfrm_index = list(spPr).index(xfrm)
                    spPr.insert(xfrm_index + 1, grad_fill)

        # Build effectLst with all effects
        effect_lst = spPr.find(qn('a:effectLst'))
        if effect_lst is None:
            effect_lst = etree.Element(qn('a:effectLst'))
            spPr.append(effect_lst)
        else:
            # Clear existing effects
            for child in list(effect_lst):
                effect_lst.remove(child)

        # Apply shadow
        if effects.shadow is not None and effects.shadow.enabled:
            shadow_effect = ShadowEffect(
                enabled=True,
                blur_radius_pt=effects.shadow.blur_radius_pt,
                distance_pt=effects.shadow.distance_inches * 72,  # Convert inches to points
                direction_deg=effects.shadow.direction_degrees,
                color=effects.shadow.color,
                opacity=1.0 - effects.shadow.transparency,
                shadow_type="outer"
            )
            shadow_elem = _create_shadow_xml(shadow_effect)
            # Extract the inner shadow element (outerShdw)
            for child in shadow_elem:
                effect_lst.append(child)

        # Apply glow
        if effects.glow is not None and effects.glow.enabled:
            glow_elem = _create_glow_xml(effects.glow)
            effect_lst.append(glow_elem)

        # Apply soft edge
        if effects.soft_edge is not None and effects.soft_edge.enabled:
            soft_edge_elem = _create_soft_edge_xml(effects.soft_edge)
            effect_lst.append(soft_edge_elem)

        # Apply reflection
        if effects.reflection is not None and effects.reflection.enabled:
            reflection_elem = _create_reflection_xml(effects.reflection)
            effect_lst.append(reflection_elem)

        # Apply bevel
        if effects.bevel is not None and effects.bevel.enabled:
            existing_3d = spPr.find(qn('a:sp3d'))
            if existing_3d is not None:
                spPr.remove(existing_3d)

            bevel_type = effects.bevel.type.value if hasattr(effects.bevel.type, 'value') else "softRound"
            width_pt = effects.bevel.width_inches * 72  # Convert inches to points
            height_pt = effects.bevel.height_inches * 72

            sp3d = _create_bevel_xml(bevel_type, width_pt, height_pt)
            spPr.append(sp3d)

    except Exception as e:
        # Silently fail - effects are optional enhancements
        pass


# =============================================================================
# PPTX RENDERER
# =============================================================================

def convert_design_style_to_visual_style(design_style) -> VisualStyle:
    """
    Convert a learned DesignStyle (from design_learner) to a VisualStyle.

    This extracts the visual effects (shadows, gradients, etc.) from a learned
    design style and creates a VisualStyle that can be used by the renderer.
    The colors from the learned style are NOT used - only the visual effects.
    """
    if design_style is None:
        return None

    shape_style = design_style.shape_style
    if shape_style is None:
        return STYLE_PROFESSIONAL

    # Convert shadow
    shadow = None
    if shape_style.shadow:
        shadow = ShadowEffect(
            enabled=shape_style.shadow.enabled,
            blur_radius_pt=shape_style.shadow.blur_radius_pt,
            distance_pt=shape_style.shadow.distance_pt,
            direction_deg=shape_style.shadow.angle_degrees,
            color=shape_style.shadow.color,
            opacity=shape_style.shadow.opacity,
            shadow_type="outer",
        )

    # Convert gradient (but colors will be replaced with user's colors)
    gradient = None
    if shape_style.gradient:
        gradient = GradientEffect(
            enabled=True,
            gradient_type=shape_style.gradient.type.value if hasattr(shape_style.gradient.type, 'value') else str(shape_style.gradient.type),
            angle_deg=shape_style.gradient.angle_degrees,
            stops=None,  # Will be filled with user's colors
        )

        # If gradient is present but no shadow, add a default subtle shadow
        # Professional diagrams with gradients typically have matching shadows
        if shadow is None:
            shadow = ShadowEffect(
                enabled=True,
                blur_radius_pt=4.0,
                distance_pt=2.0,
                direction_deg=135.0,
                color="#000000",
                opacity=0.25,
                shadow_type="outer",
            )

    # Determine shape type based on corner radius
    shape_type = ShapeType.ROUNDED_RECTANGLE
    if shape_style.corner_radius_ratio < 0.03:
        shape_type = ShapeType.RECTANGLE
    elif shape_style.corner_radius_ratio > 0.15:
        shape_type = ShapeType.ROUNDED_RECTANGLE

    return VisualStyle(
        shape_type=shape_type,
        shadow=shadow,
        gradient=gradient,
        bevel_enabled=False,  # Can be extended later
    )


class PPTXRenderer:
    """
    Renders PositionedLayout to PowerPoint files.

    The renderer is stateless — each render() call creates a new presentation.
    Supports visual styles for professional-grade output.

    STYLE PRIORITY:
    1. Learned styles from training (when visual_style is provided)
    2. Hardcoded archetype styles (FALLBACK only when no learned style)
    3. Default professional style
    """

    def __init__(self, default_style: Optional[VisualStyle] = None, visual_style=None):
        """
        Initialize renderer with optional default visual style.

        Args:
            default_style: Default visual style for blocks. If None, uses STYLE_PROFESSIONAL.
            visual_style: Optional learned DesignStyle to apply (takes precedence over hardcoded).
        """
        # Track whether a learned style was provided - this affects archetype style application
        self._has_learned_style = False

        # If a learned visual_style is provided, convert it and mark as learned
        if visual_style is not None:
            converted = convert_design_style_to_visual_style(visual_style)
            if converted:
                self.default_style = converted
                self._has_learned_style = True  # Learned style takes precedence
            else:
                self.default_style = STYLE_PROFESSIONAL
        else:
            self.default_style = default_style or STYLE_PROFESSIONAL
        self._element_styles: Dict[str, VisualStyle] = {}

    def set_element_style(self, element_id: str, style: VisualStyle) -> None:
        """Set a specific visual style for an element by ID."""
        self._element_styles[element_id] = style

    def set_archetype_style(self, archetype: str, style: VisualStyle) -> None:
        """
        Set visual style for all elements of a specific archetype.

        Args:
            archetype: Archetype name (e.g., "pyramid", "marketecture")
            style: Visual style to apply
        """
        # Store as special key
        self._element_styles[f"__archetype__{archetype}"] = style

    def get_style_for_element(
        self,
        element: PositionedElement,
        archetype: Optional[str] = None
    ) -> VisualStyle:
        """
        Get the appropriate visual style for an element.

        Priority:
        1. Element-specific style (by ID)
        2. Archetype-specific style
        3. Default renderer style
        """
        # Check element-specific style
        if element.id and element.id in self._element_styles:
            return self._element_styles[element.id]

        # Check archetype-specific style
        if archetype and f"__archetype__{archetype}" in self._element_styles:
            return self._element_styles[f"__archetype__{archetype}"]

        # Return default
        return self.default_style

    def render(
        self,
        layout: PositionedLayout,
        output: Optional[Union[str, Path, BinaryIO]] = None
    ) -> bytes:
        """
        Render a PositionedLayout to PPTX format.

        Args:
            layout: The positioned layout to render
            output: Optional output path or file-like object
                    If None, returns bytes

        Returns:
            PPTX file contents as bytes
        """
        # Create presentation with correct dimensions
        prs = Presentation()
        prs.slide_width = Emu(inches_to_emu(layout.slide_width_inches))
        prs.slide_height = Emu(inches_to_emu(layout.slide_height_inches))

        # Set archetype-specific style based on learnings or fallback to BASIC
        # Priority: 1) Learned style from JSON, 2) BASIC style (no effects)
        if layout.archetype and not self._has_learned_style:
            # Try to load style from shape_learnings.json
            learned_style = load_style_from_learnings(layout.archetype)

            if learned_style:
                # Use style learned from templates
                self.set_archetype_style(layout.archetype, learned_style)
            else:
                # No learnings available - use BASIC style (no effects)
                # This creates the "before training" look
                self.set_archetype_style(layout.archetype, STYLE_BASIC)

        # Add blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Set background color
        self._set_slide_background(slide, layout.background_color)

        # Render title and subtitle first (z-order highest)
        if layout.title:
            self._render_element(slide, layout.title)
        if layout.subtitle:
            self._render_element(slide, layout.subtitle)

        # Count pyramid levels if this is a pyramid archetype
        if layout.archetype == "pyramid":
            pyramid_levels = sum(
                1 for e in layout.elements
                if e.id and e.id.startswith("level_")
            )
            self._pyramid_level_count = max(2, pyramid_levels)

        # Count funnel stages if this is a funnel archetype
        if layout.archetype == "funnel":
            funnel_stages = sum(
                1 for e in layout.elements
                if e.id and e.id.startswith("stage_")
            )
            self._funnel_stage_count = max(2, funnel_stages)

        # Render elements sorted by z-order (lowest first = behind)
        for element in layout.elements_sorted_by_z_order():
            self._render_element(slide, element, archetype=layout.archetype)

        # Render connectors
        for connector in layout.connectors:
            self._render_connector(slide, connector)

        # Save to output
        if output is None:
            buffer = BytesIO()
            prs.save(buffer)
            return buffer.getvalue()
        elif isinstance(output, (str, Path)):
            prs.save(str(output))
            with open(output, 'rb') as f:
                return f.read()
        else:
            prs.save(output)
            output.seek(0)
            return output.read()

    def render_to_file(self, layout: PositionedLayout, filepath: Union[str, Path]) -> None:
        """
        Render layout directly to a file.

        Args:
            layout: The positioned layout to render
            filepath: Output file path
        """
        self.render(layout, output=filepath)

    def render_presentation(
        self,
        presentation: MultiSlidePresentation,
        output: Optional[Union[str, Path, BinaryIO]] = None
    ) -> bytes:
        """
        Render a multi-slide presentation to PPTX format.

        Args:
            presentation: The MultiSlidePresentation to render
            output: Optional output path or file-like object
                    If None, returns bytes

        Returns:
            PPTX file contents as bytes
        """
        if not presentation.slides:
            raise ValueError("Presentation has no slides")

        # Create presentation with correct dimensions
        prs = Presentation()
        prs.slide_width = Emu(inches_to_emu(presentation.slide_width_inches))
        prs.slide_height = Emu(inches_to_emu(presentation.slide_height_inches))

        # Render each slide
        for slide_layout in presentation.slides:
            self._render_slide(prs, slide_layout)

        # Save to output
        if output is None:
            buffer = BytesIO()
            prs.save(buffer)
            return buffer.getvalue()
        elif isinstance(output, (str, Path)):
            prs.save(str(output))
            with open(output, 'rb') as f:
                return f.read()
        else:
            prs.save(output)
            output.seek(0)
            return output.read()

    def _render_slide(self, prs: Presentation, layout: PositionedLayout) -> None:
        """Render a single slide to an existing presentation."""
        # Add blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Set background color
        self._set_slide_background(slide, layout.background_color)

        # Render title and subtitle first (z-order highest)
        if layout.title:
            self._render_element(slide, layout.title)
        if layout.subtitle:
            self._render_element(slide, layout.subtitle)

        # Render elements sorted by z-order (lowest first = behind)
        for element in layout.elements_sorted_by_z_order():
            self._render_element(slide, element)

        # Render connectors
        for connector in layout.connectors:
            self._render_connector(slide, connector)

        # Add speaker notes if present
        if layout.speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = layout.speaker_notes

    # =========================================================================
    # SLIDE BACKGROUND
    # =========================================================================

    def _set_slide_background(self, slide, color: str) -> None:
        """Set slide background color."""
        if color.lower() == "transparent" or color == "#FFFFFF":
            return  # White/transparent is default

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb_color(color)

    # =========================================================================
    # ELEMENT RENDERING
    # =========================================================================

    def _render_element(self, slide, element: PositionedElement, archetype: Optional[str] = None) -> None:
        """Render a single positioned element."""
        # Skip transparent fill elements without text (spacers)
        if element.fill_color == "transparent" and element.text is None:
            return

        # Convert to EMU
        left = Emu(inches_to_emu(element.x_inches))
        top = Emu(inches_to_emu(element.y_inches))
        width = Emu(inches_to_emu(element.width_inches))
        height = Emu(inches_to_emu(element.height_inches))

        # Choose shape type based on element type
        if element.element_type == ElementType.TITLE:
            self._render_title(slide, element, left, top, width, height)
        elif element.element_type == ElementType.SUBTITLE:
            self._render_subtitle(slide, element, left, top, width, height)
        elif element.element_type == ElementType.BAND:
            self._render_band(slide, element, left, top, width, height)
        elif element.element_type == ElementType.LABEL:
            self._render_label(slide, element, left, top, width, height)
        else:
            self._render_block(slide, element, left, top, width, height, archetype=archetype)

    def _render_block(
        self,
        slide,
        element: PositionedElement,
        left,
        top,
        width,
        height,
        visual_style: Optional[VisualStyle] = None,
        archetype: Optional[str] = None
    ) -> None:
        """
        Render a standard block with visual effects.

        Args:
            slide: PowerPoint slide
            element: Element to render
            left, top, width, height: Position/size in EMU
            visual_style: Optional visual style override
            archetype: Optional archetype name for style selection
        """
        # Check if this is a pyramid element that should use custom geometry
        if archetype == "pyramid" and CUSTOM_SHAPES_AVAILABLE:
            # total_pyramid_levels is set in render() when counting pyramid elements
            total_levels = getattr(self, '_pyramid_level_count', 4)
            self._render_pyramid_block(slide, element, archetype, total_levels)
            return

        # Check if this is a funnel element that should use learned shapes
        if archetype == "funnel" and element.id and element.id.startswith("stage_"):
            total_stages = getattr(self, '_funnel_stage_count', 4)
            self._render_funnel_block(slide, element, archetype, total_stages)
            return

        # Get visual style
        style = visual_style or self.get_style_for_element(element, archetype=archetype)

        # Check for shape_hint from canvas mode (overrides style shape type)
        shape_hint_map = {
            "trapezoid": MSO_SHAPE.TRAPEZOID,
            "chevron": MSO_SHAPE.CHEVRON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "diamond": MSO_SHAPE.DIAMOND,
            "callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
        }

        # Handle arrows with direction
        arrow_direction_map = {
            "up": MSO_SHAPE.UP_ARROW,
            "down": MSO_SHAPE.DOWN_ARROW,
            "left": MSO_SHAPE.LEFT_ARROW,
            "right": MSO_SHAPE.RIGHT_ARROW,
        }

        if element.shape_hint == "arrow":
            # Get arrow direction, default to up
            direction = getattr(element, 'arrow_direction', None) or "up"
            shape_type = arrow_direction_map.get(direction.lower(), MSO_SHAPE.UP_ARROW)
        elif element.shape_hint and element.shape_hint in shape_hint_map:
            shape_type = shape_hint_map[element.shape_hint]
        else:
            # Get shape type from style or default
            shape_type = SHAPE_TYPE_MAP.get(style.shape_type, MSO_SHAPE.ROUNDED_RECTANGLE)

        # Create shape
        shape = slide.shapes.add_shape(
            shape_type,
            left, top, width, height
        )

        # Set corner radius for rounded rectangles
        if style.shape_type == ShapeType.ROUNDED_RECTANGLE:
            corner_pct = min(50000, int((element.corner_radius_inches / element.height_inches) * 100000))
            try:
                shape.adjustments[0] = corner_pct / 100000
            except (IndexError, TypeError):
                pass  # Some shapes don't have adjustments

        # Base fill color (used for gradient computation)
        base_color = element.fill_color if element.fill_color and element.fill_color != "transparent" else "#4A90A4"

        # Apply fill (solid or gradient)
        if style.gradient and style.gradient.enabled:
            # Gradient will be applied via XML
            # First set solid fill so shape has a base
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb_color(base_color)
        elif element.fill_color and element.fill_color != "transparent":
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb_color(element.fill_color)
        else:
            shape.fill.background()  # No fill

        # Border
        if element.stroke_color:
            shape.line.color.rgb = hex_to_rgb_color(element.stroke_color)
            shape.line.width = Pt(element.stroke_width_pt)
        else:
            shape.line.fill.background()  # No border

        # Apply visual effects via XML manipulation
        if style.gradient or style.shadow or style.bevel_enabled:
            apply_visual_effects_to_shape(shape, style, base_color)

        # Add text if present
        if element.text:
            self._apply_text(shape, element.text, element.height_inches)

        # Render icon if present (Phase 1: Icon System)
        if element.icon_id and ICON_LIBRARY_AVAILABLE:
            self._render_icon_in_element(slide, element)

    def _render_icon_in_element(self, slide, element: PositionedElement) -> None:
        """
        Render an icon within or alongside an element.

        Icons are rendered as separate shapes with custGeom geometry derived
        from SVG path data in the icon library.

        Args:
            slide: PowerPoint slide
            element: Element containing icon_id and positioning info
        """
        if not ICON_LIBRARY_AVAILABLE or not element.icon_id:
            return

        icon_library = get_icon_library()
        if icon_library is None:
            return

        icon = icon_library.get(element.icon_id)
        if icon is None:
            return

        # Calculate icon size based on element dimensions and icon_size_ratio
        icon_size_ratio = element.icon_size_ratio or 0.4
        base_size = min(element.width_inches, element.height_inches)
        icon_size = base_size * icon_size_ratio

        # Determine icon position based on icon_position
        icon_pos = element.icon_position or IconPosition.LEFT

        padding = 0.05  # Padding between icon and element edge

        if icon_pos == IconPosition.LEFT:
            icon_x = element.x_inches + padding
            icon_y = element.y_inches + (element.height_inches - icon_size) / 2
        elif icon_pos == IconPosition.RIGHT:
            icon_x = element.x_inches + element.width_inches - icon_size - padding
            icon_y = element.y_inches + (element.height_inches - icon_size) / 2
        elif icon_pos == IconPosition.TOP:
            icon_x = element.x_inches + (element.width_inches - icon_size) / 2
            icon_y = element.y_inches + padding
        elif icon_pos == IconPosition.BOTTOM:
            icon_x = element.x_inches + (element.width_inches - icon_size) / 2
            icon_y = element.y_inches + element.height_inches - icon_size - padding
        elif icon_pos == IconPosition.CENTER:
            icon_x = element.x_inches + (element.width_inches - icon_size) / 2
            icon_y = element.y_inches + (element.height_inches - icon_size) / 2
        elif icon_pos == IconPosition.BACKGROUND:
            # Background icon is larger and centered with transparency
            icon_size = min(element.width_inches, element.height_inches) * 0.8
            icon_x = element.x_inches + (element.width_inches - icon_size) / 2
            icon_y = element.y_inches + (element.height_inches - icon_size) / 2
        else:
            # Default to center
            icon_x = element.x_inches + (element.width_inches - icon_size) / 2
            icon_y = element.y_inches + (element.height_inches - icon_size) / 2

        # Convert to EMU
        left_emu = inches_to_emu(icon_x)
        top_emu = inches_to_emu(icon_y)
        width_emu = inches_to_emu(icon_size)
        height_emu = inches_to_emu(icon_size)

        # Create icon shape using custGeom
        self._add_icon_shape(
            slide,
            icon.path,
            icon.viewBox,
            left_emu,
            top_emu,
            width_emu,
            height_emu,
            fill_color=element.icon_color or element.text.color if element.text else "#333333",
            is_background=(icon_pos == IconPosition.BACKGROUND)
        )

    def _add_icon_shape(
        self,
        slide,
        svg_path: str,
        view_box: str,
        left_emu: int,
        top_emu: int,
        width_emu: int,
        height_emu: int,
        fill_color: str = "#333333",
        is_background: bool = False
    ) -> None:
        """
        Add an icon as a freeform shape with custom geometry.

        Args:
            slide: PowerPoint slide
            svg_path: SVG path 'd' attribute
            view_box: SVG viewBox string
            left_emu: Left position in EMU
            top_emu: Top position in EMU
            width_emu: Width in EMU
            height_emu: Height in EMU
            fill_color: Fill color for the icon
            is_background: If True, render with low opacity as background
        """
        try:
            # Create custGeom XML
            custgeom_xml = _create_icon_custgeom_xml(svg_path, view_box, width_emu, height_emu)

            # Build the complete shape XML
            color_hex = fill_color.lstrip('#')
            opacity_attr = 'alpha val="30000"' if is_background else ''
            opacity_elem = f'<a:{opacity_attr}/>' if is_background else ''

            sp_xml = f'''
            <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvSpPr>
                <p:cNvPr id="0" name="Icon"/>
                <p:cNvSpPr/>
                <p:nvPr/>
              </p:nvSpPr>
              <p:spPr>
                <a:xfrm>
                  <a:off x="{left_emu}" y="{top_emu}"/>
                  <a:ext cx="{width_emu}" cy="{height_emu}"/>
                </a:xfrm>
                {custgeom_xml}
                <a:solidFill>
                  <a:srgbClr val="{color_hex}">
                    {opacity_elem}
                  </a:srgbClr>
                </a:solidFill>
                <a:ln w="0">
                  <a:noFill/>
                </a:ln>
              </p:spPr>
              <p:txBody>
                <a:bodyPr/>
                <a:lstStyle/>
                <a:p/>
              </p:txBody>
            </p:sp>
            '''

            # Parse and add to slide
            sp_element = etree.fromstring(sp_xml)
            slide.shapes._spTree.append(sp_element)

        except Exception as e:
            # Silently fail - icons are optional enhancements
            pass

    def _render_pyramid_block(
        self,
        slide,
        element: PositionedElement,
        archetype: str,
        total_pyramid_levels: int = 4,
    ) -> None:
        """
        Render a pyramid block using custom FREEFORM geometry.

        Creates proper triangular pyramid segments with pointed tops instead
        of basic trapezoid shapes. Extracts level information from element ID
        to generate appropriate geometry.

        Args:
            slide: PowerPoint slide
            element: The pyramid level element
            archetype: Archetype name (should be "pyramid")
            total_pyramid_levels: Total number of levels in the pyramid
        """
        # Parse level information from element ID (e.g., "level_0", "level_1")
        level_idx = 0
        total_levels = total_pyramid_levels

        if element.id and element.id.startswith("level_"):
            try:
                level_idx = int(element.id.split("_")[1])
            except (ValueError, IndexError):
                pass

        # Get shape from library
        library = get_shape_library()
        shapes = library.get_pyramid_shapes(total_levels)

        # Get the appropriate shape for this level
        if level_idx < len(shapes):
            pyramid_shape = shapes[level_idx]
        else:
            # Generate a fallback shape
            pyramid_shape = ShapeGenerator.create_pyramid_segment(
                level=level_idx,
                total_levels=total_levels,
            )

        # Fill color
        fill_color = element.fill_color if element.fill_color and element.fill_color != "transparent" else "#4472C4"

        # Determine if gradient and shadow should be used based on learnings
        learned_style = load_style_from_learnings("pyramid")
        if learned_style:
            # Use learned parameters
            use_gradient = learned_style.gradient is not None and learned_style.gradient.enabled
            use_shadow = learned_style.shadow is not None and learned_style.shadow.enabled
            shadow_blur = learned_style.shadow.blur_radius_pt if learned_style.shadow else 12.0
            shadow_opacity = learned_style.shadow.opacity if learned_style.shadow else 0.4
        else:
            # No learnings = BASIC mode (no effects)
            use_gradient = False
            use_shadow = False
            shadow_blur = 0.0
            shadow_opacity = 0.0

        # Check if stroke should be applied
        stroke_color = element.stroke_color if element.stroke_width_pt > 0 else None
        stroke_width = element.stroke_width_pt if stroke_color else 0

        # Render custom shape
        renderer = CustomShapeRenderer(slide)
        renderer.add_shape(
            shape=pyramid_shape,
            left=element.x_inches,
            top=element.y_inches,
            width=element.width_inches,
            height=element.height_inches,
            fill_color=fill_color,
            gradient=use_gradient,
            gradient_angle=270.0,  # Top to bottom
            shadow=use_shadow,
            shadow_blur_pt=shadow_blur,
            shadow_opacity=shadow_opacity,
            stroke_color=stroke_color,
            stroke_width_pt=stroke_width,
        )

        # Add text overlay if present
        if element.text:
            self._add_pyramid_text_overlay(slide, element, level_idx, total_levels)

    def _add_pyramid_text_overlay(
        self,
        slide,
        element: PositionedElement,
        level_idx: int = 0,
        total_levels: int = 4,
    ) -> None:
        """
        Add text overlay on top of a custom pyramid shape.

        For trapezoid/triangle shapes, the visual center is lower than the
        geometric center. We adjust the text position and width to fit
        within the visible shape area.

        Args:
            slide: PowerPoint slide
            element: The pyramid level element
            level_idx: Current level index (0=base, N-1=apex)
            total_levels: Total number of pyramid levels
        """
        if not element.text:
            return

        # For pyramid shapes, adjust text position based on shape geometry
        # The visible area of a trapezoid is wider at bottom than top
        # Text should be positioned in the center of the visible area

        # Calculate the average width of the shape (considering taper)
        # For trapezoids, the text area should be in the middle 60% of height
        # and horizontally centered based on average visible width

        is_apex = (level_idx == total_levels - 1)

        if is_apex:
            # Triangle apex - text goes in lower 2/3 where it's wider
            text_top = element.y_inches + element.height_inches * 0.35
            text_height = element.height_inches * 0.55
            # Width is narrower for apex - use 70% centered
            text_width = element.width_inches * 0.7
            text_left = element.x_inches + (element.width_inches - text_width) / 2
        else:
            # Trapezoid - text centered but slightly lower
            # The visible center is lower due to the taper
            text_top = element.y_inches + element.height_inches * 0.2
            text_height = element.height_inches * 0.6
            # Use 80% of width for text area (accounting for angled edges)
            text_width = element.width_inches * 0.8
            text_left = element.x_inches + (element.width_inches - text_width) / 2

        left = Emu(inches_to_emu(text_left))
        top = Emu(inches_to_emu(text_top))
        width = Emu(inches_to_emu(text_width))
        height = Emu(inches_to_emu(text_height))

        textbox = slide.shapes.add_textbox(left, top, width, height)
        self._apply_text(textbox, element.text, text_height)

    def _render_funnel_block(
        self,
        slide,
        element: PositionedElement,
        archetype: str,
        total_funnel_stages: int = 4,
    ) -> None:
        """
        Render a funnel block using custom FREEFORM geometry.

        Creates proper funnel segments with curved sides that taper from top to bottom.
        Extracts stage information from element ID to generate appropriate geometry.

        Args:
            slide: PowerPoint slide
            element: The funnel stage element
            archetype: Archetype name (should be "funnel")
            total_funnel_stages: Total number of stages in the funnel
        """
        # Parse stage information from element ID (e.g., "stage_0", "stage_1")
        stage_idx = 0
        total_stages = total_funnel_stages

        if element.id and element.id.startswith("stage_"):
            try:
                stage_idx = int(element.id.split("_")[1])
            except (ValueError, IndexError):
                pass

        # Get shape from library
        library = get_shape_library()
        shapes = library.get_funnel_shapes(total_stages)

        # Get the appropriate shape for this stage
        if stage_idx < len(shapes):
            funnel_shape = shapes[stage_idx]
        else:
            # Generate a fallback shape
            funnel_shape = ShapeGenerator.create_funnel_segment(
                stage=stage_idx,
                total_stages=total_stages,
            )

        # Fill color
        fill_color = element.fill_color if element.fill_color and element.fill_color != "transparent" else "#4472C4"

        # Determine if gradient and shadow should be used based on learnings
        learned_style = load_style_from_learnings("funnel")
        if learned_style:
            # Use learned parameters
            use_gradient = learned_style.gradient is not None and learned_style.gradient.enabled
            use_shadow = learned_style.shadow is not None and learned_style.shadow.enabled
            shadow_blur = learned_style.shadow.blur_radius_pt if learned_style.shadow else 12.0
            shadow_opacity = learned_style.shadow.opacity if learned_style.shadow else 0.4
        else:
            # No learnings = BASIC mode (no effects)
            use_gradient = False
            use_shadow = False
            shadow_blur = 0.0
            shadow_opacity = 0.0

        # Check if stroke should be applied
        stroke_color = element.stroke_color if element.stroke_width_pt > 0 else None
        stroke_width = element.stroke_width_pt if stroke_color else 0

        # Render custom shape
        renderer = CustomShapeRenderer(slide)
        renderer.add_shape(
            shape=funnel_shape,
            left=element.x_inches,
            top=element.y_inches,
            width=element.width_inches,
            height=element.height_inches,
            fill_color=fill_color,
            gradient=use_gradient,
            gradient_angle=270.0,  # Top to bottom
            shadow=use_shadow,
            shadow_blur_pt=shadow_blur,
            shadow_opacity=shadow_opacity,
            stroke_color=stroke_color,
            stroke_width_pt=stroke_width,
        )

        # Add text overlay if present
        if element.text:
            self._add_funnel_text_overlay(slide, element, stage_idx, total_stages)

    def _add_funnel_text_overlay(
        self,
        slide,
        element: PositionedElement,
        stage_idx: int = 0,
        total_stages: int = 4,
    ) -> None:
        """
        Add text overlay on top of a custom funnel shape.

        For funnel shapes (wider at top, narrower at bottom), the text should
        be positioned to fit within the visible shape area.

        Args:
            slide: PowerPoint slide
            element: The funnel stage element
            stage_idx: Current stage index (0=top/widest, N-1=bottom/narrowest)
            total_stages: Total number of funnel stages
        """
        if not element.text:
            return

        # For funnel shapes, adjust text position based on shape geometry
        # Top stages are wider, bottom stages are narrower

        is_bottom = (stage_idx == total_stages - 1)

        if is_bottom:
            # Bottom stage - narrower, center text in upper portion
            text_top = element.y_inches + element.height_inches * 0.15
            text_height = element.height_inches * 0.6
            # Width is narrower for bottom - use 60% centered
            text_width = element.width_inches * 0.6
            text_left = element.x_inches + (element.width_inches - text_width) / 2
        else:
            # Other stages - text centered with good margins
            text_top = element.y_inches + element.height_inches * 0.15
            text_height = element.height_inches * 0.7
            # Use 75% of width for text area (accounting for angled edges)
            text_width = element.width_inches * 0.75
            text_left = element.x_inches + (element.width_inches - text_width) / 2

        left = Emu(inches_to_emu(text_left))
        top = Emu(inches_to_emu(text_top))
        width = Emu(inches_to_emu(text_width))
        height = Emu(inches_to_emu(text_height))

        textbox = slide.shapes.add_textbox(left, top, width, height)
        self._apply_text(textbox, element.text, text_height)

    def _render_band(self, slide, element: PositionedElement, left, top, width, height) -> None:
        """Render a full-width band (rectangle)."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )

        # Fill with semi-transparency
        if element.fill_color and element.fill_color != "transparent":
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb_color(element.fill_color)

            # Set transparency if specified
            if element.opacity < 1.0:
                # Transparency is inverted (0 = opaque, 100000 = transparent)
                transparency = int((1 - element.opacity) * 100000)
                shape.fill.fore_color.brightness = 0  # Reset brightness
                # Note: python-pptx doesn't directly support fill transparency
                # Would need to modify XML directly for true transparency
        else:
            shape.fill.background()

        # No border for bands
        shape.line.fill.background()

        # Add text
        if element.text:
            self._apply_text(shape, element.text, element.height_inches)

    def _render_title(self, slide, element: PositionedElement, left, top, width, height) -> None:
        """Render title text."""
        if not element.text:
            return

        # Use text box for title
        textbox = slide.shapes.add_textbox(left, top, width, height)
        self._apply_text(textbox, element.text, element.height_inches, is_title=True)

    def _render_subtitle(self, slide, element: PositionedElement, left, top, width, height) -> None:
        """Render subtitle text."""
        if not element.text:
            return

        textbox = slide.shapes.add_textbox(left, top, width, height)
        self._apply_text(textbox, element.text, element.height_inches)

    def _render_label(self, slide, element: PositionedElement, left, top, width, height) -> None:
        """Render a text label."""
        if not element.text:
            return

        textbox = slide.shapes.add_textbox(left, top, width, height)
        self._apply_text(textbox, element.text, element.height_inches)

    # =========================================================================
    # TEXT APPLICATION (CRITICAL - PREVENTS OVERFLOW)
    # =========================================================================

    def _apply_text(
        self,
        shape,
        text: PositionedText,
        container_height: float,
        is_title: bool = False
    ) -> None:
        """
        Apply pre-measured text to a shape.

        CRITICAL: Uses pre-computed font sizes and line breaks.
        Never lets PowerPoint auto-size or wrap text.
        """
        tf = shape.text_frame

        # CRITICAL: Disable auto-size — we computed the font size
        tf.auto_size = None  # Explicitly None to prevent any auto-sizing

        # Set margins (minimal)
        tf.margin_left = Emu(inches_to_emu(0.05))
        tf.margin_right = Emu(inches_to_emu(0.05))
        tf.margin_top = Emu(inches_to_emu(0.03))
        tf.margin_bottom = Emu(inches_to_emu(0.03))

        # Vertical centering
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Word wrap — only if multiple lines
        tf.word_wrap = len(text.lines) > 1

        # Clear existing paragraphs and add new ones
        for i, line in enumerate(text.lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = line
            p.alignment = get_alignment(text.alignment)

            # Apply font formatting to the run
            if p.runs:
                run = p.runs[0]
            else:
                run = p.add_run()
                run.text = line

            # CRITICAL: Use pre-computed font size
            run.font.size = Pt(text.font_size_pt)
            run.font.name = text.font_family
            run.font.bold = text.bold
            run.font.italic = text.italic

            if text.color and text.color != "transparent":
                run.font.color.rgb = hex_to_rgb_color(text.color)

            # Line spacing for multi-line text
            if len(text.lines) > 1:
                p.line_spacing = 1.15

    # =========================================================================
    # CONNECTOR RENDERING
    # =========================================================================

    def _render_connector(self, slide, connector: PositionedConnector) -> None:
        """
        Render a connector line using a freeform shape for precise positioning.

        Supports both simple two-point connectors and polyline connectors with
        intermediate waypoints for orthogonal routing.
        """
        from pptx.oxml.ns import nsmap
        from pptx.oxml import parse_xml
        from lxml import etree

        # Get all points including waypoints
        all_points = connector.all_points if hasattr(connector, 'all_points') else [
            (connector.start_x, connector.start_y),
            (connector.end_x, connector.end_y)
        ]

        # Convert all points to EMU
        points_emu = [(inches_to_emu(x), inches_to_emu(y)) for x, y in all_points]

        # Calculate bounding box
        all_x = [p[0] for p in points_emu]
        all_y = [p[1] for p in points_emu]
        min_x = min(all_x)
        min_y = min(all_y)
        max_x = max(all_x)
        max_y = max(all_y)

        width = max_x - min_x
        height = max_y - min_y

        # Ensure minimum size
        if width < 1:
            width = 1
        if height < 1:
            height = 1

        # Build path elements for all segments
        path_elements = []

        # Start with moveTo
        rel_start_x = points_emu[0][0] - min_x
        rel_start_y = points_emu[0][1] - min_y
        path_elements.append(f'<a:moveTo><a:pt x="{rel_start_x}" y="{rel_start_y}"/></a:moveTo>')

        # Check if we should use curved corners for orthogonal routing
        routing_style = getattr(connector, 'routing_style', None)
        corner_radius_emu = int(connector.corner_radius_inches * 914400) if hasattr(connector, 'corner_radius_inches') else 0

        # Add lineTo for each subsequent point
        for i in range(1, len(points_emu)):
            rel_x = points_emu[i][0] - min_x
            rel_y = points_emu[i][1] - min_y

            # For orthogonal routing with corner radius, we could add curves
            # For now, use straight lines (corner smoothing requires quadBezTo)
            if routing_style == RoutingStyle.ORTHOGONAL and corner_radius_emu > 0 and i < len(points_emu) - 1:
                # This is an intermediate point - could add curve here
                # For simplicity, just use line for now
                path_elements.append(f'<a:lnTo><a:pt x="{rel_x}" y="{rel_y}"/></a:lnTo>')
            else:
                path_elements.append(f'<a:lnTo><a:pt x="{rel_x}" y="{rel_y}"/></a:lnTo>')

        path_xml = '\n                  '.join(path_elements)

        # Determine dash style
        dash_xml = ''
        if connector.style == ConnectorStyle.DASHED:
            dash_xml = '<a:prstDash val="dash"/>'

        # Build the shape XML
        sp_xml = f'''
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="0" name="Connector"/>
            <p:cNvSpPr/>
            <p:nvPr/>
          </p:nvSpPr>
          <p:spPr>
            <a:xfrm>
              <a:off x="{min_x}" y="{min_y}"/>
              <a:ext cx="{width}" cy="{height}"/>
            </a:xfrm>
            <a:custGeom>
              <a:avLst/>
              <a:gdLst/>
              <a:ahLst/>
              <a:cxnLst/>
              <a:rect l="0" t="0" r="0" b="0"/>
              <a:pathLst>
                <a:path w="{width}" h="{height}">
                  {path_xml}
                </a:path>
              </a:pathLst>
            </a:custGeom>
            <a:ln w="{int(connector.stroke_width_pt * 12700)}">
              <a:solidFill>
                <a:srgbClr val="{connector.color.lstrip('#')}"/>
              </a:solidFill>
              {dash_xml}
            </a:ln>
          </p:spPr>
          <p:txBody>
            <a:bodyPr/>
            <a:lstStyle/>
            <a:p/>
          </p:txBody>
        </p:sp>
        '''

        # Parse and add to slide
        sp_element = etree.fromstring(sp_xml)
        slide.shapes._spTree.append(sp_element)

        # For arrow heads, we need to add them via XML as well
        if connector.style in [ConnectorStyle.ARROW, ConnectorStyle.DASHED, ConnectorStyle.BIDIRECTIONAL]:
            self._add_arrow_to_xml(sp_element, connector.style)

        # Add label if present
        if connector.label:
            self._add_connector_label(slide, connector)

    def _add_arrow_to_xml(self, sp_element, style: ConnectorStyle) -> None:
        """Add arrowhead(s) to a shape via XML."""
        try:
            nsmap_a = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            ln = sp_element.find('.//a:ln', nsmap_a)
            if ln is not None:
                # Add end arrow
                tail_end = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
                tail_end.set('type', 'triangle')
                tail_end.set('w', 'med')
                tail_end.set('len', 'med')

                # Add start arrow for bidirectional
                if style == ConnectorStyle.BIDIRECTIONAL:
                    head_end = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
                    head_end.set('type', 'triangle')
                    head_end.set('w', 'med')
                    head_end.set('len', 'med')

                # Add dash for dashed style
                if style == ConnectorStyle.DASHED:
                    prstDash = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash')
                    prstDash.set('val', 'dash')
        except Exception:
            pass

    def _add_end_arrow(self, shape) -> None:
        """Add arrowhead to end of connector."""
        # Access line element and add arrow
        try:
            line = shape.line
            # Set end arrow using underlying XML
            spPr = shape._element.spPr
            ln = spPr.find(qn('a:ln'))
            if ln is not None:
                tailEnd = parse_xml(
                    '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                    'type="triangle" w="med" len="med"/>'
                )
                ln.append(tailEnd)
        except Exception:
            pass  # Arrow heads may fail in some cases

    def _add_start_arrow(self, shape) -> None:
        """Add arrowhead to start of connector."""
        try:
            spPr = shape._element.spPr
            ln = spPr.find(qn('a:ln'))
            if ln is not None:
                headEnd = parse_xml(
                    '<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                    'type="triangle" w="med" len="med"/>'
                )
                ln.append(headEnd)
        except Exception:
            pass

    def _add_connector_label(self, slide, connector: PositionedConnector) -> None:
        """Add a text label at connector midpoint."""
        if not connector.label:
            return

        # Position label at midpoint
        mid_x = connector.midpoint_x
        mid_y = connector.midpoint_y

        # Small textbox for label
        label_width = 1.5
        label_height = 0.3

        left = Emu(inches_to_emu(mid_x - label_width / 2))
        top = Emu(inches_to_emu(mid_y - label_height / 2))
        width = Emu(inches_to_emu(label_width))
        height = Emu(inches_to_emu(label_height))

        textbox = slide.shapes.add_textbox(left, top, width, height)
        self._apply_text(textbox, connector.label, label_height)


# =============================================================================
# CONVENIENCE FUNCTIONS
# =============================================================================

def render_to_pptx(
    layout: PositionedLayout,
    output_path: Optional[str] = None,
    style: Optional[VisualStyle] = None
) -> bytes:
    """
    Convenience function to render a layout to PPTX.

    Args:
        layout: PositionedLayout to render
        output_path: Optional file path to save to
        style: Optional visual style (defaults to STYLE_PROFESSIONAL)

    Returns:
        PPTX file bytes
    """
    renderer = PPTXRenderer(default_style=style)
    if output_path:
        renderer.render_to_file(layout, output_path)
        with open(output_path, 'rb') as f:
            return f.read()
    else:
        return renderer.render(layout)


def render_to_bytes(
    layout: PositionedLayout,
    style: Optional[VisualStyle] = None
) -> bytes:
    """Render layout to PPTX bytes (for API responses)."""
    renderer = PPTXRenderer(default_style=style)
    return renderer.render(layout)


def render_styled(
    layout: PositionedLayout,
    style_name: str = "professional",
    output_path: Optional[str] = None
) -> bytes:
    """
    Render a layout with a named visual style preset.

    Args:
        layout: PositionedLayout to render
        style_name: Style preset name: "flat", "subtle_3d", "professional", "executive", "pyramid"
        output_path: Optional output file path

    Returns:
        PPTX file bytes
    """
    style_map = {
        "flat": STYLE_FLAT,
        "subtle_3d": STYLE_SUBTLE_3D,
        "professional": STYLE_PROFESSIONAL,
        "executive": STYLE_EXECUTIVE,
        "pyramid": STYLE_PYRAMID_LEVEL,
    }

    style = style_map.get(style_name.lower(), STYLE_PROFESSIONAL)
    return render_to_pptx(layout, output_path, style)


def render_presentation_to_pptx(
    presentation: MultiSlidePresentation,
    output_path: Optional[str] = None
) -> bytes:
    """
    Convenience function to render a multi-slide presentation to PPTX.

    Args:
        presentation: MultiSlidePresentation to render
        output_path: Optional file path to save to

    Returns:
        PPTX file bytes
    """
    renderer = PPTXRenderer()
    if output_path:
        return renderer.render_presentation(presentation, output_path)
    else:
        return renderer.render_presentation(presentation)


def render_presentation_to_bytes(presentation: MultiSlidePresentation) -> bytes:
    """Render multi-slide presentation to PPTX bytes (for API responses)."""
    renderer = PPTXRenderer()
    return renderer.render_presentation(presentation)


def create_presentation_from_layouts(
    layouts: list,
    title: str = "Presentation"
) -> MultiSlidePresentation:
    """
    Create a MultiSlidePresentation from a list of PositionedLayout objects.

    Args:
        layouts: List of PositionedLayout objects
        title: Presentation title

    Returns:
        MultiSlidePresentation containing all slides
    """
    from datetime import datetime

    presentation = MultiSlidePresentation(
        presentation_title=title,
        created_at=datetime.now().isoformat(),
    )

    for layout in layouts:
        presentation.add_slide(layout)

    return presentation
