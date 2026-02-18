"""
shape_learning.py — Shape extraction, learning, and generation system.

This module provides a versatile system for:
1. Extracting shapes from template PPTX files (FREEFORM paths)
2. Storing learned shapes in a reusable library
3. Generating custom shapes based on parameters or context
4. Creating variations of learned shapes

The system is designed to learn from ANY template, not just pyramids.
It extracts the actual geometry (vertices, curves) and can recreate
or adapt those shapes for new contexts.
"""

import json
import math
import hashlib
from pathlib import Path
from dataclasses import dataclass, field, asdict
from typing import List, Dict, Optional, Tuple, Any, Union
from enum import Enum
import logging

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

logger = logging.getLogger(__name__)


# =============================================================================
# SHAPE PRIMITIVES
# =============================================================================

class PathCommand(Enum):
    """SVG-like path commands for shape definition."""
    MOVE_TO = "M"       # Move to point
    LINE_TO = "L"       # Line to point
    CURVE_TO = "C"      # Cubic bezier curve
    QUAD_TO = "Q"       # Quadratic bezier curve
    ARC_TO = "A"        # Arc
    CLOSE = "Z"         # Close path


@dataclass
class PathPoint:
    """A point in the shape path (normalized 0-1 coordinates)."""
    x: float  # 0.0 = left edge, 1.0 = right edge
    y: float  # 0.0 = top edge, 1.0 = bottom edge

    def scale(self, width: float, height: float) -> Tuple[float, float]:
        """Scale to actual dimensions."""
        return (self.x * width, self.y * height)

    def to_emu(self, width_emu: int, height_emu: int) -> Tuple[int, int]:
        """Convert to EMU coordinates."""
        return (int(self.x * width_emu), int(self.y * height_emu))


@dataclass
class PathSegment:
    """A segment of a shape path."""
    command: PathCommand
    points: List[PathPoint] = field(default_factory=list)
    # For arc commands
    arc_params: Optional[Dict[str, float]] = None


@dataclass
class LearnedShape:
    """
    A shape learned from a template or generated programmatically.

    Shapes are stored in normalized coordinates (0-1) so they can be
    scaled to any size while maintaining proportions.
    """
    id: str
    name: str
    description: str
    category: str  # "pyramid", "chevron", "arrow", "custom", etc.

    # The actual shape path (normalized 0-1 coordinates)
    path: List[PathSegment]

    # Bounding box aspect ratio (width/height)
    aspect_ratio: float = 1.0

    # Visual properties learned from template
    suggested_fill_type: str = "solid"  # solid, gradient, pattern
    suggested_gradient_angle: float = 270.0  # degrees
    suggested_shadow: bool = True

    # Metadata
    source_template: Optional[str] = None
    source_slide: Optional[int] = None
    tags: List[str] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        """Convert to JSON-serializable dictionary."""
        return {
            "id": self.id,
            "name": self.name,
            "description": self.description,
            "category": self.category,
            "path": [
                {
                    "command": seg.command.value,
                    "points": [{"x": p.x, "y": p.y} for p in seg.points],
                    "arc_params": seg.arc_params,
                }
                for seg in self.path
            ],
            "aspect_ratio": self.aspect_ratio,
            "suggested_fill_type": self.suggested_fill_type,
            "suggested_gradient_angle": self.suggested_gradient_angle,
            "suggested_shadow": self.suggested_shadow,
            "source_template": self.source_template,
            "source_slide": self.source_slide,
            "tags": self.tags,
        }

    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "LearnedShape":
        """Create from dictionary."""
        path = [
            PathSegment(
                command=PathCommand(seg["command"]),
                points=[PathPoint(p["x"], p["y"]) for p in seg.get("points", [])],
                arc_params=seg.get("arc_params"),
            )
            for seg in data.get("path", [])
        ]
        return cls(
            id=data["id"],
            name=data["name"],
            description=data.get("description", ""),
            category=data.get("category", "custom"),
            path=path,
            aspect_ratio=data.get("aspect_ratio", 1.0),
            suggested_fill_type=data.get("suggested_fill_type", "solid"),
            suggested_gradient_angle=data.get("suggested_gradient_angle", 270.0),
            suggested_shadow=data.get("suggested_shadow", True),
            source_template=data.get("source_template"),
            source_slide=data.get("source_slide"),
            tags=data.get("tags", []),
        )


# =============================================================================
# SHAPE EXTRACTOR - Learn from templates
# =============================================================================

class ShapeExtractor:
    """
    Extracts shapes from PPTX templates.

    Parses FREEFORM shapes and converts their DrawingML paths to
    normalized LearnedShape objects that can be reused and scaled.
    """

    def __init__(self):
        self.extracted_shapes: List[LearnedShape] = []

    def extract_from_pptx(
        self,
        pptx_path: str,
        category: str = "custom",
        slide_filter: Optional[List[int]] = None,
    ) -> List[LearnedShape]:
        """
        Extract all FREEFORM shapes from a PPTX file.

        Args:
            pptx_path: Path to the PPTX template
            category: Category to assign to extracted shapes
            slide_filter: Optional list of slide indices to process (1-indexed)

        Returns:
            List of LearnedShape objects
        """
        prs = Presentation(pptx_path)
        template_name = Path(pptx_path).stem
        shapes = []

        slide_num = 0
        for slide in prs.slides:
            slide_num += 1

            if slide_filter and slide_num not in slide_filter:
                continue

            shape_idx = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                    shape_idx += 1
                    try:
                        learned = self._extract_freeform(
                            shape,
                            template_name=template_name,
                            slide_num=slide_num,
                            shape_idx=shape_idx,
                            category=category,
                        )
                        if learned:
                            shapes.append(learned)
                    except Exception as e:
                        logger.warning(f"Failed to extract shape from slide {slide_num}: {e}")

        self.extracted_shapes.extend(shapes)
        return shapes

    def _extract_freeform(
        self,
        shape,
        template_name: str,
        slide_num: int,
        shape_idx: int,
        category: str,
    ) -> Optional[LearnedShape]:
        """Extract a single FREEFORM shape."""
        try:
            # Get shape dimensions
            width_emu = shape.width
            height_emu = shape.height

            if width_emu == 0 or height_emu == 0:
                return None

            # Get the path from XML
            spPr = shape._element.find(qn('p:spPr'))
            if spPr is None:
                return None

            custGeom = spPr.find(qn('a:custGeom'))
            if custGeom is None:
                return None

            pathLst = custGeom.find(qn('a:pathLst'))
            if pathLst is None:
                return None

            # Parse all paths
            path_segments = []
            for path_elem in pathLst.findall(qn('a:path')):
                segments = self._parse_drawingml_path(path_elem, width_emu, height_emu)
                path_segments.extend(segments)

            if not path_segments:
                return None

            # Check for gradient
            has_gradient = spPr.find('.//' + qn('a:gradFill')) is not None
            gradient_angle = 270.0
            if has_gradient:
                lin = spPr.find('.//' + qn('a:lin'))
                if lin is not None:
                    ang = lin.get('ang')
                    if ang:
                        gradient_angle = int(ang) / 60000

            # Check for shadow
            has_shadow = spPr.find('.//' + qn('a:outerShdw')) is not None

            # Generate ID
            shape_id = f"{category}_{template_name}_s{slide_num}_{shape_idx}"
            shape_hash = hashlib.md5(str(path_segments).encode()).hexdigest()[:8]
            shape_id = f"{shape_id}_{shape_hash}"

            return LearnedShape(
                id=shape_id,
                name=f"{category.title()} from {template_name} (Slide {slide_num})",
                description=f"Shape extracted from {template_name}, slide {slide_num}",
                category=category,
                path=path_segments,
                aspect_ratio=width_emu / height_emu if height_emu else 1.0,
                suggested_fill_type="gradient" if has_gradient else "solid",
                suggested_gradient_angle=gradient_angle,
                suggested_shadow=has_shadow,
                source_template=template_name,
                source_slide=slide_num,
                tags=[category, template_name],
            )

        except Exception as e:
            logger.warning(f"Error extracting freeform: {e}")
            return None

    def _parse_drawingml_path(
        self,
        path_elem,
        width_emu: int,
        height_emu: int,
    ) -> List[PathSegment]:
        """Parse DrawingML path element to PathSegments."""
        segments = []

        # Get path dimensions (may differ from shape dimensions)
        path_w = int(path_elem.get('w', width_emu))
        path_h = int(path_elem.get('h', height_emu))

        for child in path_elem:
            tag = child.tag.replace('{http://schemas.openxmlformats.org/drawingml/2006/main}', '')

            if tag == 'moveTo':
                pt = child.find(qn('a:pt'))
                if pt is not None:
                    x = int(pt.get('x', 0)) / path_w
                    y = int(pt.get('y', 0)) / path_h
                    segments.append(PathSegment(
                        command=PathCommand.MOVE_TO,
                        points=[PathPoint(x, y)]
                    ))

            elif tag == 'lnTo':
                pt = child.find(qn('a:pt'))
                if pt is not None:
                    x = int(pt.get('x', 0)) / path_w
                    y = int(pt.get('y', 0)) / path_h
                    segments.append(PathSegment(
                        command=PathCommand.LINE_TO,
                        points=[PathPoint(x, y)]
                    ))

            elif tag == 'cubicBezTo':
                pts = child.findall(qn('a:pt'))
                if len(pts) >= 3:
                    points = []
                    for pt in pts[:3]:
                        x = int(pt.get('x', 0)) / path_w
                        y = int(pt.get('y', 0)) / path_h
                        points.append(PathPoint(x, y))
                    segments.append(PathSegment(
                        command=PathCommand.CURVE_TO,
                        points=points
                    ))

            elif tag == 'quadBezTo':
                pts = child.findall(qn('a:pt'))
                if len(pts) >= 2:
                    points = []
                    for pt in pts[:2]:
                        x = int(pt.get('x', 0)) / path_w
                        y = int(pt.get('y', 0)) / path_h
                        points.append(PathPoint(x, y))
                    segments.append(PathSegment(
                        command=PathCommand.QUAD_TO,
                        points=points
                    ))

            elif tag == 'arcTo':
                # Arc parameters
                segments.append(PathSegment(
                    command=PathCommand.ARC_TO,
                    arc_params={
                        'wR': int(child.get('wR', 0)) / path_w,
                        'hR': int(child.get('hR', 0)) / path_h,
                        'stAng': int(child.get('stAng', 0)) / 60000,
                        'swAng': int(child.get('swAng', 0)) / 60000,
                    }
                ))

            elif tag == 'close':
                segments.append(PathSegment(command=PathCommand.CLOSE))

        return segments


# =============================================================================
# SHAPE GENERATOR - Create parametric shapes
# =============================================================================

class ShapeGenerator:
    """
    Generates shapes programmatically based on parameters.

    Can create:
    - Basic geometric shapes (triangles, polygons)
    - Parametric shapes (pyramids with N levels, arrows, chevrons)
    - Variations of learned shapes
    """

    @staticmethod
    def create_triangle(
        pointed_top: bool = True,
        corner_radius: float = 0.0,
    ) -> LearnedShape:
        """Create a basic triangle shape."""
        if pointed_top:
            # Point at top, base at bottom
            path = [
                PathSegment(PathCommand.MOVE_TO, [PathPoint(0.5, 0.0)]),   # Top center
                PathSegment(PathCommand.LINE_TO, [PathPoint(1.0, 1.0)]),   # Bottom right
                PathSegment(PathCommand.LINE_TO, [PathPoint(0.0, 1.0)]),   # Bottom left
                PathSegment(PathCommand.CLOSE),
            ]
        else:
            # Point at bottom (inverted)
            path = [
                PathSegment(PathCommand.MOVE_TO, [PathPoint(0.0, 0.0)]),   # Top left
                PathSegment(PathCommand.LINE_TO, [PathPoint(1.0, 0.0)]),   # Top right
                PathSegment(PathCommand.LINE_TO, [PathPoint(0.5, 1.0)]),   # Bottom center
                PathSegment(PathCommand.CLOSE),
            ]

        return LearnedShape(
            id="triangle_basic",
            name="Basic Triangle",
            description="Simple triangular shape",
            category="basic",
            path=path,
            aspect_ratio=1.0,
            tags=["triangle", "basic", "geometric"],
        )

    @staticmethod
    def create_pyramid_segment(
        level: int,
        total_levels: int,
        base_ratio: float = 0.85,  # Base width as ratio of total (matches archetype default)
        top_ratio: float = 0.25,   # Top width as ratio of total (matches archetype default)
    ) -> LearnedShape:
        """
        Create a single pyramid level segment with aligned edges.

        The key for proper alignment: each shape's TOP edge must match the
        WIDTH of the next level up. This creates continuous diagonal edges
        when shapes are stacked at decreasing widths.

        Math:
        - Archetype uses linear width stepping from base_ratio to top_ratio
        - W_i = base - i * step, where step = (base - top) / (N-1)
        - This shape's top edge ratio = W_{i+1} / W_i

        Args:
            level: Current level (0 = base, total_levels-1 = apex)
            total_levels: Total number of levels
            base_ratio: Base width ratio (should match archetype config)
            top_ratio: Top width ratio (should match archetype config)
        """
        # Calculate width stepping (matches archetype logic)
        step = (base_ratio - top_ratio) / (total_levels - 1) if total_levels > 1 else 0

        # Width of this level and next level
        width_this = base_ratio - level * step
        width_next = base_ratio - (level + 1) * step if level < total_levels - 1 else 0

        # For apex level, create a pointed triangle
        if level == total_levels - 1:
            path = [
                PathSegment(PathCommand.MOVE_TO, [PathPoint(0.5, 0.0)]),   # Apex point
                PathSegment(PathCommand.LINE_TO, [PathPoint(1.0, 1.0)]),   # Bottom right
                PathSegment(PathCommand.LINE_TO, [PathPoint(0.0, 1.0)]),   # Bottom left
                PathSegment(PathCommand.CLOSE),
            ]
        else:
            # Calculate top edge ratio: next level's width / this level's width
            # This ensures diagonal edges align across all levels
            top_edge_ratio = width_next / width_this if width_this > 0 else 0.5

            # Inset on each side to center the narrower top edge
            inset = (1.0 - top_edge_ratio) / 2
            left_top = inset
            right_top = 1.0 - inset

            path = [
                PathSegment(PathCommand.MOVE_TO, [PathPoint(left_top, 0.0)]),   # Top left
                PathSegment(PathCommand.LINE_TO, [PathPoint(right_top, 0.0)]),  # Top right
                PathSegment(PathCommand.LINE_TO, [PathPoint(1.0, 1.0)]),        # Bottom right
                PathSegment(PathCommand.LINE_TO, [PathPoint(0.0, 1.0)]),        # Bottom left
                PathSegment(PathCommand.CLOSE),
            ]

        return LearnedShape(
            id=f"pyramid_segment_L{level}_of_{total_levels}",
            name=f"Pyramid Level {level + 1} of {total_levels}",
            description=f"Pyramid segment for level {level + 1}",
            category="pyramid",
            path=path,
            aspect_ratio=2.0,  # Wider than tall
            suggested_fill_type="gradient",
            suggested_gradient_angle=270.0,
            suggested_shadow=True,
            tags=["pyramid", "segment", f"level_{level}"],
        )

    @staticmethod
    def create_full_pyramid(
        num_levels: int = 4,
        gap_ratio: float = 0.02,  # Gap between levels as ratio of total height
    ) -> List[LearnedShape]:
        """
        Create all segments for a complete pyramid.

        Returns list of shapes from base to apex.
        """
        shapes = []
        for level in range(num_levels):
            shape = ShapeGenerator.create_pyramid_segment(level, num_levels)
            shapes.append(shape)
        return shapes

    @staticmethod
    def create_chevron(
        point_depth: float = 0.2,  # How deep the chevron point goes (0-0.5)
    ) -> LearnedShape:
        """Create a chevron/arrow shape."""
        path = [
            PathSegment(PathCommand.MOVE_TO, [PathPoint(0.0, 0.0)]),
            PathSegment(PathCommand.LINE_TO, [PathPoint(1.0 - point_depth, 0.0)]),
            PathSegment(PathCommand.LINE_TO, [PathPoint(1.0, 0.5)]),  # Point
            PathSegment(PathCommand.LINE_TO, [PathPoint(1.0 - point_depth, 1.0)]),
            PathSegment(PathCommand.LINE_TO, [PathPoint(0.0, 1.0)]),
            PathSegment(PathCommand.LINE_TO, [PathPoint(point_depth, 0.5)]),  # Indent
            PathSegment(PathCommand.CLOSE),
        ]

        return LearnedShape(
            id="chevron_basic",
            name="Chevron",
            description="Arrow-like chevron shape",
            category="arrow",
            path=path,
            aspect_ratio=2.0,
            tags=["chevron", "arrow", "process"],
        )

    @staticmethod
    def create_hexagon() -> LearnedShape:
        """Create a regular hexagon."""
        # Regular hexagon vertices
        points = []
        for i in range(6):
            angle = math.pi / 6 + i * math.pi / 3  # Start at 30 degrees
            x = 0.5 + 0.5 * math.cos(angle)
            y = 0.5 + 0.5 * math.sin(angle)
            points.append(PathPoint(x, y))

        path = [PathSegment(PathCommand.MOVE_TO, [points[0]])]
        for pt in points[1:]:
            path.append(PathSegment(PathCommand.LINE_TO, [pt]))
        path.append(PathSegment(PathCommand.CLOSE))

        return LearnedShape(
            id="hexagon_regular",
            name="Regular Hexagon",
            description="Six-sided polygon",
            category="basic",
            path=path,
            aspect_ratio=1.155,  # Width/height for regular hexagon
            tags=["hexagon", "polygon", "geometric"],
        )

    @staticmethod
    def create_rounded_rectangle(
        corner_radius: float = 0.1,  # As ratio of smaller dimension
    ) -> LearnedShape:
        """Create a rectangle with rounded corners using curves."""
        r = corner_radius

        # Using cubic bezier for corners (approximation of arc)
        k = 0.5523  # Magic number for bezier approximation of circle

        path = [
            # Start at top-left, after corner
            PathSegment(PathCommand.MOVE_TO, [PathPoint(r, 0.0)]),
            # Top edge
            PathSegment(PathCommand.LINE_TO, [PathPoint(1.0 - r, 0.0)]),
            # Top-right corner
            PathSegment(PathCommand.CURVE_TO, [
                PathPoint(1.0 - r + r * k, 0.0),
                PathPoint(1.0, r - r * k),
                PathPoint(1.0, r),
            ]),
            # Right edge
            PathSegment(PathCommand.LINE_TO, [PathPoint(1.0, 1.0 - r)]),
            # Bottom-right corner
            PathSegment(PathCommand.CURVE_TO, [
                PathPoint(1.0, 1.0 - r + r * k),
                PathPoint(1.0 - r + r * k, 1.0),
                PathPoint(1.0 - r, 1.0),
            ]),
            # Bottom edge
            PathSegment(PathCommand.LINE_TO, [PathPoint(r, 1.0)]),
            # Bottom-left corner
            PathSegment(PathCommand.CURVE_TO, [
                PathPoint(r - r * k, 1.0),
                PathPoint(0.0, 1.0 - r + r * k),
                PathPoint(0.0, 1.0 - r),
            ]),
            # Left edge
            PathSegment(PathCommand.LINE_TO, [PathPoint(0.0, r)]),
            # Top-left corner
            PathSegment(PathCommand.CURVE_TO, [
                PathPoint(0.0, r - r * k),
                PathPoint(r - r * k, 0.0),
                PathPoint(r, 0.0),
            ]),
            PathSegment(PathCommand.CLOSE),
        ]

        return LearnedShape(
            id="rounded_rect",
            name="Rounded Rectangle",
            description="Rectangle with rounded corners",
            category="basic",
            path=path,
            aspect_ratio=1.5,
            tags=["rectangle", "rounded", "basic"],
        )


# =============================================================================
# SHAPE LIBRARY - Store and manage shapes
# =============================================================================

class ShapeLibrary:
    """
    Manages a library of learned and generated shapes.

    Provides storage, retrieval, and search functionality.
    """

    def __init__(self, storage_path: Optional[str] = None):
        self.shapes: Dict[str, LearnedShape] = {}
        self.storage_path = Path(storage_path) if storage_path else None

        # Load from storage if exists
        if self.storage_path and self.storage_path.exists():
            self.load()

        # Add built-in shapes
        self._register_builtin_shapes()

    def _register_builtin_shapes(self):
        """Register built-in generated shapes."""
        # Basic shapes
        self.add(ShapeGenerator.create_triangle(pointed_top=True))
        self.add(ShapeGenerator.create_triangle(pointed_top=False))
        self.add(ShapeGenerator.create_hexagon())
        self.add(ShapeGenerator.create_chevron())
        self.add(ShapeGenerator.create_rounded_rectangle())

        # Pyramid segments (for different pyramid sizes)
        for num_levels in [3, 4, 5, 6]:
            for shape in ShapeGenerator.create_full_pyramid(num_levels):
                self.add(shape)

    def add(self, shape: LearnedShape) -> None:
        """Add a shape to the library."""
        self.shapes[shape.id] = shape

    def get(self, shape_id: str) -> Optional[LearnedShape]:
        """Get a shape by ID."""
        return self.shapes.get(shape_id)

    def find_by_category(self, category: str) -> List[LearnedShape]:
        """Find all shapes in a category."""
        return [s for s in self.shapes.values() if s.category == category]

    def find_by_tag(self, tag: str) -> List[LearnedShape]:
        """Find all shapes with a specific tag."""
        return [s for s in self.shapes.values() if tag in s.tags]

    def search(self, query: str) -> List[LearnedShape]:
        """Search shapes by name, description, or tags."""
        query_lower = query.lower()
        results = []
        for shape in self.shapes.values():
            if (query_lower in shape.name.lower() or
                query_lower in shape.description.lower() or
                any(query_lower in tag.lower() for tag in shape.tags)):
                results.append(shape)
        return results

    def get_pyramid_shapes(self, num_levels: int) -> List[LearnedShape]:
        """Get pyramid segment shapes for a specific number of levels."""
        shapes = []
        for level in range(num_levels):
            shape_id = f"pyramid_segment_L{level}_of_{num_levels}"
            shape = self.get(shape_id)
            if shape:
                shapes.append(shape)
            else:
                # Generate if not exists
                shape = ShapeGenerator.create_pyramid_segment(level, num_levels)
                self.add(shape)
                shapes.append(shape)
        return shapes

    def learn_from_template(
        self,
        pptx_path: str,
        category: str = "custom",
    ) -> List[LearnedShape]:
        """
        Learn shapes from a PPTX template and add to library.

        Returns the newly learned shapes.
        """
        extractor = ShapeExtractor()
        shapes = extractor.extract_from_pptx(pptx_path, category=category)

        for shape in shapes:
            self.add(shape)

        # Save if storage path is set
        if self.storage_path:
            self.save()

        return shapes

    def save(self) -> None:
        """Save library to storage."""
        if not self.storage_path:
            return

        self.storage_path.parent.mkdir(parents=True, exist_ok=True)

        data = {
            "shapes": {sid: s.to_dict() for sid, s in self.shapes.items()}
        }

        with open(self.storage_path, 'w') as f:
            json.dump(data, f, indent=2)

    def load(self) -> None:
        """Load library from storage."""
        if not self.storage_path or not self.storage_path.exists():
            return

        try:
            with open(self.storage_path, 'r') as f:
                data = json.load(f)

            for sid, shape_data in data.get("shapes", {}).items():
                shape = LearnedShape.from_dict(shape_data)
                self.shapes[sid] = shape
        except Exception as e:
            logger.warning(f"Failed to load shape library: {e}")

    def list_categories(self) -> List[str]:
        """List all shape categories."""
        return list(set(s.category for s in self.shapes.values()))

    def stats(self) -> Dict[str, Any]:
        """Get library statistics."""
        categories = {}
        for shape in self.shapes.values():
            categories[shape.category] = categories.get(shape.category, 0) + 1

        return {
            "total_shapes": len(self.shapes),
            "categories": categories,
            "learned_shapes": len([s for s in self.shapes.values() if s.source_template]),
            "generated_shapes": len([s for s in self.shapes.values() if not s.source_template]),
        }


# =============================================================================
# SINGLETON LIBRARY INSTANCE
# =============================================================================

_library_instance: Optional[ShapeLibrary] = None


def get_shape_library(storage_path: Optional[str] = None) -> ShapeLibrary:
    """Get the global shape library instance."""
    global _library_instance
    if _library_instance is None:
        _library_instance = ShapeLibrary(storage_path)
    return _library_instance


# =============================================================================
# CONVENIENCE FUNCTIONS
# =============================================================================

def extract_shapes_from_template(
    pptx_path: str,
    category: str = "custom",
) -> List[LearnedShape]:
    """
    Extract shapes from a PPTX template.

    Convenience function that creates an extractor and runs extraction.
    """
    extractor = ShapeExtractor()
    return extractor.extract_from_pptx(pptx_path, category=category)


def generate_pyramid_shapes(num_levels: int) -> List[LearnedShape]:
    """Generate pyramid segment shapes."""
    return ShapeGenerator.create_full_pyramid(num_levels)
