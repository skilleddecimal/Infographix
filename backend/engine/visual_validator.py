"""
visual_validator.py — Visual quality validation for generated PPTX files.

This module provides visual validation by:
1. Converting PPTX to images (using LibreOffice or similar)
2. Analyzing images with Claude's vision capability
3. Returning structured validation results

Validation checks:
- Text alignment and readability
- Shape alignment and visual balance
- Color contrast and accessibility
- Overall professional appearance
- Context appropriateness
"""

import os
import subprocess
import tempfile
import base64
from pathlib import Path
from typing import Optional, List, Dict, Any
from dataclasses import dataclass, field
from enum import Enum
import logging

logger = logging.getLogger(__name__)


class ValidationSeverity(Enum):
    """Severity levels for validation issues."""
    INFO = "info"
    WARNING = "warning"
    ERROR = "error"


@dataclass
class ValidationIssue:
    """A single validation issue found during review."""
    severity: ValidationSeverity
    category: str  # e.g., "text_alignment", "color_contrast", "layout"
    message: str
    suggestion: Optional[str] = None
    location: Optional[str] = None  # e.g., "slide 1, level 2"


@dataclass
class ValidationResult:
    """Complete validation result for a PPTX file."""
    is_valid: bool
    score: float  # 0-100 quality score
    issues: List[ValidationIssue] = field(default_factory=list)
    summary: str = ""
    details: Dict[str, Any] = field(default_factory=dict)

    def add_issue(self, issue: ValidationIssue) -> None:
        """Add a validation issue."""
        self.issues.append(issue)
        # Recalculate validity
        error_count = sum(1 for i in self.issues if i.severity == ValidationSeverity.ERROR)
        self.is_valid = error_count == 0


class PPTXToImageConverter:
    """
    Converts PPTX files to images for visual validation.

    Supports multiple backends:
    - LibreOffice (soffice) - free, cross-platform
    - PowerPoint COM automation (Windows only)
    - Cloud APIs (future)
    """

    def __init__(self):
        self.libreoffice_path = self._find_libreoffice()

    def _find_libreoffice(self) -> Optional[str]:
        """Find LibreOffice installation."""
        # Common paths on Windows
        possible_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "/usr/bin/soffice",
            "/usr/bin/libreoffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        ]

        for path in possible_paths:
            if os.path.exists(path):
                return path

        # Try finding in PATH
        try:
            result = subprocess.run(
                ["where" if os.name == "nt" else "which", "soffice"],
                capture_output=True,
                text=True
            )
            if result.returncode == 0:
                return result.stdout.strip().split('\n')[0]
        except Exception:
            pass

        return None

    def convert_to_images(
        self,
        pptx_path: str,
        output_dir: Optional[str] = None,
        dpi: int = 150,
    ) -> List[str]:
        """
        Convert PPTX to PNG images.

        Args:
            pptx_path: Path to PPTX file
            output_dir: Output directory for images (uses temp if None)
            dpi: Resolution for output images

        Returns:
            List of paths to generated PNG images
        """
        if output_dir is None:
            output_dir = tempfile.mkdtemp(prefix="pptx_validation_")

        if self.libreoffice_path:
            return self._convert_with_libreoffice(pptx_path, output_dir)
        else:
            logger.warning("LibreOffice not found. Visual validation limited.")
            return []

    def _convert_with_libreoffice(
        self,
        pptx_path: str,
        output_dir: str,
    ) -> List[str]:
        """Convert using LibreOffice headless mode."""
        try:
            # Convert to PDF first (more reliable), then to images
            cmd = [
                self.libreoffice_path,
                "--headless",
                "--convert-to", "png",
                "--outdir", output_dir,
                pptx_path
            ]

            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=60
            )

            if result.returncode != 0:
                logger.error(f"LibreOffice conversion failed: {result.stderr}")
                return []

            # Find generated images
            images = list(Path(output_dir).glob("*.png"))
            return [str(img) for img in sorted(images)]

        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timed out")
            return []
        except Exception as e:
            logger.error(f"LibreOffice conversion error: {e}")
            return []


class VisualValidator:
    """
    Validates PPTX visual quality using AI vision.

    Uses Claude's vision capability to analyze rendered slides
    and identify potential issues with alignment, readability,
    color contrast, and overall professional appearance.
    """

    def __init__(self, anthropic_client=None):
        """
        Initialize validator.

        Args:
            anthropic_client: Anthropic client for vision API calls.
                            If None, validation will be limited to basic checks.
        """
        self.anthropic_client = anthropic_client
        self.converter = PPTXToImageConverter()

    def validate_pptx(
        self,
        pptx_path: str,
        context: Optional[str] = None,
        strict: bool = False,
    ) -> ValidationResult:
        """
        Validate a PPTX file for visual quality.

        Args:
            pptx_path: Path to the PPTX file
            context: Original prompt/context for contextual validation
            strict: If True, applies stricter validation rules

        Returns:
            ValidationResult with issues and score
        """
        result = ValidationResult(is_valid=True, score=100.0)

        # Check file exists
        if not os.path.exists(pptx_path):
            result.add_issue(ValidationIssue(
                severity=ValidationSeverity.ERROR,
                category="file",
                message="PPTX file not found",
            ))
            result.score = 0
            return result

        # Convert to images
        images = self.converter.convert_to_images(pptx_path)

        if not images:
            # Fallback to basic validation without images
            return self._validate_without_images(pptx_path, result)

        # Validate each slide image
        for i, image_path in enumerate(images):
            slide_issues = self._validate_slide_image(
                image_path,
                slide_num=i + 1,
                context=context,
                strict=strict,
            )
            for issue in slide_issues:
                result.add_issue(issue)

        # Calculate final score
        result.score = self._calculate_score(result.issues)
        result.summary = self._generate_summary(result)

        # Cleanup temp images
        for img in images:
            try:
                os.remove(img)
            except Exception:
                pass

        return result

    def _validate_slide_image(
        self,
        image_path: str,
        slide_num: int,
        context: Optional[str] = None,
        strict: bool = False,
    ) -> List[ValidationIssue]:
        """Validate a single slide image using vision AI."""
        issues = []

        if not self.anthropic_client:
            logger.warning("No Anthropic client - skipping AI vision validation")
            return issues

        try:
            # Read and encode image
            with open(image_path, "rb") as f:
                image_data = base64.standard_b64encode(f.read()).decode("utf-8")

            # Build validation prompt
            prompt = self._build_validation_prompt(context, strict)

            # Call Claude's vision API
            response = self.anthropic_client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=1024,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image",
                                "source": {
                                    "type": "base64",
                                    "media_type": "image/png",
                                    "data": image_data,
                                },
                            },
                            {
                                "type": "text",
                                "text": prompt,
                            },
                        ],
                    }
                ],
            )

            # Parse response into issues
            issues = self._parse_validation_response(
                response.content[0].text,
                slide_num,
            )

        except Exception as e:
            logger.error(f"Vision validation error: {e}")
            issues.append(ValidationIssue(
                severity=ValidationSeverity.WARNING,
                category="validation",
                message=f"Could not perform AI vision validation: {str(e)}",
            ))

        return issues

    def _build_validation_prompt(
        self,
        context: Optional[str] = None,
        strict: bool = False,
    ) -> str:
        """Build the validation prompt for Claude."""
        base_prompt = """Analyze this PowerPoint slide image for visual quality issues.

Check for:
1. TEXT ALIGNMENT: Is text properly centered within shapes? Is text readable and not cut off?
2. SHAPE ALIGNMENT: Are shapes properly aligned? Do edges line up correctly?
3. COLOR CONTRAST: Is there sufficient contrast between text and background?
4. VISUAL BALANCE: Is the layout balanced and professional looking?
5. SPACING: Is there appropriate spacing between elements?

"""
        if context:
            base_prompt += f"""
The original request was: "{context}"
Check if the diagram appropriately represents this request.

"""

        strictness = "Apply strict professional standards." if strict else "Apply reasonable professional standards."

        base_prompt += f"""{strictness}

Respond in this exact format:
ISSUES:
- [SEVERITY:category] Description of issue | Suggestion to fix
- [SEVERITY:category] Description of issue | Suggestion to fix

SCORE: X/100

Where SEVERITY is one of: ERROR, WARNING, INFO
And category is one of: text_alignment, shape_alignment, color_contrast, visual_balance, spacing, context

If no issues found, respond with:
ISSUES: None
SCORE: 100/100
"""
        return base_prompt

    def _parse_validation_response(
        self,
        response: str,
        slide_num: int,
    ) -> List[ValidationIssue]:
        """Parse Claude's response into ValidationIssues."""
        issues = []

        lines = response.strip().split('\n')
        in_issues_section = False

        for line in lines:
            line = line.strip()

            if line.startswith("ISSUES:"):
                in_issues_section = True
                if "None" in line:
                    break
                continue

            if line.startswith("SCORE:"):
                break

            if in_issues_section and line.startswith("- ["):
                try:
                    # Parse: - [SEVERITY:category] Message | Suggestion
                    bracket_end = line.index("]")
                    severity_cat = line[3:bracket_end]
                    severity_str, category = severity_cat.split(":")

                    rest = line[bracket_end + 1:].strip()
                    if "|" in rest:
                        message, suggestion = rest.split("|", 1)
                    else:
                        message = rest
                        suggestion = None

                    severity = ValidationSeverity[severity_str.upper()]

                    issues.append(ValidationIssue(
                        severity=severity,
                        category=category.strip(),
                        message=message.strip(),
                        suggestion=suggestion.strip() if suggestion else None,
                        location=f"slide {slide_num}",
                    ))
                except Exception as e:
                    logger.warning(f"Could not parse validation line: {line} - {e}")

        return issues

    def _validate_without_images(
        self,
        pptx_path: str,
        result: ValidationResult,
    ) -> ValidationResult:
        """Perform basic validation without image conversion."""
        result.add_issue(ValidationIssue(
            severity=ValidationSeverity.INFO,
            category="validation",
            message="Visual validation skipped (LibreOffice not available)",
            suggestion="Install LibreOffice for full visual validation",
        ))

        # Basic file checks
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)

            # Check slide count
            if len(prs.slides) == 0:
                result.add_issue(ValidationIssue(
                    severity=ValidationSeverity.ERROR,
                    category="content",
                    message="Presentation has no slides",
                ))

            # Check for shapes
            for i, slide in enumerate(prs.slides):
                if len(slide.shapes) == 0:
                    result.add_issue(ValidationIssue(
                        severity=ValidationSeverity.WARNING,
                        category="content",
                        message=f"Slide {i + 1} has no shapes",
                    ))

            result.score = self._calculate_score(result.issues)
            result.summary = "Basic validation completed (visual checks skipped)"

        except Exception as e:
            result.add_issue(ValidationIssue(
                severity=ValidationSeverity.ERROR,
                category="file",
                message=f"Could not read PPTX file: {e}",
            ))
            result.score = 0

        return result

    def _calculate_score(self, issues: List[ValidationIssue]) -> float:
        """Calculate quality score based on issues."""
        score = 100.0

        for issue in issues:
            if issue.severity == ValidationSeverity.ERROR:
                score -= 25
            elif issue.severity == ValidationSeverity.WARNING:
                score -= 10
            elif issue.severity == ValidationSeverity.INFO:
                score -= 2

        return max(0, score)

    def _generate_summary(self, result: ValidationResult) -> str:
        """Generate a human-readable summary."""
        error_count = sum(1 for i in result.issues if i.severity == ValidationSeverity.ERROR)
        warning_count = sum(1 for i in result.issues if i.severity == ValidationSeverity.WARNING)

        if result.score >= 90:
            quality = "excellent"
        elif result.score >= 70:
            quality = "good"
        elif result.score >= 50:
            quality = "acceptable"
        else:
            quality = "needs improvement"

        summary = f"Visual quality is {quality} (score: {result.score:.0f}/100)."

        if error_count > 0 or warning_count > 0:
            summary += f" Found {error_count} errors and {warning_count} warnings."

        return summary


# Convenience function for quick validation
def validate_pptx(
    pptx_path: str,
    context: Optional[str] = None,
    anthropic_client=None,
) -> ValidationResult:
    """
    Quick validation of a PPTX file.

    Args:
        pptx_path: Path to PPTX file
        context: Original prompt for context validation
        anthropic_client: Anthropic client for vision API

    Returns:
        ValidationResult
    """
    validator = VisualValidator(anthropic_client)
    return validator.validate_pptx(pptx_path, context)


# Integration with API routes
async def validate_before_delivery(
    pptx_path: str,
    context: str,
    anthropic_client,
    min_score: float = 70.0,
) -> Dict[str, Any]:
    """
    Validate PPTX before delivering to user.

    Returns dict with:
    - approved: bool - whether file meets quality threshold
    - result: ValidationResult details
    - suggestions: list of improvement suggestions
    """
    validator = VisualValidator(anthropic_client)
    result = validator.validate_pptx(pptx_path, context)

    approved = result.score >= min_score and result.is_valid

    suggestions = [
        issue.suggestion
        for issue in result.issues
        if issue.suggestion and issue.severity in [ValidationSeverity.ERROR, ValidationSeverity.WARNING]
    ]

    return {
        "approved": approved,
        "score": result.score,
        "summary": result.summary,
        "issues": [
            {
                "severity": issue.severity.value,
                "category": issue.category,
                "message": issue.message,
                "suggestion": issue.suggestion,
                "location": issue.location,
            }
            for issue in result.issues
        ],
        "suggestions": suggestions,
    }
