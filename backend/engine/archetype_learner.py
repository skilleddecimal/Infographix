"""
archetype_learner.py — Machine learning component for discovering new archetypes.

The ArchetypeLearner analyzes PPTX templates to discover layout patterns
and automatically generate ArchetypeRules that can be used without code changes.

Pipeline:
1. PPTX Template → Shape Extractor (extracts shapes, positions, colors)
2. Shape Data → Pattern Analyzer (detects arrangement: grid, radial, stack, etc.)
3. Pattern → Rule Generator (creates ArchetypeRules JSON)
4. Rules → Learned Archetype (usable immediately)

Key features:
- Detects layout patterns (grid, stack, radial, flow, tree)
- Learns size progressions (funnel, pyramid)
- Captures color schemes and styling
- Generates confidence scores
"""

from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional, Tuple
from enum import Enum
import math
from collections import defaultdict

from .archetype_rules import (
    ArchetypeRules,
    ElementTemplate,
    ConnectorTemplate,
    LayoutConstraint,
    LayoutStrategy,
    LayoutDirection,
    ElementShape,
    PositionRule,
    SizeRule,
    ColorRule,
    ConnectorPattern,
    SizeParams,
    ColorParams,
    LearnedArchetypeResult,
)


# =============================================================================
# DATA CLASSES FOR EXTRACTED SHAPES
# =============================================================================

@dataclass
class ExtractedShape:
    """A shape extracted from a PPTX file."""
    shape_id: str
    shape_type: str                     # "rectangle", "rounded_rect", "ellipse", etc.
    x: float                            # Left edge in inches
    y: float                            # Top edge in inches
    width: float                        # Width in inches
    height: float                       # Height in inches
    fill_color: Optional[str] = None    # Hex color
    stroke_color: Optional[str] = None
    stroke_width: float = 1.0
    text: Optional[str] = None
    font_size: Optional[int] = None
    rotation: float = 0.0               # Rotation in degrees
    z_order: int = 0

    @property
    def center_x(self) -> float:
        return self.x + self.width / 2

    @property
    def center_y(self) -> float:
        return self.y + self.height / 2

    @property
    def area(self) -> float:
        return self.width * self.height


@dataclass
class ExtractedConnector:
    """A connector extracted from a PPTX file."""
    connector_id: str
    start_x: float
    start_y: float
    end_x: float
    end_y: float
    connector_type: str = "straight"    # "straight", "elbow", "curved"
    has_arrow: bool = True
    color: Optional[str] = None


@dataclass
class ExtractedSlide:
    """Extracted data from a single slide."""
    slide_number: int
    shapes: List[ExtractedShape] = field(default_factory=list)
    connectors: List[ExtractedConnector] = field(default_factory=list)
    title: Optional[str] = None
    background_color: str = "#FFFFFF"


# =============================================================================
# PATTERN DETECTION
# =============================================================================

class DetectedPattern(Enum):
    """Detected layout patterns."""
    GRID = "grid"
    STACK_VERTICAL = "stack_vertical"
    STACK_HORIZONTAL = "stack_horizontal"
    RADIAL = "radial"
    TREE = "tree"
    FLOW_HORIZONTAL = "flow_horizontal"
    FLOW_VERTICAL = "flow_vertical"
    SCATTERED = "scattered"             # No clear pattern
    UNKNOWN = "unknown"


@dataclass
class PatternAnalysis:
    """Result of pattern analysis."""
    pattern: DetectedPattern
    confidence: float                   # 0.0 to 1.0
    details: Dict[str, Any] = field(default_factory=dict)

    # Pattern-specific metrics
    grid_rows: int = 0
    grid_cols: int = 0
    radial_center: Optional[Tuple[float, float]] = None
    radial_radius: float = 0.0
    stack_direction: str = "vertical"
    size_progression: Optional[float] = None  # Ratio between consecutive elements


# =============================================================================
# ARCHETYPE LEARNER
# =============================================================================

class ArchetypeLearner:
    """
    Learns archetype rules from PPTX templates.

    Usage:
        learner = ArchetypeLearner()
        result = learner.learn_from_pptx("template.pptx")
        if result.confidence_score > 0.7:
            # Use result.rules as a new archetype
            save_rules(result.rules, "learned/my_archetype.json")
    """

    def __init__(self):
        """Initialize the learner."""
        pass

    def learn_from_pptx(
        self,
        pptx_path: str,
        archetype_name: Optional[str] = None,
    ) -> LearnedArchetypeResult:
        """
        Learn archetype rules from a PPTX template.

        Args:
            pptx_path: Path to the PPTX file
            archetype_name: Optional name for the learned archetype

        Returns:
            LearnedArchetypeResult with rules and metadata
        """
        # Extract shapes from PPTX
        slides = self._extract_shapes(pptx_path)

        if not slides:
            return LearnedArchetypeResult(
                rules=self._create_fallback_rules(archetype_name or "unknown"),
                source_file=pptx_path,
                confidence_score=0.0,
                warnings=["Could not extract any slides from PPTX"],
            )

        # Use the first slide with shapes
        main_slide = None
        for slide in slides:
            if slide.shapes:
                main_slide = slide
                break

        if not main_slide:
            return LearnedArchetypeResult(
                rules=self._create_fallback_rules(archetype_name or "unknown"),
                source_file=pptx_path,
                confidence_score=0.0,
                warnings=["No shapes found in PPTX"],
            )

        # Analyze pattern
        pattern_analysis = self._analyze_pattern(main_slide.shapes)

        # Detect size progression
        size_progression = self._detect_size_progression(main_slide.shapes)

        # Extract color scheme
        colors = self._extract_color_scheme(main_slide.shapes)

        # Generate rules
        rules = self._generate_rules(
            shapes=main_slide.shapes,
            connectors=main_slide.connectors,
            pattern=pattern_analysis,
            size_progression=size_progression,
            colors=colors,
            archetype_name=archetype_name,
        )

        return LearnedArchetypeResult(
            rules=rules,
            source_file=pptx_path,
            extraction_metadata={
                "slide_count": len(slides),
                "shape_count": len(main_slide.shapes),
                "connector_count": len(main_slide.connectors),
            },
            confidence_score=pattern_analysis.confidence,
            detected_pattern=pattern_analysis.pattern.value,
            element_count=len(main_slide.shapes),
        )

    def _extract_shapes(self, pptx_path: str) -> List[ExtractedSlide]:
        """Extract shapes from PPTX file."""
        slides = []

        try:
            from pptx import Presentation
            from pptx.util import Emu
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(pptx_path)

            for slide_num, slide in enumerate(prs.slides, 1):
                extracted = ExtractedSlide(slide_number=slide_num)

                for shape in slide.shapes:
                    # Skip placeholders and text boxes with no shape
                    if not hasattr(shape, 'left') or shape.left is None:
                        continue

                    # Convert EMU to inches
                    x = shape.left / 914400
                    y = shape.top / 914400
                    width = shape.width / 914400
                    height = shape.height / 914400

                    # Determine shape type
                    shape_type = "rectangle"
                    if hasattr(shape, 'shape_type'):
                        if shape.shape_type == MSO_SHAPE_TYPE.OVAL:
                            shape_type = "ellipse"
                        elif shape.shape_type == MSO_SHAPE_TYPE.ROUNDED_RECTANGLE:
                            shape_type = "rounded_rect"

                    # Extract colors
                    fill_color = None
                    if hasattr(shape, 'fill') and shape.fill:
                        try:
                            if shape.fill.type is not None:
                                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                                    rgb = shape.fill.fore_color.rgb
                                    if rgb:
                                        fill_color = f"#{rgb}"
                        except:
                            pass

                    # Extract text
                    text = None
                    font_size = None
                    if hasattr(shape, 'text_frame'):
                        try:
                            text = shape.text_frame.text
                            if shape.text_frame.paragraphs:
                                for para in shape.text_frame.paragraphs:
                                    if para.runs:
                                        font = para.runs[0].font
                                        if font.size:
                                            font_size = int(font.size.pt)
                                        break
                        except:
                            pass

                    extracted_shape = ExtractedShape(
                        shape_id=f"shape_{slide_num}_{len(extracted.shapes)}",
                        shape_type=shape_type,
                        x=x,
                        y=y,
                        width=width,
                        height=height,
                        fill_color=fill_color,
                        text=text,
                        font_size=font_size,
                    )
                    extracted.shapes.append(extracted_shape)

                if extracted.shapes:
                    slides.append(extracted)

        except Exception as e:
            # Return empty list on error
            pass

        return slides

    def _analyze_pattern(self, shapes: List[ExtractedShape]) -> PatternAnalysis:
        """Analyze the layout pattern of shapes."""
        if len(shapes) < 2:
            return PatternAnalysis(
                pattern=DetectedPattern.UNKNOWN,
                confidence=0.0,
            )

        # Try each pattern detector
        patterns = [
            self._detect_grid_pattern(shapes),
            self._detect_stack_pattern(shapes),
            self._detect_radial_pattern(shapes),
            self._detect_flow_pattern(shapes),
        ]

        # Return the pattern with highest confidence
        best = max(patterns, key=lambda p: p.confidence)

        if best.confidence < 0.3:
            return PatternAnalysis(
                pattern=DetectedPattern.SCATTERED,
                confidence=0.3,
            )

        return best

    def _detect_grid_pattern(self, shapes: List[ExtractedShape]) -> PatternAnalysis:
        """Detect if shapes are arranged in a grid."""
        if len(shapes) < 4:
            return PatternAnalysis(pattern=DetectedPattern.GRID, confidence=0.0)

        # Get unique x and y positions (with tolerance)
        tolerance = 0.3  # inches
        x_positions = self._cluster_values([s.center_x for s in shapes], tolerance)
        y_positions = self._cluster_values([s.center_y for s in shapes], tolerance)

        cols = len(x_positions)
        rows = len(y_positions)

        # Check if grid is well-formed
        expected_count = rows * cols
        actual_count = len(shapes)

        # Allow for some empty cells
        fill_ratio = actual_count / max(expected_count, 1)

        if rows >= 2 and cols >= 2 and fill_ratio > 0.6:
            confidence = min(0.9, fill_ratio * 0.9)
            return PatternAnalysis(
                pattern=DetectedPattern.GRID,
                confidence=confidence,
                grid_rows=rows,
                grid_cols=cols,
                details={"fill_ratio": fill_ratio},
            )

        return PatternAnalysis(pattern=DetectedPattern.GRID, confidence=0.0)

    def _detect_stack_pattern(self, shapes: List[ExtractedShape]) -> PatternAnalysis:
        """Detect if shapes are stacked vertically or horizontally."""
        if len(shapes) < 2:
            return PatternAnalysis(pattern=DetectedPattern.STACK_VERTICAL, confidence=0.0)

        # Check vertical alignment (same x center)
        x_centers = [s.center_x for s in shapes]
        x_variance = self._calculate_variance(x_centers)

        # Check horizontal alignment (same y center)
        y_centers = [s.center_y for s in shapes]
        y_variance = self._calculate_variance(y_centers)

        # Normalize variances by content area
        x_variance_norm = x_variance / 100  # Rough normalization
        y_variance_norm = y_variance / 100

        if x_variance_norm < 0.1:  # Low x variance = vertically stacked
            # Check if y positions are evenly spaced
            y_sorted = sorted(y_centers)
            spacings = [y_sorted[i+1] - y_sorted[i] for i in range(len(y_sorted)-1)]
            if spacings:
                spacing_variance = self._calculate_variance(spacings)
                regularity = 1.0 / (1.0 + spacing_variance)
                confidence = min(0.9, 0.5 + regularity * 0.4)
                return PatternAnalysis(
                    pattern=DetectedPattern.STACK_VERTICAL,
                    confidence=confidence,
                    stack_direction="vertical",
                )

        if y_variance_norm < 0.1:  # Low y variance = horizontally stacked
            x_sorted = sorted(x_centers)
            spacings = [x_sorted[i+1] - x_sorted[i] for i in range(len(x_sorted)-1)]
            if spacings:
                spacing_variance = self._calculate_variance(spacings)
                regularity = 1.0 / (1.0 + spacing_variance)
                confidence = min(0.9, 0.5 + regularity * 0.4)
                return PatternAnalysis(
                    pattern=DetectedPattern.STACK_HORIZONTAL,
                    confidence=confidence,
                    stack_direction="horizontal",
                )

        return PatternAnalysis(pattern=DetectedPattern.STACK_VERTICAL, confidence=0.0)

    def _detect_radial_pattern(self, shapes: List[ExtractedShape]) -> PatternAnalysis:
        """Detect if shapes are arranged radially around a center."""
        if len(shapes) < 3:
            return PatternAnalysis(pattern=DetectedPattern.RADIAL, confidence=0.0)

        # Calculate centroid
        cx = sum(s.center_x for s in shapes) / len(shapes)
        cy = sum(s.center_y for s in shapes) / len(shapes)

        # Calculate distances from centroid
        distances = []
        for s in shapes:
            d = math.sqrt((s.center_x - cx)**2 + (s.center_y - cy)**2)
            distances.append(d)

        # Check if distances are similar (excluding potential center element)
        sorted_distances = sorted(distances)

        # Exclude the smallest distance (potential center)
        outer_distances = sorted_distances[1:] if len(sorted_distances) > 3 else sorted_distances

        if outer_distances:
            distance_variance = self._calculate_variance(outer_distances)
            avg_distance = sum(outer_distances) / len(outer_distances)

            # Coefficient of variation
            cv = math.sqrt(distance_variance) / max(avg_distance, 0.01)

            if cv < 0.3:  # Low variation = radial pattern
                confidence = max(0.0, 0.9 - cv)
                return PatternAnalysis(
                    pattern=DetectedPattern.RADIAL,
                    confidence=confidence,
                    radial_center=(cx, cy),
                    radial_radius=avg_distance,
                )

        return PatternAnalysis(pattern=DetectedPattern.RADIAL, confidence=0.0)

    def _detect_flow_pattern(self, shapes: List[ExtractedShape]) -> PatternAnalysis:
        """Detect if shapes form a flow (sequential with possible wrapping)."""
        # Similar to stack but with connectors indicating sequence
        stack_result = self._detect_stack_pattern(shapes)

        if stack_result.confidence > 0.5:
            if stack_result.stack_direction == "horizontal":
                return PatternAnalysis(
                    pattern=DetectedPattern.FLOW_HORIZONTAL,
                    confidence=stack_result.confidence * 0.9,
                )
            else:
                return PatternAnalysis(
                    pattern=DetectedPattern.FLOW_VERTICAL,
                    confidence=stack_result.confidence * 0.9,
                )

        return PatternAnalysis(pattern=DetectedPattern.FLOW_HORIZONTAL, confidence=0.0)

    def _detect_size_progression(self, shapes: List[ExtractedShape]) -> Optional[float]:
        """
        Detect if shapes have a progressive size change.

        Returns the average ratio between consecutive elements, or None if no clear progression.
        """
        if len(shapes) < 3:
            return None

        # Sort by y position (for vertical progressions like funnel/pyramid)
        sorted_shapes = sorted(shapes, key=lambda s: s.center_y)

        widths = [s.width for s in sorted_shapes]
        ratios = []
        for i in range(len(widths) - 1):
            if widths[i] > 0:
                ratios.append(widths[i + 1] / widths[i])

        if not ratios:
            return None

        # Check if ratios are consistent
        ratio_variance = self._calculate_variance(ratios)
        avg_ratio = sum(ratios) / len(ratios)

        # If variance is low and ratio is not ~1.0, we have a progression
        if ratio_variance < 0.1 and abs(avg_ratio - 1.0) > 0.1:
            return avg_ratio

        return None

    def _extract_color_scheme(self, shapes: List[ExtractedShape]) -> List[str]:
        """Extract the color scheme from shapes."""
        colors = []
        for s in shapes:
            if s.fill_color and s.fill_color not in colors:
                colors.append(s.fill_color)

        return colors[:4]  # Return up to 4 colors

    def _generate_rules(
        self,
        shapes: List[ExtractedShape],
        connectors: List[ExtractedConnector],
        pattern: PatternAnalysis,
        size_progression: Optional[float],
        colors: List[str],
        archetype_name: Optional[str],
    ) -> ArchetypeRules:
        """Generate archetype rules from analyzed data."""
        # Determine layout strategy
        strategy_map = {
            DetectedPattern.GRID: LayoutStrategy.GRID,
            DetectedPattern.STACK_VERTICAL: LayoutStrategy.STACK,
            DetectedPattern.STACK_HORIZONTAL: LayoutStrategy.STACK,
            DetectedPattern.RADIAL: LayoutStrategy.RADIAL,
            DetectedPattern.FLOW_HORIZONTAL: LayoutStrategy.FLOW,
            DetectedPattern.FLOW_VERTICAL: LayoutStrategy.FLOW,
            DetectedPattern.TREE: LayoutStrategy.TREE,
            DetectedPattern.SCATTERED: LayoutStrategy.FREEFORM,
            DetectedPattern.UNKNOWN: LayoutStrategy.FREEFORM,
        }
        strategy = strategy_map.get(pattern.pattern, LayoutStrategy.FREEFORM)

        # Determine direction
        direction_map = {
            DetectedPattern.STACK_VERTICAL: LayoutDirection.VERTICAL,
            DetectedPattern.STACK_HORIZONTAL: LayoutDirection.HORIZONTAL,
            DetectedPattern.FLOW_HORIZONTAL: LayoutDirection.HORIZONTAL,
            DetectedPattern.FLOW_VERTICAL: LayoutDirection.VERTICAL,
            DetectedPattern.RADIAL: LayoutDirection.RADIAL,
        }
        direction = direction_map.get(pattern.pattern, LayoutDirection.VERTICAL)

        # Determine element shape
        shape_types = [s.shape_type for s in shapes]
        most_common_shape = max(set(shape_types), key=shape_types.count) if shape_types else "rounded_rect"
        element_shape_map = {
            "rectangle": ElementShape.RECTANGLE,
            "rounded_rect": ElementShape.ROUNDED_RECT,
            "ellipse": ElementShape.ELLIPSE,
            "trapezoid": ElementShape.TRAPEZOID,
        }
        element_shape = element_shape_map.get(most_common_shape, ElementShape.ROUNDED_RECT)

        # Build size params
        size_rule = SizeRule.TEXT_FIT
        size_params = SizeParams()
        if size_progression is not None:
            size_rule = SizeRule.PROGRESSIVE
            size_params.width_progression = size_progression

        # Build color params
        color_params = ColorParams()
        if colors:
            color_params.uniform_color = colors[0]

        # Create element template
        element_template = ElementTemplate(
            element_type=element_shape,
            position_rule=self._get_position_rule(pattern.pattern),
            size_rule=size_rule,
            size_params=size_params,
            color_rule=ColorRule.PALETTE_SEQUENCE if len(colors) > 1 else ColorRule.UNIFORM,
            color_params=color_params,
        )

        # Create connector template
        connector_pattern = ConnectorPattern.NONE
        if connectors:
            connector_pattern = ConnectorPattern.SEQUENTIAL

        connector_template = ConnectorTemplate(
            pattern=connector_pattern,
            style="arrow" if any(c.has_arrow for c in connectors) else "plain",
        )

        # Build strategy-specific params
        grid_params = {"columns": pattern.grid_cols, "rows": pattern.grid_rows} if pattern.grid_cols > 0 else {}
        stack_params = {}
        if size_progression is not None:
            if size_progression < 1.0:
                stack_params = {"top_width_ratio": 0.9, "bottom_width_ratio": 0.25}
            else:
                stack_params = {"top_width_ratio": 0.25, "bottom_width_ratio": 0.9, "direction": "top_narrow"}

        radial_params = {}
        if pattern.radial_center:
            radial_params = {
                "center_element": True,
                "radius_ratio": pattern.radial_radius / 3.5,  # Normalize
            }

        # Generate archetype ID
        archetype_id = archetype_name or f"learned_{pattern.pattern.value}"

        return ArchetypeRules(
            archetype_id=archetype_id,
            display_name=archetype_id.replace("_", " ").title(),
            description=f"Learned archetype from template ({pattern.pattern.value} pattern)",
            layout_strategy=strategy,
            primary_direction=direction,
            element_template=element_template,
            connector_template=connector_template,
            grid_params=grid_params or {},
            stack_params=stack_params or {},
            radial_params=radial_params or {},
            min_elements=max(2, len(shapes) - 2),
            max_elements=min(50, len(shapes) + 10),
            learned_from=[archetype_name or "template"],
            confidence_score=pattern.confidence,
            category=self._infer_category(pattern.pattern),
        )

    def _get_position_rule(self, pattern: DetectedPattern) -> PositionRule:
        """Get position rule for pattern."""
        rule_map = {
            DetectedPattern.GRID: PositionRule.GRID_FILL,
            DetectedPattern.STACK_VERTICAL: PositionRule.STACKED_CENTERED,
            DetectedPattern.STACK_HORIZONTAL: PositionRule.STACKED_LEFT,
            DetectedPattern.RADIAL: PositionRule.RADIAL_EVEN,
            DetectedPattern.FLOW_HORIZONTAL: PositionRule.FLOW_LINEAR,
            DetectedPattern.FLOW_VERTICAL: PositionRule.FLOW_LINEAR,
            DetectedPattern.TREE: PositionRule.TREE_BALANCED,
        }
        return rule_map.get(pattern, PositionRule.FREEFORM)

    def _infer_category(self, pattern: DetectedPattern) -> str:
        """Infer archetype category from pattern."""
        category_map = {
            DetectedPattern.GRID: "list",
            DetectedPattern.STACK_VERTICAL: "hierarchy",
            DetectedPattern.STACK_HORIZONTAL: "list",
            DetectedPattern.RADIAL: "relationship",
            DetectedPattern.FLOW_HORIZONTAL: "process",
            DetectedPattern.FLOW_VERTICAL: "process",
            DetectedPattern.TREE: "hierarchy",
        }
        return category_map.get(pattern, "other")

    def _create_fallback_rules(self, name: str) -> ArchetypeRules:
        """Create fallback rules when learning fails."""
        return ArchetypeRules(
            archetype_id=name,
            display_name=name.replace("_", " ").title(),
            description="Fallback rules (learning failed)",
            layout_strategy=LayoutStrategy.FREEFORM,
            confidence_score=0.0,
        )

    # =========================================================================
    # UTILITY METHODS
    # =========================================================================

    def _cluster_values(self, values: List[float], tolerance: float) -> List[float]:
        """Cluster similar values within tolerance."""
        if not values:
            return []

        sorted_vals = sorted(values)
        clusters = [[sorted_vals[0]]]

        for val in sorted_vals[1:]:
            if val - clusters[-1][-1] <= tolerance:
                clusters[-1].append(val)
            else:
                clusters.append([val])

        # Return cluster centers
        return [sum(c) / len(c) for c in clusters]

    def _calculate_variance(self, values: List[float]) -> float:
        """Calculate variance of a list of values."""
        if len(values) < 2:
            return 0.0

        mean = sum(values) / len(values)
        return sum((v - mean) ** 2 for v in values) / len(values)


# =============================================================================
# CONVENIENCE FUNCTIONS
# =============================================================================

def learn_archetype(pptx_path: str, name: Optional[str] = None) -> LearnedArchetypeResult:
    """
    Convenience function to learn an archetype from a PPTX file.

    Args:
        pptx_path: Path to the PPTX template
        name: Optional name for the learned archetype

    Returns:
        LearnedArchetypeResult with rules and metadata
    """
    learner = ArchetypeLearner()
    return learner.learn_from_pptx(pptx_path, name)
