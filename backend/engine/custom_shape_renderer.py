"""
custom_shape_renderer.py — Render custom shapes to PowerPoint.

This module converts LearnedShape objects to DrawingML XML and
adds them to PowerPoint slides. It supports:
- Custom polygon paths (FREEFORM shapes)
- Bezier curves
- Gradient fills
- Shadows and effects

This is the bridge between the shape learning system and PPTX output.
"""

from typing import List, Optional, Tuple, Dict, Any
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
from pptx.shapes.autoshape import Shape

from .shape_learning import (
    LearnedShape,
    PathSegment,
    PathCommand,
    PathPoint,
    ShapeLibrary,
    get_shape_library,
)

# DrawingML namespace
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _emu(inches: float) -> int:
    """Convert inches to EMU."""
    return int(inches * 914400)


class CustomShapeRenderer:
    """
    Renders LearnedShape objects to PowerPoint slides.
    """

    def __init__(self, slide):
        """
        Initialize with a PowerPoint slide.

        Args:
            slide: pptx.slide.Slide object
        """
        self.slide = slide
        self.spTree = slide.shapes._spTree

    def add_shape(
        self,
        shape: LearnedShape,
        left: float,
        top: float,
        width: float,
        height: float,
        fill_color: str = "#4472C4",
        gradient: bool = True,
        gradient_angle: float = 270.0,
        shadow: bool = True,
        shadow_blur_pt: float = 40.0,
        shadow_opacity: float = 0.6,
        stroke_color: Optional[str] = None,
        stroke_width_pt: float = 0.0,
    ) -> None:
        """
        Add a custom shape to the slide.

        Args:
            shape: LearnedShape to render
            left, top: Position in inches
            width, height: Size in inches
            fill_color: Hex color for fill
            gradient: Whether to apply gradient fill
            gradient_angle: Gradient angle in degrees
            shadow: Whether to apply shadow
            shadow_blur_pt: Shadow blur radius in points
            shadow_opacity: Shadow opacity (0-1)
            stroke_color: Optional stroke color (None = no stroke)
            stroke_width_pt: Stroke width in points
        """
        # Convert to EMU
        left_emu = _emu(left)
        top_emu = _emu(top)
        width_emu = _emu(width)
        height_emu = _emu(height)

        # Get next shape ID
        shape_id = len(self.slide.shapes) + 1

        # Build the shape XML
        sp_xml = self._build_shape_xml(
            shape=shape,
            shape_id=shape_id,
            left_emu=left_emu,
            top_emu=top_emu,
            width_emu=width_emu,
            height_emu=height_emu,
            fill_color=fill_color,
            gradient=gradient,
            gradient_angle=gradient_angle,
            shadow=shadow,
            shadow_blur_pt=shadow_blur_pt,
            shadow_opacity=shadow_opacity,
            stroke_color=stroke_color,
            stroke_width_pt=stroke_width_pt,
        )

        # Parse and add to slide
        sp_element = etree.fromstring(sp_xml)
        self.spTree.append(sp_element)

    def _build_shape_xml(
        self,
        shape: LearnedShape,
        shape_id: int,
        left_emu: int,
        top_emu: int,
        width_emu: int,
        height_emu: int,
        fill_color: str,
        gradient: bool,
        gradient_angle: float,
        shadow: bool,
        shadow_blur_pt: float,
        shadow_opacity: float,
        stroke_color: Optional[str],
        stroke_width_pt: float,
    ) -> str:
        """Build the complete shape XML string."""

        # Build path XML
        path_xml = self._build_path_xml(shape, width_emu, height_emu)

        # Build fill XML
        if gradient:
            fill_xml = self._build_gradient_fill_xml(fill_color, gradient_angle)
        else:
            fill_xml = self._build_solid_fill_xml(fill_color)

        # Build stroke XML
        if stroke_color and stroke_width_pt > 0:
            stroke_xml = self._build_stroke_xml(stroke_color, stroke_width_pt)
        else:
            stroke_xml = '<a:ln><a:noFill/></a:ln>'

        # Build effects XML (shadow)
        if shadow:
            effects_xml = self._build_shadow_xml(shadow_blur_pt, shadow_opacity)
        else:
            effects_xml = ''

        # Assemble full shape XML
        shape_xml = f'''
        <p:sp xmlns:a="{A_NS}" xmlns:p="{P_NS}" xmlns:r="{R_NS}">
            <p:nvSpPr>
                <p:cNvPr id="{shape_id}" name="CustomShape {shape_id}"/>
                <p:cNvSpPr/>
                <p:nvPr/>
            </p:nvSpPr>
            <p:spPr>
                <a:xfrm>
                    <a:off x="{left_emu}" y="{top_emu}"/>
                    <a:ext cx="{width_emu}" cy="{height_emu}"/>
                </a:xfrm>
                <a:custGeom>
                    <a:avLst/>
                    <a:gdLst/>
                    <a:ahLst/>
                    <a:cxnLst/>
                    <a:rect l="0" t="0" r="0" b="0"/>
                    <a:pathLst>
                        {path_xml}
                    </a:pathLst>
                </a:custGeom>
                {fill_xml}
                {stroke_xml}
                {effects_xml}
            </p:spPr>
            <p:txBody>
                <a:bodyPr/>
                <a:lstStyle/>
                <a:p>
                    <a:endParaRPr lang="en-US"/>
                </a:p>
            </p:txBody>
        </p:sp>
        '''

        return shape_xml

    def _build_path_xml(
        self,
        shape: LearnedShape,
        width_emu: int,
        height_emu: int,
    ) -> str:
        """Build the path XML for custom geometry."""
        # Use shape dimensions as path dimensions
        path_w = width_emu
        path_h = height_emu

        path_commands = []

        for segment in shape.path:
            if segment.command == PathCommand.MOVE_TO:
                if segment.points:
                    pt = segment.points[0]
                    x, y = pt.to_emu(path_w, path_h)
                    path_commands.append(f'<a:moveTo><a:pt x="{x}" y="{y}"/></a:moveTo>')

            elif segment.command == PathCommand.LINE_TO:
                if segment.points:
                    pt = segment.points[0]
                    x, y = pt.to_emu(path_w, path_h)
                    path_commands.append(f'<a:lnTo><a:pt x="{x}" y="{y}"/></a:lnTo>')

            elif segment.command == PathCommand.CURVE_TO:
                if len(segment.points) >= 3:
                    pts = segment.points[:3]
                    pts_xml = ''.join([
                        f'<a:pt x="{int(p.x * path_w)}" y="{int(p.y * path_h)}"/>'
                        for p in pts
                    ])
                    path_commands.append(f'<a:cubicBezTo>{pts_xml}</a:cubicBezTo>')

            elif segment.command == PathCommand.QUAD_TO:
                if len(segment.points) >= 2:
                    pts = segment.points[:2]
                    pts_xml = ''.join([
                        f'<a:pt x="{int(p.x * path_w)}" y="{int(p.y * path_h)}"/>'
                        for p in pts
                    ])
                    path_commands.append(f'<a:quadBezTo>{pts_xml}</a:quadBezTo>')

            elif segment.command == PathCommand.ARC_TO:
                if segment.arc_params:
                    params = segment.arc_params
                    wR = int(params.get('wR', 0) * path_w)
                    hR = int(params.get('hR', 0) * path_h)
                    stAng = int(params.get('stAng', 0) * 60000)
                    swAng = int(params.get('swAng', 0) * 60000)
                    path_commands.append(
                        f'<a:arcTo wR="{wR}" hR="{hR}" stAng="{stAng}" swAng="{swAng}"/>'
                    )

            elif segment.command == PathCommand.CLOSE:
                path_commands.append('<a:close/>')

        return f'<a:path w="{path_w}" h="{path_h}">{" ".join(path_commands)}</a:path>'

    def _build_solid_fill_xml(self, color: str) -> str:
        """Build solid fill XML."""
        color_hex = color.lstrip('#')
        return f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'

    def _build_gradient_fill_xml(self, base_color: str, angle: float) -> str:
        """Build gradient fill XML with automatic color variation."""
        color_hex = base_color.lstrip('#')

        # Calculate lighter and darker variants
        r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)

        # Lighter (for top)
        lr = min(255, int(r + (255 - r) * 0.3))
        lg = min(255, int(g + (255 - g) * 0.3))
        lb = min(255, int(b + (255 - b) * 0.3))
        light_hex = f"{lr:02X}{lg:02X}{lb:02X}"

        # Darker (for bottom)
        dr = max(0, int(r * 0.7))
        dg = max(0, int(g * 0.7))
        db = max(0, int(b * 0.7))
        dark_hex = f"{dr:02X}{dg:02X}{db:02X}"

        # Angle in PowerPoint units (60000 = 1 degree)
        angle_ppt = int(angle * 60000)

        return f'''
        <a:gradFill rotWithShape="1">
            <a:gsLst>
                <a:gs pos="0">
                    <a:srgbClr val="{light_hex}"/>
                </a:gs>
                <a:gs pos="50000">
                    <a:srgbClr val="{color_hex}"/>
                </a:gs>
                <a:gs pos="100000">
                    <a:srgbClr val="{dark_hex}"/>
                </a:gs>
            </a:gsLst>
            <a:lin ang="{angle_ppt}" scaled="1"/>
        </a:gradFill>
        '''

    def _build_stroke_xml(self, color: str, width_pt: float) -> str:
        """Build stroke/line XML."""
        color_hex = color.lstrip('#')
        width_emu = int(width_pt * 12700)  # 1 point = 12700 EMU

        return f'''
        <a:ln w="{width_emu}">
            <a:solidFill>
                <a:srgbClr val="{color_hex}"/>
            </a:solidFill>
        </a:ln>
        '''

    def _build_shadow_xml(self, blur_pt: float, opacity: float) -> str:
        """Build shadow effects XML."""
        blur_emu = int(blur_pt * 12700)
        dist_emu = int(4 * 12700)  # 4pt distance
        direction = int(135 * 60000)  # 135 degrees (bottom-right)
        alpha = int(opacity * 100000)

        return f'''
        <a:effectLst>
            <a:outerShdw blurRad="{blur_emu}" dist="{dist_emu}" dir="{direction}" algn="ctr" rotWithShape="0">
                <a:srgbClr val="000000">
                    <a:alpha val="{alpha}"/>
                </a:srgbClr>
            </a:outerShdw>
        </a:effectLst>
        '''


def add_custom_shape_to_slide(
    slide,
    shape: LearnedShape,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str = "#4472C4",
    **kwargs,
) -> None:
    """
    Convenience function to add a custom shape to a slide.

    Args:
        slide: pptx.slide.Slide object
        shape: LearnedShape to render
        left, top: Position in inches
        width, height: Size in inches
        fill_color: Hex color
        **kwargs: Additional style options
    """
    renderer = CustomShapeRenderer(slide)
    renderer.add_shape(
        shape=shape,
        left=left,
        top=top,
        width=width,
        height=height,
        fill_color=fill_color,
        **kwargs,
    )


def render_pyramid_to_slide(
    slide,
    num_levels: int,
    left: float,
    top: float,
    width: float,
    height: float,
    colors: List[str],
    level_labels: Optional[List[str]] = None,
    shadow: bool = True,
    shadow_blur_pt: float = 40.0,
) -> None:
    """
    Render a complete pyramid with proper triangular geometry.

    Args:
        slide: pptx.slide.Slide object
        num_levels: Number of pyramid levels
        left, top: Position in inches
        width, height: Size in inches
        colors: List of colors for each level (base to apex)
        level_labels: Optional labels for each level
        shadow: Whether to apply shadows
        shadow_blur_pt: Shadow blur radius
    """
    library = get_shape_library()
    shapes = library.get_pyramid_shapes(num_levels)

    # Calculate level heights
    level_height = height / num_levels
    gap = level_height * 0.05  # 5% gap between levels

    renderer = CustomShapeRenderer(slide)

    for i, shape in enumerate(shapes):
        # Position from bottom to top
        level_top = top + height - (i + 1) * level_height + gap / 2
        level_h = level_height - gap

        # Get color (wrap if not enough colors)
        color = colors[i % len(colors)] if colors else "#4472C4"

        renderer.add_shape(
            shape=shape,
            left=left,
            top=level_top,
            width=width,
            height=level_h,
            fill_color=color,
            gradient=True,
            gradient_angle=270.0,
            shadow=shadow,
            shadow_blur_pt=shadow_blur_pt,
            shadow_opacity=0.6,
        )


def render_learned_shape_to_slide(
    slide,
    shape_id: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str,
    **kwargs,
) -> bool:
    """
    Render a shape from the library by ID.

    Returns True if shape was found and rendered, False otherwise.
    """
    library = get_shape_library()
    shape = library.get(shape_id)

    if not shape:
        return False

    add_custom_shape_to_slide(
        slide=slide,
        shape=shape,
        left=left,
        top=top,
        width=width,
        height=height,
        fill_color=fill_color,
        **kwargs,
    )
    return True
