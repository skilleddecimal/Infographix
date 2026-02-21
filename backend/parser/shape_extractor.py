"""Extract shapes from PowerPoint slides into DSL format."""

import uuid
from typing import Any

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape as PPTXAutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.shapes.shapetree import SlideShapes

from backend.dsl.schema import (
    BoundingBox,
    Effects,
    NoFill,
    Shape,
    ShapeType,
    TextContent,
    TextRun,
    Transform,
)
from backend.parser.style_extractor import StyleExtractor


class ShapeExtractor:
    """Extracts shapes from PPTX slides."""

    def __init__(self) -> None:
        """Initialize the shape extractor."""
        self.style_extractor = StyleExtractor()
        self._z_index_counter = 0

    def extract_shapes(
        self,
        shapes: SlideShapes,
        group_path: list[str] | None = None,
    ) -> list[Shape]:
        """Extract all shapes from a slide.

        Args:
            shapes: SlideShapes collection from python-pptx.
            group_path: Current group nesting path.

        Returns:
            List of DSL Shape objects.
        """
        if group_path is None:
            group_path = ["root"]
            self._z_index_counter = 0

        result: list[Shape] = []

        for pptx_shape in shapes:
            dsl_shape = self._extract_shape(pptx_shape, group_path)
            if dsl_shape:
                result.append(dsl_shape)

        return result

    def _extract_shape(
        self,
        pptx_shape: BaseShape,
        group_path: list[str],
    ) -> Shape | None:
        """Extract a single shape to DSL format.

        Args:
            pptx_shape: The python-pptx shape object.
            group_path: Current group nesting path.

        Returns:
            DSL Shape object or None if unsupported.
        """
        shape_type = self._get_shape_type(pptx_shape)
        if shape_type is None:
            return None

        shape_id = self._generate_id(pptx_shape)
        self._z_index_counter += 1

        # Extract bounding box
        bbox = BoundingBox(
            x=pptx_shape.left,
            y=pptx_shape.top,
            width=pptx_shape.width,
            height=pptx_shape.height,
        )

        # Extract transform
        transform = self._extract_transform(pptx_shape)

        # Base shape properties
        shape_dict: dict[str, Any] = {
            "id": shape_id,
            "type": shape_type,
            "name": pptx_shape.name,
            "group_path": group_path.copy(),
            "z_index": self._z_index_counter,
            "bbox": bbox,
            "transform": transform,
            "effects": Effects(),
        }

        # Type-specific extraction
        if shape_type == ShapeType.AUTO_SHAPE:
            self._extract_auto_shape(pptx_shape, shape_dict)
        elif shape_type == ShapeType.FREEFORM:
            self._extract_freeform(pptx_shape, shape_dict)
        elif shape_type == ShapeType.TEXT:
            self._extract_text_shape(pptx_shape, shape_dict)
        elif shape_type == ShapeType.IMAGE:
            self._extract_image(pptx_shape, shape_dict)
        elif shape_type == ShapeType.GROUP:
            self._extract_group(pptx_shape, shape_dict, group_path)

        return Shape(**shape_dict)

    def _get_shape_type(self, pptx_shape: BaseShape) -> ShapeType | None:
        """Determine the DSL shape type from a python-pptx shape.

        Args:
            pptx_shape: The python-pptx shape object.

        Returns:
            ShapeType enum value or None if unsupported.
        """
        if isinstance(pptx_shape, GroupShape):
            return ShapeType.GROUP
        if isinstance(pptx_shape, Picture):
            return ShapeType.IMAGE

        shape_type = pptx_shape.shape_type

        if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return ShapeType.AUTO_SHAPE
        if shape_type == MSO_SHAPE_TYPE.FREEFORM:
            return ShapeType.FREEFORM
        if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            return ShapeType.TEXT
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            return ShapeType.IMAGE
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            return ShapeType.GROUP
        if shape_type == MSO_SHAPE_TYPE.LINE:
            return ShapeType.CONNECTOR

        # Unsupported types
        return None

    def _generate_id(self, pptx_shape: BaseShape) -> str:
        """Generate a unique ID for a shape.

        Args:
            pptx_shape: The python-pptx shape object.

        Returns:
            Unique string identifier.
        """
        return f"shape_{pptx_shape.shape_id}_{uuid.uuid4().hex[:8]}"

    def _extract_transform(self, pptx_shape: BaseShape) -> Transform:
        """Extract transformation properties.

        Args:
            pptx_shape: The python-pptx shape object.

        Returns:
            Transform object.
        """
        rotation = 0.0
        if hasattr(pptx_shape, "rotation"):
            rotation = float(pptx_shape.rotation)

        return Transform(
            rotation=rotation,
            flip_h=False,  # TODO: Extract from PPTX XML
            flip_v=False,
            scale_x=1.0,
            scale_y=1.0,
        )

    def _extract_auto_shape(self, pptx_shape: BaseShape, shape_dict: dict[str, Any]) -> None:
        """Extract auto shape specific properties.

        Args:
            pptx_shape: The python-pptx shape object.
            shape_dict: Dictionary to populate with properties.
        """
        if hasattr(pptx_shape, "auto_shape_type") and pptx_shape.auto_shape_type:
            shape_dict["auto_shape_type"] = str(pptx_shape.auto_shape_type).split(".")[-1].lower()

        # Extract fill
        if isinstance(pptx_shape, PPTXAutoShape):
            shape_dict["fill"] = self.style_extractor.extract_fill(pptx_shape.fill)
            if pptx_shape.line:
                shape_dict["stroke"] = self.style_extractor.extract_stroke(pptx_shape.line)
            shape_dict["effects"] = self.style_extractor.extract_effects(pptx_shape)

        # Extract text if present
        if hasattr(pptx_shape, "has_text_frame") and pptx_shape.has_text_frame:
            shape_dict["text"] = self._extract_text_content(pptx_shape.text_frame)

    def _extract_freeform(self, pptx_shape: BaseShape, shape_dict: dict[str, Any]) -> None:
        """Extract freeform path data.

        Args:
            pptx_shape: The python-pptx shape object.
            shape_dict: Dictionary to populate with properties.
        """
        # TODO: Extract path commands from shape XML
        shape_dict["path"] = []

        # Extract fill and stroke
        if hasattr(pptx_shape, "fill"):
            shape_dict["fill"] = self.style_extractor.extract_fill(pptx_shape.fill)
        if hasattr(pptx_shape, "line") and pptx_shape.line:
            shape_dict["stroke"] = self.style_extractor.extract_stroke(pptx_shape.line)

    def _extract_text_shape(self, pptx_shape: BaseShape, shape_dict: dict[str, Any]) -> None:
        """Extract text box properties.

        Args:
            pptx_shape: The python-pptx shape object.
            shape_dict: Dictionary to populate with properties.
        """
        shape_dict["fill"] = NoFill()

        if hasattr(pptx_shape, "text_frame"):
            shape_dict["text"] = self._extract_text_content(pptx_shape.text_frame)

    def _extract_image(self, pptx_shape: BaseShape, shape_dict: dict[str, Any]) -> None:
        """Extract image properties.

        Args:
            pptx_shape: The python-pptx shape object.
            shape_dict: Dictionary to populate with properties.
        """
        if isinstance(pptx_shape, Picture):
            shape_dict["image_path"] = pptx_shape.image.filename
            shape_dict["fill"] = NoFill()

    def _extract_group(
        self,
        pptx_shape: BaseShape,
        shape_dict: dict[str, Any],
        group_path: list[str],
    ) -> None:
        """Extract group and its children.

        Args:
            pptx_shape: The python-pptx group shape object.
            shape_dict: Dictionary to populate with properties.
            group_path: Current group nesting path.
        """
        shape_dict["fill"] = NoFill()

        if isinstance(pptx_shape, GroupShape):
            new_group_path = group_path + [shape_dict["id"]]
            shape_dict["children"] = self.extract_shapes(pptx_shape.shapes, new_group_path)

    def _extract_text_content(self, text_frame: Any) -> TextContent:
        """Extract formatted text content.

        Args:
            text_frame: The python-pptx TextFrame object.

        Returns:
            TextContent with formatted runs.
        """
        runs: list[TextRun] = []

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                text_run = TextRun(
                    text=run.text,
                    font_family=font.name or "Calibri",
                    font_size=int(font.size.pt * 100) if font.size else 1400,
                    bold=bool(font.bold),
                    italic=bool(font.italic),
                    underline=bool(font.underline),
                    color=self._rgb_to_hex(font.color.rgb) if font.color and font.color.rgb else "#000000",
                )
                runs.append(text_run)

        return TextContent(
            runs=runs,
            alignment="left",  # TODO: Extract alignment
        )

    def _rgb_to_hex(self, rgb: Any) -> str:
        """Convert RGB color to hex string.

        Args:
            rgb: RGBColor object.

        Returns:
            Hex color string.
        """
        if rgb is None:
            return "#000000"
        return f"#{rgb}"
