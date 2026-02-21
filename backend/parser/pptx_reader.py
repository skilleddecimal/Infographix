"""High-level PPTX reading and parsing."""

from pathlib import Path
from typing import BinaryIO, Union

from pptx import Presentation
from pptx.slide import Slide

from backend.dsl.schema import Canvas, SlideMetadata, SlideScene, ThemeColors
from backend.parser.shape_extractor import ShapeExtractor
from backend.parser.style_extractor import StyleExtractor
from backend.parser.theme_parser import ThemeParser


class PPTXReader:
    """Reads PPTX files and extracts DSL scene graphs."""

    def __init__(self) -> None:
        """Initialize the PPTX reader."""
        self.shape_extractor = ShapeExtractor()
        self.style_extractor = StyleExtractor()
        self.theme_parser = ThemeParser()

    def read(self, source: Union[str, Path, BinaryIO]) -> list[SlideScene]:
        """Read a PPTX file and extract scene graphs for all slides.

        Args:
            source: Path to PPTX file or file-like object.

        Returns:
            List of SlideScene objects, one per slide.
        """
        prs = Presentation(source)
        scenes: list[SlideScene] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            scene = self._extract_slide(slide, slide_idx, prs)
            scenes.append(scene)

        return scenes

    def read_slide(self, source: Union[str, Path, BinaryIO], slide_number: int = 1) -> SlideScene:
        """Read a specific slide from a PPTX file.

        Args:
            source: Path to PPTX file or file-like object.
            slide_number: 1-based slide number to extract.

        Returns:
            SlideScene for the specified slide.

        Raises:
            IndexError: If slide_number is out of range.
        """
        scenes = self.read(source)
        if slide_number < 1 or slide_number > len(scenes):
            raise IndexError(f"Slide {slide_number} not found. File has {len(scenes)} slides.")
        return scenes[slide_number - 1]

    def _extract_slide(self, slide: Slide, slide_number: int, prs: Presentation) -> SlideScene:
        """Extract a single slide to a SlideScene.

        Args:
            slide: The python-pptx Slide object.
            slide_number: 1-based slide number.
            prs: The parent Presentation object.

        Returns:
            SlideScene representing the slide.
        """
        # Extract canvas dimensions
        canvas = Canvas(
            width=prs.slide_width,
            height=prs.slide_height,
            background=self.style_extractor.extract_background(slide),
        )

        # Extract theme colors
        theme = self._extract_theme(prs)

        # Extract shapes
        shapes = self.shape_extractor.extract_shapes(slide.shapes)

        # Build metadata
        metadata = SlideMetadata(
            slide_number=slide_number,
            layout_name=slide.slide_layout.name if slide.slide_layout else None,
            notes=self._extract_notes(slide),
        )

        return SlideScene(
            canvas=canvas,
            shapes=shapes,
            theme=theme,
            metadata=metadata,
        )

    def _extract_theme(self, prs: Presentation) -> ThemeColors:
        """Extract theme colors from presentation.

        Uses ThemeParser to extract colors from the slide master's
        color scheme XML.

        Args:
            prs: The Presentation object.

        Returns:
            ThemeColors with the theme palette.
        """
        return self.theme_parser.extract_theme(prs)

    def _extract_notes(self, slide: Slide) -> str | None:
        """Extract speaker notes from a slide.

        Args:
            slide: The Slide object.

        Returns:
            Speaker notes text or None.
        """
        if slide.has_notes_slide and slide.notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                return notes_frame.text
        return None
